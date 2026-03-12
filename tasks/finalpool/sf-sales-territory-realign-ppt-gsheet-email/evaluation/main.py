"""Evaluation for sf-sales-territory-realign-ppt-gsheet-email.

Checks:
1. Excel file Territory_Analysis.xlsx with 3 sheets matching groundtruth
2. Google Sheet "Territory Dashboard" in gsheet DB
3. PPT file Territory_Review.pptx with 5+ slides
4. 5 regional manager emails
"""
import argparse
import json
import os
import sys

import psycopg2

DB = {"host": os.environ.get("PGHOST", "localhost"), "port": 5432, "dbname": "toolathlon_gym", "user": "eigent", "password": "camel"}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = f": {str(detail)[:200]}" if detail else ""
        print(f"  [FAIL] {name}{detail_str}")


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def check_excel(agent_workspace, gt_dir):
    print("\n=== Checking Excel ===")
    try:
        import openpyxl
    except ImportError:
        check("openpyxl available", False, "openpyxl not installed")
        return

    agent_file = os.path.join(agent_workspace, "Territory_Analysis.xlsx")
    gt_file = os.path.join(gt_dir, "Territory_Analysis.xlsx")

    check("Territory_Analysis.xlsx exists", os.path.isfile(agent_file),
          f"Expected at {agent_file}")
    if not os.path.isfile(agent_file):
        return
    if not os.path.isfile(gt_file):
        check("Groundtruth Excel exists", False, f"Not at {gt_file}")
        return

    agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
    gt_wb = openpyxl.load_workbook(gt_file, data_only=True)

    # Sheet 1: Regional Performance
    a_rows = load_sheet_rows(agent_wb, "Regional Performance")
    g_rows = load_sheet_rows(gt_wb, "Regional Performance")
    check("Sheet 'Regional Performance' exists", a_rows is not None)
    if a_rows and g_rows:
        a_data = [r for r in (a_rows[1:] if len(a_rows) > 1 else []) if r and r[0] is not None]
        g_data = [r for r in (g_rows[1:] if g_rows and len(g_rows) > 1 else []) if r and r[0] is not None]
        check("Regional Performance has 5 data rows", len(a_data) >= 5, f"Got {len(a_data)}")

        if len(a_data) >= 5:
            a_lookup = {str(r[0]).strip().lower(): r for r in a_data if r[0]}
            for g_row in g_data:
                key = str(g_row[0]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    check(f"Region '{g_row[0]}' present", False)
                    continue
                # Total_Orders col 1
                if len(a_row) > 1 and len(g_row) > 1:
                    check(f"{g_row[0]} Total_Orders",
                          num_close(a_row[1], g_row[1], 10),
                          f"got {a_row[1]}, expected {g_row[1]}")
                # Total_Revenue col 2
                if len(a_row) > 2 and len(g_row) > 2:
                    check(f"{g_row[0]} Total_Revenue",
                          num_close(a_row[2], g_row[2], 5.0),
                          f"got {a_row[2]}, expected {g_row[2]}")
                # Avg_Order_Value col 3
                if len(a_row) > 3 and len(g_row) > 3:
                    check(f"{g_row[0]} Avg_Order_Value",
                          num_close(a_row[3], g_row[3], 1.0),
                          f"got {a_row[3]}, expected {g_row[3]}")
                # Revenue_Share_Pct col 6
                if len(a_row) > 6 and len(g_row) > 6:
                    check(f"{g_row[0]} Revenue_Share_Pct",
                          num_close(a_row[6], g_row[6], 0.5),
                          f"got {a_row[6]}, expected {g_row[6]}")
                # Status col 9
                if len(a_row) > 9 and len(g_row) > 9:
                    check(f"{g_row[0]} Status",
                          str_match(a_row[9], g_row[9]),
                          f"got {a_row[9]}, expected {g_row[9]}")

    # Sheet 2: Segment Mix
    a_seg = load_sheet_rows(agent_wb, "Segment Mix")
    check("Sheet 'Segment Mix' exists", a_seg is not None)
    if a_seg:
        a_seg_data = [r for r in (a_seg[1:] if len(a_seg) > 1 else []) if r and r[0] is not None]
        check("Segment Mix has 20 data rows", len(a_seg_data) >= 20,
              f"Got {len(a_seg_data)}")

    # Sheet 3: Realignment Recommendations
    a_rec = load_sheet_rows(agent_wb, "Realignment Recommendations")
    g_rec = load_sheet_rows(gt_wb, "Realignment Recommendations")
    check("Sheet 'Realignment Recommendations' exists", a_rec is not None)
    if a_rec and g_rec:
        a_rec_data = [r for r in (a_rec[1:] if len(a_rec) > 1 else []) if r and r[0] is not None]
        g_rec_data = [r for r in (g_rec[1:] if g_rec and len(g_rec) > 1 else []) if r and r[0] is not None]
        check("Recommendations has 5 data rows", len(a_rec_data) >= 5,
              f"Got {len(a_rec_data)}")

        if len(a_rec_data) >= 5:
            a_rec_lookup = {str(r[0]).strip().lower(): r for r in a_rec_data if r[0]}
            for g_row in g_rec_data:
                key = str(g_row[0]).strip().lower()
                a_row = a_rec_lookup.get(key)
                if a_row is None:
                    continue
                # Priority col 4
                if len(a_row) > 4 and len(g_row) > 4:
                    check(f"{g_row[0]} Priority",
                          str_match(a_row[4], g_row[4]),
                          f"got {a_row[4]}, expected {g_row[4]}")


def check_gsheet():
    print("\n=== Checking Google Sheet ===")
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT s.id, s.title
            FROM gsheet.spreadsheets s
            WHERE LOWER(s.title) LIKE '%territory%' AND LOWER(s.title) LIKE '%dashboard%'
        """)
        sheets = cur.fetchall()
        check("Territory Dashboard spreadsheet exists", len(sheets) >= 1,
              f"Found {len(sheets)} matching spreadsheets")

        if sheets:
            ss_id = sheets[0][0]
            cur.execute("""
                SELECT COUNT(DISTINCT c.row_index)
                FROM gsheet.cells c
                WHERE c.spreadsheet_id = %s AND c.row_index > 0
            """, (ss_id,))
            row_count = cur.fetchone()[0]
            check("Territory Dashboard has at least 5 data rows", row_count >= 5,
                  f"Got {row_count} data rows")

            # Check for key region values
            cur.execute("""
                SELECT value FROM gsheet.cells
                WHERE spreadsheet_id = %s
            """, (ss_id,))
            all_values = " ".join(str(r[0]).lower() for r in cur.fetchall() if r[0])
            check("Dashboard contains 'asia pacific'", "asia pacific" in all_values)
            check("Dashboard contains 'europe'", "europe" in all_values)
            check("Dashboard contains 'latin america'", "latin america" in all_values)

        cur.close()
        conn.close()
    except Exception as e:
        check("GSheet DB check", False, str(e))


def check_ppt(agent_workspace):
    print("\n=== Checking PPT File ===")
    ppt_path = os.path.join(agent_workspace, "Territory_Review.pptx")
    check("Territory_Review.pptx exists", os.path.isfile(ppt_path),
          f"Expected at {ppt_path}")
    if not os.path.isfile(ppt_path):
        return

    try:
        from pptx import Presentation
        prs = Presentation(ppt_path)
    except Exception as e:
        check("PPT file readable", False, str(e))
        return

    check("PPT has at least 5 slides", len(prs.slides) >= 5,
          f"Found {len(prs.slides)} slides")

    all_text = " ".join(
        shape.text.lower()
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )

    check("PPT contains 'territory' and 'realignment' or 'review'",
          "territory" in all_text and ("realign" in all_text or "review" in all_text),
          "Title keywords not found")
    check("PPT mentions Europe", "europe" in all_text)
    check("PPT mentions Latin America", "latin america" in all_text)
    check("PPT contains revenue figures",
          "648" in all_text or "642" in all_text or "549" in all_text,
          "Revenue figures not found")
    check("PPT contains recommendations or actions",
          "recommend" in all_text or "action" in all_text or "strategy" in all_text,
          "No recommendations found")


def check_emails():
    print("\n=== Checking Emails ===")
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT subject, from_addr, to_addr, body_text
            FROM email.messages
            WHERE folder_id != 0 OR folder_id IS NULL
        """)
        all_emails = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        check("Email DB query", False, str(e))
        return

    def parse_recipients(to_addr):
        if to_addr is None:
            return []
        if isinstance(to_addr, list):
            return [str(r).strip().lower() for r in to_addr]
        to_str = str(to_addr).strip()
        try:
            parsed = json.loads(to_str)
            if isinstance(parsed, list):
                return [str(r).strip().lower() for r in parsed]
            return [to_str.lower()]
        except (json.JSONDecodeError, TypeError):
            return [to_str.lower()]

    # Check for 5 regional emails
    region_managers = {
        "asiapacific_manager@company.com": "asia pacific",
        "europe_manager@company.com": "europe",
        "latam_manager@company.com": "latin america",
        "middleeast_manager@company.com": "middle east",
        "northamerica_manager@company.com": "north america",
    }

    found_count = 0
    all_recipients = []
    for subj, from_addr, to_addr, body in all_emails:
        recipients = parse_recipients(to_addr)
        all_recipients.extend(recipients)
        for mgr_email in region_managers:
            if mgr_email in recipients:
                found_count += 1
                # Check subject contains territory and review
                subj_lower = (subj or "").lower()
                region_name = region_managers[mgr_email]
                has_region_or_territory = ("territory" in subj_lower or
                                           region_name.split()[0] in subj_lower or
                                           "review" in subj_lower)
                check(f"Email to {mgr_email} has appropriate subject",
                      has_region_or_territory,
                      f"Subject: {subj}")
                break

    check("At least 5 regional manager emails sent", found_count >= 5,
          f"Found {found_count} emails to regional managers. All recipients: {all_recipients[:10]}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    print("=" * 70)
    print("SF SALES TERRITORY REALIGN PPT GSHEET EMAIL - EVALUATION")
    print("=" * 70)

    check_excel(args.agent_workspace, gt_dir)
    check_gsheet()
    check_ppt(args.agent_workspace)
    check_emails()

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    overall = FAIL_COUNT == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump({"passed": PASS_COUNT, "failed": FAIL_COUNT, "success": overall}, f, indent=2)

    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
