"""Evaluation for yf-peer-comparison-excel-ppt-email.

Checks:
1. Excel file (Peer_Comparison.xlsx) - 3 sheets with correct data
2. PowerPoint (Investor_Presentation.pptx) - >=6 slides
3. Emails - 3 emails to correct recipients
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

FILE_PASS = 0
FILE_FAIL = 0
DB_PASS = 0
DB_FAIL = 0


def check(name, condition, detail="", db=False):
    global FILE_PASS, FILE_FAIL, DB_PASS, DB_FAIL
    if condition:
        if db:
            DB_PASS += 1
        else:
            FILE_PASS += 1
        print(f"  [PASS] {name}")
    else:
        if db:
            DB_FAIL += 1
        else:
            FILE_FAIL += 1
        detail_str = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{detail_str}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def load_gt_sheet(gt_ws, sheet_name):
    """Load groundtruth sheet rows as list of tuples."""
    gt_path = os.path.join(gt_ws, "Peer_Comparison.xlsx")
    if not os.path.exists(gt_path):
        return None
    wb = openpyxl.load_workbook(gt_path, data_only=True)
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def check_excel(agent_workspace, groundtruth_workspace):
    print("\n=== Checking Excel Output ===")
    excel_path = os.path.join(agent_workspace, "Peer_Comparison.xlsx")
    check("Excel file exists", os.path.isfile(excel_path), f"Expected {excel_path}")
    if not os.path.isfile(excel_path):
        return

    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
    except Exception as e:
        check("Excel file readable", False, str(e))
        return

    # --- Sheet 1: Company Profiles ---
    ws1 = None
    for s in wb.sheetnames:
        if "company" in s.lower() and "profile" in s.lower():
            ws1 = wb[s]
            break
    if ws1 is None:
        for s in wb.sheetnames:
            if "company" in s.lower() or "profile" in s.lower():
                ws1 = wb[s]
                break

    check("Sheet 'Company Profiles' exists", ws1 is not None, f"Sheets: {wb.sheetnames}")
    if ws1 is not None:
        rows = list(ws1.iter_rows(min_row=2, values_only=True))
        data_rows = [r for r in rows if r and r[0] is not None]
        check("Company Profiles has 5 rows", len(data_rows) == 5, f"Got {len(data_rows)}")

        # Check symbols present
        symbols_found = {str(r[0]).strip().upper() for r in data_rows if r[0]}
        for sym in ["AMZN", "GOOGL", "JNJ", "JPM", "XOM"]:
            check(f"{sym} in Company Profiles", sym in symbols_found,
                  f"Found: {symbols_found}")

        # Check alphabetical order
        sym_list = [str(r[0]).strip().upper() for r in data_rows if r[0]]
        check("Symbols sorted alphabetically", sym_list == sorted(sym_list),
              f"Order: {sym_list}")

        # Compare with groundtruth values
        gt_rows = load_gt_sheet(groundtruth_workspace, "Company Profiles")
        if gt_rows:
            gt_data = {}
            for r in gt_rows[1:]:
                if r and r[0]:
                    gt_data[str(r[0]).strip().upper()] = r

            agent_data = {}
            for r in data_rows:
                if r and r[0]:
                    agent_data[str(r[0]).strip().upper()] = r

            for sym in ["AMZN", "GOOGL", "JNJ", "JPM", "XOM"]:
                if sym in agent_data and sym in gt_data:
                    ar = agent_data[sym]
                    gr = gt_data[sym]
                    # Market Cap (col 3, tol 1M)
                    if len(ar) > 3 and len(gr) > 3 and ar[3] is not None:
                        check(f"{sym} Market_Cap", num_close(ar[3], gr[3], 1e6),
                              f"Agent={ar[3]}, GT={gr[3]}")
                    # Trailing PE (col 4, tol 0.5)
                    if len(ar) > 4 and len(gr) > 4 and ar[4] is not None:
                        check(f"{sym} Trailing_PE", num_close(ar[4], gr[4], 0.5),
                              f"Agent={ar[4]}, GT={gr[4]}")
                    # YTD Return (col 10, tol 0.5)
                    if len(ar) > 10 and len(gr) > 10 and ar[10] is not None:
                        check(f"{sym} YTD_Return_Pct", num_close(ar[10], gr[10], 0.5),
                              f"Agent={ar[10]}, GT={gr[10]}")

    # --- Sheet 2: Financial Comparison ---
    ws2 = None
    for s in wb.sheetnames:
        if "financial" in s.lower() and "comparison" in s.lower():
            ws2 = wb[s]
            break
    if ws2 is None:
        for s in wb.sheetnames:
            if "financial" in s.lower():
                ws2 = wb[s]
                break

    check("Sheet 'Financial Comparison' exists", ws2 is not None, f"Sheets: {wb.sheetnames}")
    if ws2 is not None:
        rows2 = list(ws2.iter_rows(min_row=2, values_only=True))
        data_rows2 = [r for r in rows2 if r and r[0] is not None]
        check("Financial Comparison has 5 rows", len(data_rows2) == 5, f"Got {len(data_rows2)}")

        # Compare revenue for spot check
        gt_rows2 = load_gt_sheet(groundtruth_workspace, "Financial Comparison")
        if gt_rows2:
            gt_fin = {str(r[0]).strip().upper(): r for r in gt_rows2[1:] if r and r[0]}
            agent_fin = {str(r[0]).strip().upper(): r for r in data_rows2 if r and r[0]}
            for sym in ["AMZN", "GOOGL"]:
                if sym in agent_fin and sym in gt_fin:
                    ar = agent_fin[sym]
                    gr = gt_fin[sym]
                    # Revenue (col 1, relative tolerance 5%)
                    if len(ar) > 1 and len(gr) > 1 and ar[1] is not None and gr[1] is not None:
                        rel_err = abs(float(ar[1]) - float(gr[1])) / float(gr[1]) if float(gr[1]) != 0 else 1
                        check(f"{sym} Revenue within 5%", rel_err < 0.05,
                              f"Agent={ar[1]}, GT={gr[1]}, err={rel_err:.4f}")

    # --- Sheet 3: Scoring ---
    ws3 = None
    for s in wb.sheetnames:
        if "scor" in s.lower():
            ws3 = wb[s]
            break

    check("Sheet 'Scoring' exists", ws3 is not None, f"Sheets: {wb.sheetnames}")
    if ws3 is not None:
        rows3 = list(ws3.iter_rows(min_row=2, values_only=True))
        data_rows3 = [r for r in rows3 if r and r[0] is not None]
        check("Scoring has 5 rows", len(data_rows3) == 5, f"Got {len(data_rows3)}")

        gt_rows3 = load_gt_sheet(groundtruth_workspace, "Scoring")
        if gt_rows3:
            gt_score = {str(r[0]).strip().upper(): r for r in gt_rows3[1:] if r and r[0]}
            agent_score = {str(r[0]).strip().upper(): r for r in data_rows3 if r and r[0]}

            for sym in ["AMZN", "GOOGL", "JNJ", "JPM", "XOM"]:
                if sym in agent_score and sym in gt_score:
                    ar = agent_score[sym]
                    gr = gt_score[sym]
                    # Weighted Score (col 6, tol 0.5)
                    if len(ar) > 6 and len(gr) > 6 and ar[6] is not None:
                        check(f"{sym} Weighted_Score", num_close(ar[6], gr[6], 0.5),
                              f"Agent={ar[6]}, GT={gr[6]}")
                    # Overall Rating (col 7, exact match)
                    if len(ar) > 7 and len(gr) > 7 and ar[7] is not None:
                        agent_rating = str(ar[7]).strip().lower()
                        gt_rating = str(gr[7]).strip().lower()
                        check(f"{sym} Overall_Rating", agent_rating == gt_rating,
                              f"Agent='{ar[7]}', GT='{gr[7]}'")


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint Output ===")
    pptx_path = os.path.join(agent_workspace, "Investor_Presentation.pptx")
    check("PPTX file exists", os.path.isfile(pptx_path), f"Expected {pptx_path}")
    if not os.path.isfile(pptx_path):
        return

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        check("PPTX file readable", False, str(e))
        return

    slide_count = len(prs.slides)
    check("PPTX has >= 6 slides", slide_count >= 6, f"Got {slide_count} slides")

    # Check for key content in slides
    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text.lower() + " "

    check("PPTX contains 'peer' or 'comparison'",
          "peer" in all_text or "comparison" in all_text,
          f"Text sample: {all_text[:200]}")
    check("PPTX mentions company symbols",
          "amzn" in all_text or "amazon" in all_text,
          f"Text sample: {all_text[:200]}")
    check("PPTX contains scoring or recommendation content",
          "score" in all_text or "rank" in all_text or "recommendation" in all_text or "buy" in all_text or "hold" in all_text,
          f"Text sample: {all_text[:200]}")


def check_emails():
    print("\n=== Checking Emails ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Check email to portfolio_managers@firm.com
    cur.execute("""
        SELECT subject, to_addr FROM email.messages
        WHERE to_addr::text ILIKE '%portfolio_managers@firm.com%'
        ORDER BY id DESC LIMIT 5
    """)
    pm_rows = cur.fetchall()
    check("Email to portfolio_managers@firm.com", len(pm_rows) > 0,
          "No email found", db=True)

    # Check email to research_team@firm.com
    cur.execute("""
        SELECT subject, to_addr FROM email.messages
        WHERE to_addr::text ILIKE '%research_team@firm.com%'
        ORDER BY id DESC LIMIT 5
    """)
    rt_rows = cur.fetchall()
    check("Email to research_team@firm.com", len(rt_rows) > 0,
          "No email found", db=True)

    # Check email to compliance@firm.com
    cur.execute("""
        SELECT subject, to_addr FROM email.messages
        WHERE to_addr::text ILIKE '%compliance@firm.com%'
        ORDER BY id DESC LIMIT 5
    """)
    comp_rows = cur.fetchall()
    check("Email to compliance@firm.com", len(comp_rows) > 0,
          "No email found", db=True)

    cur.close()
    conn.close()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    gt_ws = args.groundtruth_workspace or os.path.join(
        os.path.dirname(__file__), "..", "groundtruth_workspace")

    check_excel(args.agent_workspace, gt_ws)
    check_pptx(args.agent_workspace)
    check_emails()

    total_pass = FILE_PASS + DB_PASS
    total_fail = FILE_FAIL + DB_FAIL
    file_ok = FILE_FAIL == 0

    print(f"\n=== SUMMARY ===")
    print(f"  File checks - Passed: {FILE_PASS}, Failed: {FILE_FAIL}")
    print(f"  DB checks   - Passed: {DB_PASS}, Failed: {DB_FAIL}")
    if DB_FAIL > 0:
        print(f"  WARNING: {DB_FAIL} DB checks failed (not blocking)")
    print(f"  Overall: {'PASS' if file_ok else 'FAIL'}")

    if args.res_log_file:
        result = {"passed": total_pass, "failed": total_fail, "success": file_ok}
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if file_ok else 1)


if __name__ == "__main__":
    main()
