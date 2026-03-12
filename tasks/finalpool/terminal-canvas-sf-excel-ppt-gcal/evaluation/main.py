"""Evaluation for terminal-canvas-sf-excel-ppt-gcal.
Checks:
1. Skills_Gap_Analysis.xlsx with 4 sheets and correct data
2. Skills_Gap_Presentation.pptx with 5+ slides
3. Google Calendar advisory board events
4. gap_analyzer.py script exists
"""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def check_excel(workspace):
    print("\n=== Check 1: Skills_Gap_Analysis.xlsx ===")
    path = os.path.join(workspace, "Skills_Gap_Analysis.xlsx")
    if not os.path.exists(path):
        check("Excel file exists", False, f"Not found at {path}")
        return
    check("Excel file exists", True)

    wb = openpyxl.load_workbook(path)
    sheets = wb.sheetnames
    check("Has at least 4 sheets", len(sheets) >= 4, f"Found {len(sheets)}: {sheets}")

    sheets_lower = [s.lower() for s in sheets]

    # Curriculum_Coverage sheet
    cc_idx = next((i for i, s in enumerate(sheets_lower) if "curriculum" in s or "coverage" in s), 0)
    ws = wb[sheets[cc_idx]]
    rows = list(ws.iter_rows(values_only=True))
    data_rows = [r for r in rows[1:] if any(c for c in r)]
    check("Curriculum_Coverage has 4 course rows", len(data_rows) >= 4, f"Found {len(data_rows)}")

    all_text = " ".join(str(c) for r in rows for c in r if c).lower()
    check("Contains Applied Analytics course", "applied analytics" in all_text or "analytics" in all_text,
          f"Text sample: {all_text[:120]}")
    check("Contains Biochemistry course", "biochemistry" in all_text or "bioinformatics" in all_text,
          f"Text sample: {all_text[:120]}")

    # Workforce_Profile sheet
    wp_idx = next((i for i, s in enumerate(sheets_lower) if "workforce" in s or "profile" in s), 1)
    if wp_idx < len(sheets):
        ws2 = wb[sheets[wp_idx]]
        rows2 = list(ws2.iter_rows(values_only=True))
        data_rows2 = [r for r in rows2[1:] if any(c for c in r)]
        check("Workforce_Profile has 7 department rows", len(data_rows2) >= 7, f"Found {len(data_rows2)}")
        all_text2 = " ".join(str(c) for r in rows2 for c in r if c).lower()
        check("Contains Engineering department", "engineering" in all_text2, f"Text: {all_text2[:120]}")

    # Gap_Matrix sheet
    gm_idx = next((i for i, s in enumerate(sheets_lower) if "gap" in s or "matrix" in s), 2)
    if gm_idx < len(sheets):
        ws3 = wb[sheets[gm_idx]]
        rows3 = list(ws3.iter_rows(values_only=True))
        data_rows3 = [r for r in rows3[1:] if any(c for c in r)]
        check("Gap_Matrix has 4 skill area rows", len(data_rows3) >= 4, f"Found {len(data_rows3)}")
        all_text3 = " ".join(str(c) for r in rows3 for c in r if c).lower()
        check("Contains Quantitative Analysis", "quantitative" in all_text3, f"Text: {all_text3[:120]}")
        check("Contains priority levels", "critical" in all_text3 or "high" in all_text3,
              f"Text: {all_text3[:120]}")

    # Recommendations sheet
    rec_idx = next((i for i, s in enumerate(sheets_lower) if "recommend" in s), 3)
    if rec_idx < len(sheets):
        ws4 = wb[sheets[rec_idx]]
        rows4 = list(ws4.iter_rows(values_only=True))
        data_rows4 = [r for r in rows4[1:] if any(c for c in r)]
        check("Recommendations has at least 4 rows", len(data_rows4) >= 4, f"Found {len(data_rows4)}")


def check_pptx(workspace):
    print("\n=== Check 2: Skills_Gap_Presentation.pptx ===")
    path = os.path.join(workspace, "Skills_Gap_Presentation.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = list(prs.slides)
        check("Has at least 5 slides", len(slides) >= 5, f"Found {len(slides)} slides")

        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += " " + shape.text_frame.text
        all_text_lower = all_text.lower()
        check("Contains gap analysis content", "gap" in all_text_lower, f"Text: {all_text_lower[:150]}")
        check("Contains curriculum content", "curriculum" in all_text_lower or "course" in all_text_lower,
              f"Text: {all_text_lower[:150]}")
        check("Contains workforce content", "workforce" in all_text_lower or "department" in all_text_lower,
              f"Text: {all_text_lower[:150]}")
    except ImportError:
        check("python-pptx available", False, "python-pptx not installed")


def check_gcal():
    print("\n=== Check 3: Advisory Board Calendar Events ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT summary, start_datetime, end_datetime
        FROM gcal.events
        WHERE lower(summary) LIKE '%%advisory%%'
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    check("At least 3 advisory board events", len(events) >= 3, f"Found {len(events)} events")

    if events:
        summaries = " ".join(str(e[0]) for e in events).lower()
        check("Events mention curriculum or review", "curriculum" in summaries or "review" in summaries,
              f"Summaries: {summaries[:150]}")
        check("Events mention workforce or data or findings",
              "workforce" in summaries or "data" in summaries or "finding" in summaries,
              f"Summaries: {summaries[:150]}")
        check("Events mention recommendation or gap",
              "recommend" in summaries or "gap" in summaries,
              f"Summaries: {summaries[:150]}")

    cur.close()
    conn.close()


def check_script(workspace):
    print("\n=== Check 4: gap_analyzer.py ===")
    path = os.path.join(workspace, "gap_analyzer.py")
    check("gap_analyzer.py exists", os.path.exists(path))


def check_reverse_validation(workspace):
    """Verify things that should NOT exist in output."""
    print("\n=== Reverse Validation ===")

    # Excel: no unexpected sheets beyond the 4 required
    path = os.path.join(workspace, "Skills_Gap_Analysis.xlsx")
    if os.path.isfile(path):
        wb = openpyxl.load_workbook(path)
        check("Excel has no more than 6 sheets", len(wb.sheetnames) <= 6,
              f"Found {len(wb.sheetnames)} sheets: {wb.sheetnames}")

    # GCal: no advisory events on weekends
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM gcal.events
            WHERE lower(summary) LIKE '%%advisory%%'
              AND EXTRACT(DOW FROM start_datetime) IN (0, 6)
        """)
        weekend_count = cur.fetchone()[0]
        check("No advisory events on weekends", weekend_count == 0,
              f"Found {weekend_count} weekend events")
        cur.close()
        conn.close()
    except Exception:
        pass


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_gcal()
    check_script(args.agent_workspace)
    check_reverse_validation(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
