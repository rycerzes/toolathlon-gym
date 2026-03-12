"""Evaluation for sf-sales-customer-segment-excel-ppt-gform."""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    agent_ws = args.agent_workspace or task_root
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    all_errors = []

    # --- Check 1: Excel workbook ---
    print("Checking Excel workbook...")
    agent_file = os.path.join(agent_ws, "Customer_Segment_Analysis.xlsx")
    gt_file = os.path.join(gt_dir, "Customer_Segment_Analysis.xlsx")

    if not os.path.exists(agent_file):
        all_errors.append("Customer_Segment_Analysis.xlsx not found in agent workspace")
    elif not os.path.exists(gt_file):
        all_errors.append("Customer_Segment_Analysis.xlsx not found in groundtruth")
    else:
        agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_file, data_only=True)

        # Sheet 1: Segment Performance
        print("  Checking Segment Performance...")
        a_rows = load_sheet_rows(agent_wb, "Segment Performance")
        g_rows = load_sheet_rows(gt_wb, "Segment Performance")
        if a_rows is None:
            all_errors.append("Sheet 'Segment Performance' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Segment Performance' not found in groundtruth")
        else:
            errors = []
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []

            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None:
                    a_lookup[str(row[0]).strip().lower()] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = str(g_row[0]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    errors.append(f"Missing row: {g_row[0]}")
                    continue

                # Customer_Count (col 1)
                if len(a_row) > 1 and len(g_row) > 1:
                    if not num_close(a_row[1], g_row[1], 5):
                        errors.append(f"{key}.Customer_Count: {a_row[1]} vs {g_row[1]}")
                # Total_Orders (col 2)
                if len(a_row) > 2 and len(g_row) > 2:
                    if not num_close(a_row[2], g_row[2], 10):
                        errors.append(f"{key}.Total_Orders: {a_row[2]} vs {g_row[2]}")
                # Total_Revenue (col 3)
                if len(a_row) > 3 and len(g_row) > 3:
                    if not num_close(a_row[3], g_row[3], 500):
                        errors.append(f"{key}.Total_Revenue: {a_row[3]} vs {g_row[3]}")
                # Avg_Order_Value (col 4)
                if len(a_row) > 4 and len(g_row) > 4:
                    if not num_close(a_row[4], g_row[4], 5):
                        errors.append(f"{key}.Avg_Order_Value: {a_row[4]} vs {g_row[4]}")
                # Revenue_Share_Pct (col 7)
                if len(a_row) > 7 and len(g_row) > 7:
                    if not num_close(a_row[7], g_row[7], 2):
                        errors.append(f"{key}.Revenue_Share_Pct: {a_row[7]} vs {g_row[7]}")

            if errors:
                all_errors.extend(errors)
                print(f"    ERRORS: {len(errors)}")
                for e in errors[:5]:
                    print(f"      {e}")
            else:
                print("    PASS")

        # Sheet 2: Segment by Region
        print("  Checking Segment by Region...")
        a_rows = load_sheet_rows(agent_wb, "Segment by Region")
        g_rows = load_sheet_rows(gt_wb, "Segment by Region")
        if a_rows is None:
            all_errors.append("Sheet 'Segment by Region' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Segment by Region' not found in groundtruth")
        else:
            errors = []
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []

            # Build lookup by region+segment
            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None and row[1] is not None:
                    key = f"{str(row[0]).strip().lower()}|{str(row[1]).strip().lower()}"
                    a_lookup[key] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = f"{str(g_row[0]).strip().lower()}|{str(g_row[1]).strip().lower()}"
                a_row = a_lookup.get(key)
                if a_row is None:
                    errors.append(f"Missing row: {g_row[0]}|{g_row[1]}")
                    continue

                # Revenue (col 4)
                if len(a_row) > 4 and len(g_row) > 4:
                    if not num_close(a_row[4], g_row[4], 500):
                        errors.append(f"{key}.Revenue: {a_row[4]} vs {g_row[4]}")
                # Profitability_Index (col 6)
                if len(a_row) > 6 and len(g_row) > 6:
                    if not num_close(a_row[6], g_row[6], 2):
                        errors.append(f"{key}.Profitability_Index: {a_row[6]} vs {g_row[6]}")

            if errors:
                all_errors.extend(errors)
                print(f"    ERRORS: {len(errors)}")
                for e in errors[:5]:
                    print(f"      {e}")
            else:
                print("    PASS")

        # Sheet 3: Strategic Matrix
        print("  Checking Strategic Matrix...")
        a_rows = load_sheet_rows(agent_wb, "Strategic Matrix")
        g_rows = load_sheet_rows(gt_wb, "Strategic Matrix")
        if a_rows is None:
            all_errors.append("Sheet 'Strategic Matrix' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Strategic Matrix' not found in groundtruth")
        else:
            errors = []
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []

            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None:
                    a_lookup[str(row[0]).strip().lower()] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = str(g_row[0]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    errors.append(f"Missing row: {g_row[0]}")
                    continue

                # Revenue_Contribution_Pct (col 1)
                if len(a_row) > 1 and len(g_row) > 1:
                    if not num_close(a_row[1], g_row[1], 2):
                        errors.append(f"{key}.Revenue_Contribution_Pct: {a_row[1]} vs {g_row[1]}")
                # Growth_Indicator (col 2)
                if len(a_row) > 2 and len(g_row) > 2:
                    if not str_match(a_row[2], g_row[2]):
                        errors.append(f"{key}.Growth_Indicator: {a_row[2]} vs {g_row[2]}")
                # Strategic_Category (col 4)
                if len(a_row) > 4 and len(g_row) > 4:
                    if not str_match(a_row[4], g_row[4]):
                        errors.append(f"{key}.Strategic_Category: {a_row[4]} vs {g_row[4]}")

            if errors:
                all_errors.extend(errors)
                print(f"    ERRORS: {len(errors)}")
                for e in errors[:5]:
                    print(f"      {e}")
            else:
                print("    PASS")

    # --- Check 2: PowerPoint ---
    print("Checking PowerPoint presentation...")
    pptx_path = os.path.join(agent_ws, "QBR_Presentation.pptx")
    if not os.path.exists(pptx_path):
        all_errors.append("QBR_Presentation.pptx not found in agent workspace")
    else:
        from pptx import Presentation
        prs = Presentation(pptx_path)

        if len(prs.slides) < 6:
            all_errors.append(f"PPT has only {len(prs.slides)} slides, expected at least 6")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += " " + shape.text

        for seg in ["Consumer", "Enterprise", "Government", "SMB"]:
            if seg not in all_text:
                all_errors.append(f"PPT missing '{seg}' segment")

        # Check for revenue figure
        revenue_present = any(x in all_text for x in ["839609", "839,609"])
        if not revenue_present:
            all_errors.append("PPT missing Consumer revenue figure (~839609)")

        # Check for strategic categories
        for cat in ["Star", "Cash Cow"]:
            if cat not in all_text:
                all_errors.append(f"PPT missing strategic category '{cat}'")

        print(f"    PPT checks done ({len(prs.slides)} slides)")

    # --- Check 3: Google Form ---
    print("Checking Google Form...")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT id FROM gform.forms WHERE LOWER(title) LIKE '%customer experience%'")
        rows = cur.fetchall()
        if not rows:
            all_errors.append("Google Form 'Customer Experience Survey' not found in gform.forms")
        else:
            form_id = rows[0][0]
            cur.execute("SELECT COUNT(*) FROM gform.questions WHERE form_id = %s", (form_id,))
            q_count = cur.fetchone()[0]
            if q_count < 6:
                all_errors.append(f"Google Form has only {q_count} questions, expected at least 6")
            else:
                print(f"    GForm found with {q_count} questions")
        cur.close()
        conn.close()
    except Exception as e:
        all_errors.append(f"Error checking GForm: {e}")

    # --- Final result ---
    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:15]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
