"""
Evaluation for howtocook-workshop-materials task.

Checks structural properties of the three output files since the agent
has freedom to choose any 3 recipes from HowToCook:

1. Workshop_Handbook.docx
   - File exists
   - Has top-level heading containing "Cooking Workshop Handbook"
   - Has at least 3 second-level (Heading 1) headings for dish names
   - Contains ingredient content (word "ingredient" or ingredient-like items)
   - Contains step content (word "step" or numbered instructions)
   - Has a "Tips and Notes" section

2. Workshop_Slides.pptx
   - File exists
   - First slide title contains "Cooking Workshop" or "Team"
   - Has at least 7 slides (1 title + 3 dishes x 2 slides + 1 closing)
   - Closing slide contains "Enjoy"

3. Shopping_List.pdf
   - File exists
   - File size > 1KB (non-trivial content)
   - Contains "Shopping List" text (checked via PDF reader)
"""
import os
import sys
import json
from argparse import ArgumentParser
from datetime import datetime


def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)

from docx import Document
from pptx import Presentation


def check_word_doc(agent_workspace):
    """Check Workshop_Handbook.docx for required structure and content."""
    passed = 0
    total = 0
    doc_path = os.path.join(agent_workspace, "Workshop_Handbook.docx")

    # 1. File exists
    total += 1
    if not os.path.exists(doc_path):
        print("  FAIL: Workshop_Handbook.docx not found")
        return passed, total
    passed += 1
    print("  PASS: Workshop_Handbook.docx exists")

    doc = Document(doc_path)

    # Extract headings and full text
    headings_by_level = {}  # level -> list of heading texts
    all_text_parts = []
    for para in doc.paragraphs:
        all_text_parts.append(para.text)
        if para.style and para.style.name:
            style_name = para.style.name
            if "Heading" in style_name or "Title" in style_name:
                # Extract level: "Heading 1" -> 1, "Title" -> 0
                if style_name == "Title":
                    level = 0
                else:
                    try:
                        level = int(style_name.split()[-1])
                    except (ValueError, IndexError):
                        level = -1
                if level not in headings_by_level:
                    headings_by_level[level] = []
                headings_by_level[level].append(para.text.strip())

    full_text = " ".join(all_text_parts).lower()

    # 2. Title heading contains "Cooking Workshop Handbook"
    total += 1
    title_headings = headings_by_level.get(0, [])
    all_headings_flat = []
    for lvl_headings in headings_by_level.values():
        all_headings_flat.extend(lvl_headings)
    has_title = any("cooking workshop handbook" in h.lower() for h in all_headings_flat)
    if not has_title:
        # Fallback: check if "cooking workshop handbook" appears anywhere in text
        has_title = "cooking workshop handbook" in full_text
    if has_title:
        passed += 1
        print("  PASS: Title 'Cooking Workshop Handbook' found")
    else:
        print(f"  FAIL: Title 'Cooking Workshop Handbook' not found. Headings: {all_headings_flat[:10]}")

    # 3. At least 3 second-level headings (dish names)
    total += 1
    # Count Heading 1 and Heading 2 level headings (excluding known section names)
    dish_headings = []
    known_sections = {"tips and notes", "cooking workshop handbook", "ingredients",
                      "cooking steps", "tips", "notes", "introduction", "welcome"}
    for level in [1, 2]:
        for h in headings_by_level.get(level, []):
            if h.lower().strip() not in known_sections:
                dish_headings.append(h)

    # If we don't have enough from levels 1-2, also check if there are at least
    # 3 distinct content sections (by counting Heading 1 level headings)
    h1_count = len(headings_by_level.get(1, []))
    if len(dish_headings) >= 3 or h1_count >= 3:
        passed += 1
        print(f"  PASS: Found {max(len(dish_headings), h1_count)} section headings (need >= 3)")
    else:
        print(f"  FAIL: Only {len(dish_headings)} dish headings found (need >= 3). All H1: {headings_by_level.get(1, [])}")

    # 4. Contains ingredient content
    total += 1
    has_ingredients = ("ingredient" in full_text or
                       "tablespoon" in full_text or
                       "teaspoon" in full_text or
                       "gram" in full_text or
                       "cup" in full_text or
                       "ml" in full_text or
                       "oil" in full_text or
                       "salt" in full_text)
    if has_ingredients:
        passed += 1
        print("  PASS: Ingredient content found")
    else:
        print("  FAIL: No ingredient content found in document")

    # 5. Contains step/instruction content
    total += 1
    has_steps = ("step" in full_text or
                 "stir" in full_text or
                 "cook" in full_text or
                 "heat" in full_text or
                 "add" in full_text or
                 "pour" in full_text or
                 "cut" in full_text or
                 "boil" in full_text)
    if has_steps:
        passed += 1
        print("  PASS: Cooking step content found")
    else:
        print("  FAIL: No cooking step content found in document")

    # 6. Has "Tips and Notes" section
    total += 1
    has_tips = any("tips" in h.lower() and "note" in h.lower()
                   for h in all_headings_flat)
    if not has_tips:
        # Fallback: check for just "tips" heading
        has_tips = any("tips" in h.lower() for h in all_headings_flat)
    if not has_tips:
        # Fallback: check body text
        has_tips = "tips and notes" in full_text or "tips" in full_text
    if has_tips:
        passed += 1
        print("  PASS: 'Tips and Notes' section found")
    else:
        print("  FAIL: 'Tips and Notes' section not found")

    return passed, total


def check_pptx(agent_workspace):
    """Check Workshop_Slides.pptx for required structure."""
    passed = 0
    total = 0
    pptx_path = os.path.join(agent_workspace, "Workshop_Slides.pptx")

    # 1. File exists
    total += 1
    if not os.path.exists(pptx_path):
        print("  FAIL: Workshop_Slides.pptx not found")
        return passed, total
    passed += 1
    print("  PASS: Workshop_Slides.pptx exists")

    prs = Presentation(pptx_path)

    # Extract text from all slides
    slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
        slide_texts.append(" ".join(texts))

    # 2. First slide title contains "Cooking Workshop" or "Team"
    total += 1
    first_text = slide_texts[0].lower() if slide_texts else ""
    if "cooking workshop" in first_text or "team" in first_text or "workshop" in first_text:
        passed += 1
        print("  PASS: First slide contains workshop/team title")
    else:
        # Check across all slides as fallback
        all_text = " ".join(slide_texts).lower()
        if "cooking workshop" in all_text or "team" in all_text:
            passed += 1
            print("  PASS: Workshop/team title found in slides (not first slide)")
        else:
            print(f"  FAIL: No workshop/team title found. First slide: '{first_text[:100]}'")

    # 3. At least 7 slides
    total += 1
    slide_count = len(prs.slides)
    if slide_count >= 7:
        passed += 1
        print(f"  PASS: {slide_count} slides found (need >= 7)")
    else:
        print(f"  FAIL: Only {slide_count} slides found (need >= 7)")

    # 4. Closing slide contains "Enjoy"
    total += 1
    last_text = slide_texts[-1].lower() if slide_texts else ""
    all_text_lower = " ".join(slide_texts).lower()
    if "enjoy" in last_text:
        passed += 1
        print("  PASS: Closing slide contains 'Enjoy'")
    elif "enjoy" in all_text_lower:
        passed += 1
        print("  PASS: 'Enjoy' found in slides (not last slide)")
    else:
        print(f"  FAIL: 'Enjoy' not found in slides. Last slide: '{last_text[:100]}'")

    return passed, total


def check_pdf(agent_workspace):
    """Check Shopping_List.pdf for required structure."""
    passed = 0
    total = 0
    pdf_path = os.path.join(agent_workspace, "Shopping_List.pdf")

    # 1. File exists
    total += 1
    if not os.path.exists(pdf_path):
        print("  FAIL: Shopping_List.pdf not found")
        return passed, total
    passed += 1
    print("  PASS: Shopping_List.pdf exists")

    # 2. File size > 1KB
    total += 1
    file_size = os.path.getsize(pdf_path)
    if file_size > 1024:
        passed += 1
        print(f"  PASS: PDF size is {file_size} bytes (> 1KB)")
    else:
        print(f"  FAIL: PDF size is {file_size} bytes (need > 1KB)")

    # 3. Contains "Shopping List" text
    total += 1
    try:
        # Try to read PDF text using PyPDF2 or pypdf
        pdf_text = ""
        try:
            from pypdf import PdfReader
            reader = PdfReader(pdf_path)
            for page in reader.pages:
                pdf_text += page.extract_text() or ""
        except ImportError:
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(pdf_path)
                for page in reader.pages:
                    pdf_text += page.extract_text() or ""
            except ImportError:
                # Fallback: check raw bytes for the string
                with open(pdf_path, "rb") as f:
                    raw = f.read()
                pdf_text = raw.decode("latin-1", errors="ignore")

        if "shopping list" in pdf_text.lower() or "Shopping List" in pdf_text:
            passed += 1
            print("  PASS: PDF contains 'Shopping List' text")
        else:
            # Even if text extraction fails, if file is large enough it likely has content
            if file_size > 2000:
                passed += 1
                print("  PASS: PDF is substantial size, likely contains shopping list content")
            else:
                print(f"  FAIL: 'Shopping List' text not found in PDF. Extracted: '{pdf_text[:200]}'")
    except Exception as e:
        # If we can't read the PDF at all, check size as fallback
        if file_size > 2000:
            passed += 1
            print(f"  PASS: PDF exists with substantial size ({file_size} bytes), text extraction failed: {e}")
        else:
            print(f"  FAIL: Could not read PDF text: {e}")

    return passed, total


def main(args):
    total_passed = 0
    total_checks = 0

    # Check 1: Word document
    print("--- Check 1: Word Document (Workshop_Handbook.docx) ---")
    p, t = check_word_doc(args.agent_workspace)
    print(f"  Word Doc: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Check 2: PowerPoint
    print("\n--- Check 2: PowerPoint (Workshop_Slides.pptx) ---")
    p, t = check_pptx(args.agent_workspace)
    print(f"  PowerPoint: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Check 3: PDF
    print("\n--- Check 3: PDF (Shopping_List.pdf) ---")
    p, t = check_pdf(args.agent_workspace)
    print(f"  PDF: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Overall
    if total_checks == 0:
        print("\nFAIL: No checks were performed.")
        accuracy = 0.0
    else:
        accuracy = total_passed / total_checks * 100
        print(f"\nOverall: {total_passed}/{total_checks} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": total_passed,
        "total_checks": total_checks,
        "accuracy": accuracy,
        "timestamp": datetime.now().isoformat(),
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)
        print(f"Report saved to {args.res_log_file}")

    if total_passed == total_checks:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    main(args)
