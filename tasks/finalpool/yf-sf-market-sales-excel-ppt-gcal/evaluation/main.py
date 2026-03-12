"""Evaluation script for yf-sf-market-sales-excel-ppt-gcal."""
import os
import argparse, json, os, sys
import openpyxl

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0

def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {name}: {detail_str}")

def safe_float(val, default=None):
    try:
        if val is None: return default
        return float(str(val).replace(",", "").replace("%", "").replace("$", "").strip())
    except (ValueError, TypeError):
        return default

def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)

def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT
    PASS_COUNT = 0
    FAIL_COUNT = 0

    # Check Market_Sales_Correlation.xlsx
    excel_path = os.path.join(agent_workspace, "Market_Sales_Correlation.xlsx")
    check("Market_Sales_Correlation.xlsx exists", os.path.exists(excel_path))
    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)
        gt_path = os.path.join(groundtruth_workspace, "Market_Sales_Correlation.xlsx")
        gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

        if gt_wb:
            for sheet_name in gt_wb.sheetnames:
                check(f"{sheet_name} sheet exists", sheet_name in wb.sheetnames)
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    gt_ws = gt_wb[sheet_name]
                    # Check headers
                    gt_headers = [str(c.value).strip().lower() if c.value else "" for c in gt_ws[1]]
                    headers = [str(c.value).strip().lower() if c.value else "" for c in ws[1]]
                    for h in gt_headers:
                        if h:
                            check(f"{sheet_name} has {h} column", h in headers, f"headers: {headers[:10]}")
                    # Check row count
                    gt_rows = list(gt_ws.iter_rows(min_row=2, values_only=True))
                    data_rows = list(ws.iter_rows(min_row=2, values_only=True))
                    min_rows = max(1, len(gt_rows) - 2)
                    check(f"{sheet_name} has >= {min_rows} data rows", len(data_rows) >= min_rows, f"got {len(data_rows)}")

                    # Cell value comparison against groundtruth
                    header_map = {h: i for i, h in enumerate(headers)}
                    gt_header_map = {h: i for i, h in enumerate(gt_headers)}
                    for ri in range(min(3, len(gt_rows), len(data_rows))):
                        gt_row = gt_rows[ri]
                        agent_row = data_rows[ri]
                        for ci, gt_h in enumerate(gt_headers):
                            if not gt_h or ci >= len(gt_row):
                                continue
                            gv = gt_row[ci]
                            agent_ci = header_map.get(gt_h)
                            if agent_ci is None or agent_ci >= len(agent_row):
                                continue
                            av = agent_row[agent_ci]
                            gf = safe_float(gv)
                            af = safe_float(av)
                            if gf is not None and af is not None:
                                tol = max(0.5, abs(gf) * 0.15)
                                check(f"{sheet_name} R{ri+2} {gt_h} ~{gf:.1f}",
                                      abs(gf - af) <= tol, f"got {af}")
                            elif gv is not None and av is not None:
                                gs = str(gv).strip().lower()
                                avs = str(av).strip().lower()
                                if gs:
                                    check(f"{sheet_name} R{ri+2} {gt_h} text",
                                          gs == avs or gs in avs or avs in gs,
                                          f"expected {gs[:50]}, got {avs[:50]}")

    # Check Market_Sales_Strategy.pptx
    pptx_path = os.path.join(agent_workspace, "Market_Sales_Strategy.pptx")
    check("Market_Sales_Strategy.pptx exists", os.path.exists(pptx_path))
    if os.path.exists(pptx_path):
        from pptx import Presentation
        prs = Presentation(pptx_path)
        gt_pptx_path = os.path.join(groundtruth_workspace, "Market_Sales_Strategy.pptx")
        if os.path.exists(gt_pptx_path):
            gt_prs = Presentation(gt_pptx_path)
            gt_slide_count = len(gt_prs.slides)
            check(f"Market_Sales_Strategy.pptx has >= {gt_slide_count} slides", len(prs.slides) >= gt_slide_count, f"got {len(prs.slides)} slides")
            # Compare slide titles
            gt_titles = []
            for s in gt_prs.slides:
                if s.shapes.title:
                    gt_titles.append(s.shapes.title.text.strip().lower())
            agent_titles = []
            for s in prs.slides:
                if s.shapes.title:
                    agent_titles.append(s.shapes.title.text.strip().lower())
            for gt in gt_titles:
                found = any(gt in at or at in gt for at in agent_titles)
                check(f"Market_Sales_Strategy.pptx has slide \"{gt[:40]}\"", found, f"agent titles: {agent_titles[:5]}")
        else:
            check("Market_Sales_Strategy.pptx has >= 3 slides", len(prs.slides) >= 3, f"got {len(prs.slides)} slides")

    # Check Python script exists (terminal usage)
    py_files = [f for f in os.listdir(agent_workspace) if f.endswith(".py")]
    check("Python analysis script exists", len(py_files) >= 1, f"found: {py_files}")

    # Database checks
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute("SELECT summary, start_datetime FROM gcal.events WHERE summary ILIKE '%market%'")
        event_row = cur.fetchone()
        check("Calendar event with correct summary", event_row is not None, "no matching event found")
        # Reverse verification: noise events should not match task keyword
        cur.execute("SELECT COUNT(*) FROM gcal.events WHERE summary ILIKE '%standup%' OR summary ILIKE '%lunch%'")
        noise_events = cur.fetchone()[0]
        check("Noise events exist (not deleted by agent)", noise_events >= 1, f"noise events: {noise_events}")
        conn.close()
    except Exception as e:
        check("DB checks", False, str(e))

    return FAIL_COUNT == 0, f"Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()