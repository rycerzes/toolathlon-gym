"""
Evaluation for academic-presentation-builder task.

Checks:
1. LLM_Reasoning_Review.docx: headings, summary table, required keywords
2. LLM_Reasoning_Slides.pptx: slide count, title, key terms
3. word_count.txt: exists with two lines containing word counts
"""
import os
import sys
import json
import re
from argparse import ArgumentParser
from datetime import datetime

from docx import Document
from pptx import Presentation

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)



def check_word_doc(agent_workspace, gt_data):
    """Check the Word document for required structure and content."""
    passed = 0
    total = 0
    filename = gt_data["review_doc"]["filename"]
    doc_path = os.path.join(agent_workspace, filename)

    # Check file exists
    total += 1
    if not os.path.exists(doc_path):
        print(f"  FAIL: {filename} not found at {doc_path}")
        return passed, total

    passed += 1
    print(f"  PASS: {filename} exists")

    doc = Document(doc_path)

    # Extract all text and headings
    all_text = []
    headings_found = []
    for para in doc.paragraphs:
        all_text.append(para.text)
        if para.style and para.style.name and "Heading" in para.style.name:
            headings_found.append(para.text.strip())

    full_text = " ".join(all_text).lower()

    # Check required headings
    required_headings = gt_data["review_doc"]["required_headings"]
    for heading in required_headings:
        total += 1
        found = any(heading.lower() in h.lower() for h in headings_found)
        if found:
            passed += 1
            print(f"  PASS: Heading '{heading}' found")
        else:
            print(f"  FAIL: Heading '{heading}' not found. Found headings: {headings_found}")

    # Check summary table
    total += 1
    tables = doc.tables
    if len(tables) == 0:
        print("  FAIL: No tables found in document")
    else:
        table = tables[0]
        # Check row count (header + data rows)
        expected_data_rows = gt_data["review_doc"]["required_table_rows"]
        # Count non-empty data rows
        data_rows = 0
        for row in table.rows[1:]:  # skip header
            cell_text = " ".join(cell.text.strip() for cell in row.cells)
            if cell_text.strip():
                data_rows += 1

        if data_rows >= expected_data_rows:
            passed += 1
            print(f"  PASS: Summary table has {data_rows} data rows (expected >= {expected_data_rows})")
        else:
            print(f"  FAIL: Summary table has {data_rows} data rows (expected >= {expected_data_rows})")

    # Check table columns
    total += 1
    if len(tables) > 0:
        header_cells = [cell.text.strip().lower() for cell in tables[0].rows[0].cells]
        expected_cols = [c.lower() for c in gt_data["review_doc"]["required_table_columns"]]
        cols_found = sum(1 for ec in expected_cols if any(ec in hc for hc in header_cells))
        if cols_found >= len(expected_cols):
            passed += 1
            print(f"  PASS: All {len(expected_cols)} required table columns found")
        else:
            print(f"  FAIL: Only {cols_found}/{len(expected_cols)} required columns found. Headers: {header_cells}")
    else:
        print("  FAIL: No tables to check columns")

    # Check required mentions (keywords)
    required_mentions = gt_data["review_doc"]["required_mentions"]
    for mention in required_mentions:
        total += 1
        if mention.lower() in full_text:
            passed += 1
            print(f"  PASS: Keyword '{mention}' found in document")
        else:
            print(f"  FAIL: Keyword '{mention}' not found in document text")

    return passed, total


def check_pptx(agent_workspace, gt_data):
    """Check the PowerPoint presentation for required structure and content."""
    passed = 0
    total = 0
    filename = gt_data["slides"]["filename"]
    pptx_path = os.path.join(agent_workspace, filename)

    # Check file exists
    total += 1
    if not os.path.exists(pptx_path):
        print(f"  FAIL: {filename} not found at {pptx_path}")
        return passed, total

    passed += 1
    print(f"  PASS: {filename} exists")

    prs = Presentation(pptx_path)

    # Check slide count
    total += 1
    slide_count = len(prs.slides)
    required_count = gt_data["slides"]["required_slide_count"]
    if slide_count >= required_count:
        passed += 1
        print(f"  PASS: Presentation has {slide_count} slides (required >= {required_count})")
    else:
        print(f"  FAIL: Presentation has {slide_count} slides (required >= {required_count})")

    # Extract all text from slides
    all_slide_text = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    slide_texts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        slide_texts.append(cell.text)
        all_slide_text.append(" ".join(slide_texts))

    full_slides_text = " ".join(all_slide_text).lower()

    # Check title slide contains required title
    total += 1
    required_title = gt_data["slides"]["required_title"].lower()
    first_slide_text = all_slide_text[0].lower() if all_slide_text else ""
    # Also check across all slides in case title is on a different slide
    if required_title in first_slide_text or required_title in full_slides_text:
        passed += 1
        print(f"  PASS: Required title '{gt_data['slides']['required_title']}' found")
    else:
        # Try partial match
        title_words = required_title.split()
        match_count = sum(1 for w in title_words if w in full_slides_text)
        if match_count >= len(title_words) - 1:
            passed += 1
            print(f"  PASS: Required title approximately matched ({match_count}/{len(title_words)} words)")
        else:
            print(f"  FAIL: Required title '{gt_data['slides']['required_title']}' not found. First slide: {first_slide_text[:100]}")

    # Check required mentions across slides
    required_mentions = gt_data["slides"]["required_mentions"]
    for mention in required_mentions:
        total += 1
        if mention.lower() in full_slides_text:
            passed += 1
            print(f"  PASS: Keyword '{mention}' found in slides")
        else:
            print(f"  FAIL: Keyword '{mention}' not found in slides")

    return passed, total


def check_word_count(agent_workspace):
    """Check word_count.txt exists and has valid content."""
    passed = 0
    total = 0

    wc_path = os.path.join(agent_workspace, "word_count.txt")

    # Check file exists
    total += 1
    if not os.path.exists(wc_path):
        print(f"  FAIL: word_count.txt not found at {wc_path}")
        return passed, total

    passed += 1
    print("  PASS: word_count.txt exists")

    with open(wc_path, "r") as f:
        content = f.read().strip()

    lines = [line.strip() for line in content.split("\n") if line.strip()]

    # Check at least 2 lines
    total += 1
    if len(lines) >= 2:
        passed += 1
        print(f"  PASS: word_count.txt has {len(lines)} lines (expected >= 2)")
    else:
        print(f"  FAIL: word_count.txt has {len(lines)} lines (expected >= 2)")

    # Check each line has a word count > 0
    for line in lines:
        total += 1
        # Try to extract a number from the line
        numbers = re.findall(r'\d+', line)
        if numbers and any(int(n) > 0 for n in numbers):
            passed += 1
            print(f"  PASS: Line contains valid word count: {line}")
        else:
            print(f"  FAIL: Line does not contain valid word count: {line}")

    # Check both filenames are mentioned
    total += 1
    if "LLM_Reasoning_Review" in content and "LLM_Reasoning_Slides" in content:
        passed += 1
        print("  PASS: Both filenames referenced in word_count.txt")
    else:
        print("  FAIL: Not both filenames found in word_count.txt")

    return passed, total


def main(args):
    gt_path = os.path.join(args.groundtruth_workspace, "expected_results.json")
    if not os.path.exists(gt_path):
        print(f"FAIL: expected_results.json not found at {gt_path}")
        sys.exit(1)

    with open(gt_path, "r") as f:
        gt_data = json.load(f)

    total_passed = 0
    total_checks = 0

    # Check 1: Word document
    print("--- Check 1: Word Document (LLM_Reasoning_Review.docx) ---")
    p, t = check_word_doc(args.agent_workspace, gt_data)
    print(f"  Word Doc: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Check 2: PowerPoint
    print("\n--- Check 2: PowerPoint (LLM_Reasoning_Slides.pptx) ---")
    p, t = check_pptx(args.agent_workspace, gt_data)
    print(f"  PowerPoint: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Check 3: Word count file
    print("\n--- Check 3: Word Count File (word_count.txt) ---")
    p, t = check_word_count(args.agent_workspace)
    print(f"  Word Count: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Overall
    if total_checks == 0:
        print("\nFAIL: No checks were performed.")
        accuracy = 0.0
    else:
        accuracy = total_passed / total_checks * 100
        print(f"\nOverall: {total_passed}/{total_checks} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": total_passed,
        "total_checks": total_checks,
        "accuracy": accuracy,
        "timestamp": datetime.now().isoformat(),
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)
        print(f"Report saved to {args.res_log_file}")

    if accuracy >= 80:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    main(args)
