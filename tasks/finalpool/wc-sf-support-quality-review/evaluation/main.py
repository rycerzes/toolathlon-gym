"""Evaluation for wc-sf-support-quality-review."""
import argparse
import os
import sys

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)



def nums_close(a, b, abs_tol=1.0, rel_tol=0.05):
    try:
        a, b = float(a), float(b)
    except (TypeError, ValueError):
        return False
    if abs(a - b) <= abs_tol:
        return True
    if b != 0 and abs(a - b) / abs(b) <= rel_tol:
        return True
    return False


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def check_excel(agent_workspace):
    errors = []
    import openpyxl

    path = os.path.join(agent_workspace, "CX_Quality_Review.xlsx")
    if not os.path.exists(path):
        return ["CX_Quality_Review.xlsx not found"]
    try:
        wb = openpyxl.load_workbook(path, data_only=True)

        # Sheet 1: Order Issues
        rows = load_sheet_rows(wb, "Order Issues")
        if rows is None:
            errors.append("Sheet 'Order Issues' not found")
        else:
            data_rows = [r for r in rows[1:] if r and r[0] is not None]
            if len(data_rows) < 4:
                errors.append(
                    f"Order Issues has {len(data_rows)} rows, expected >= 4"
                )
            # Check Electronics category
            elec_rows = [
                r for r in data_rows if r[0] and "electronics" in str(r[0]).lower()
            ]
            if not elec_rows:
                errors.append("Electronics category not found in Order Issues")
            else:
                if len(elec_rows[0]) > 1 and elec_rows[0][1]:
                    if not nums_close(elec_rows[0][1], 83, abs_tol=10):
                        errors.append(
                            f"Electronics total orders {elec_rows[0][1]}, expected ~83"
                        )

            # Check issue rate column exists with reasonable values
            has_rate = False
            for r in data_rows:
                if len(r) > 3 and r[3] is not None:
                    try:
                        rate = float(r[3])
                        if 0 < rate < 100:
                            has_rate = True
                            break
                    except (ValueError, TypeError):
                        pass
            if not has_rate:
                errors.append("Issue rate column missing or has no valid values")

        # Sheet 2: Support Metrics
        rows2 = load_sheet_rows(wb, "Support Metrics")
        if rows2 is None:
            errors.append("Sheet 'Support Metrics' not found")
        else:
            data_rows2 = [r for r in rows2[1:] if r and r[0] is not None]
            if len(data_rows2) < 3:
                errors.append(
                    f"Support Metrics has {len(data_rows2)} rows, expected 3"
                )
            # Check High priority row
            high_rows = [
                r for r in data_rows2 if r[0] and "high" in str(r[0]).lower()
            ]
            if not high_rows:
                errors.append("High priority row not found in Support Metrics")
            else:
                if len(high_rows[0]) > 2 and high_rows[0][2]:
                    if not nums_close(high_rows[0][2], 6.23, abs_tol=1.0):
                        errors.append(
                            f"High priority avg response time {high_rows[0][2]}, expected ~6.23"
                        )

        # Sheet 3: Cross Reference
        rows3 = load_sheet_rows(wb, "Cross Reference")
        if rows3 is None:
            errors.append("Sheet 'Cross Reference' not found")
        else:
            data_rows3 = [r for r in rows3[1:] if r and r[0] is not None]
            if len(data_rows3) < 5:
                errors.append(
                    f"Cross Reference has {len(data_rows3)} rows, expected >= 5"
                )
            # Check Bug issue type exists
            bug_rows = [
                r for r in data_rows3 if r[0] and "bug" in str(r[0]).lower()
            ]
            if not bug_rows:
                errors.append("Bug issue type not found in Cross Reference")

        # Sheet 4: Executive Summary
        rows4 = load_sheet_rows(wb, "Executive Summary")
        if rows4 is None:
            errors.append("Sheet 'Executive Summary' not found")
        else:
            data_rows4 = [r for r in rows4[1:] if r and r[0] is not None]
            if len(data_rows4) < 5:
                errors.append(
                    f"Executive Summary has {len(data_rows4)} rows, expected >= 5"
                )

    except Exception as e:
        errors.append(f"Error reading Excel: {e}")
    return errors


def check_pptx(agent_workspace):
    errors = []
    from pptx import Presentation

    path = os.path.join(agent_workspace, "Quality_Review.pptx")
    if not os.path.exists(path):
        return ["Quality_Review.pptx not found"]
    try:
        prs = Presentation(path)
        slide_count = len(prs.slides)
        if slide_count < 6:
            errors.append(f"PowerPoint has {slide_count} slides, expected >= 6")

        # Check content across slides
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        if "benchmark" not in all_text and "target" not in all_text:
            errors.append("PowerPoint does not mention benchmarks or targets")
        if "issue" not in all_text and "order" not in all_text:
            errors.append("PowerPoint does not discuss order issues")
        if "support" not in all_text and "ticket" not in all_text:
            errors.append("PowerPoint does not discuss support metrics")
        if "recommend" not in all_text:
            errors.append("PowerPoint does not include recommendations")

    except Exception as e:
        errors.append(f"Error reading PowerPoint: {e}")
    return errors


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    agent_ws = args.agent_workspace or os.path.join(
        os.path.dirname(__file__), "..", "groundtruth_workspace"
    )

    all_errors = []

    print("  Checking Excel file...")
    errs = check_excel(agent_ws)
    if errs:
        all_errors.extend(errs)
        for e in errs[:5]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking PowerPoint...")
    errs = check_pptx(agent_ws)
    if errs:
        all_errors.extend(errs)
        for e in errs[:3]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
