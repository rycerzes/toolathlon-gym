"""Evaluation for canvas-semester-summary-ppt."""
import argparse
import os
import sys
import openpyxl
from pptx import Presentation


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")
    all_errors = []

    # ---- Check Excel ----
    agent_excel = os.path.join(args.agent_workspace, "Semester_Summary.xlsx")
    gt_excel = os.path.join(gt_dir, "Semester_Summary.xlsx")

    if not os.path.exists(agent_excel):
        print(f"FAIL: Agent output Semester_Summary.xlsx not found")
        sys.exit(1)
    if not os.path.exists(gt_excel):
        print(f"FAIL: Groundtruth Semester_Summary.xlsx not found")
        sys.exit(1)

    agent_wb = openpyxl.load_workbook(agent_excel, data_only=True)
    gt_wb = openpyxl.load_workbook(gt_excel, data_only=True)

    # Check Course Overview sheet
    print("  Checking Course Overview...")
    a_rows = load_sheet_rows(agent_wb, "Course Overview")
    g_rows = load_sheet_rows(gt_wb, "Course Overview")
    if a_rows is None:
        all_errors.append("Sheet 'Course Overview' not found in agent output")
    elif g_rows is None:
        all_errors.append("Sheet 'Course Overview' not found in groundtruth")
    else:
        a_data = a_rows[1:] if len(a_rows) > 1 else []
        g_data = g_rows[1:] if len(g_rows) > 1 else []
        if len(a_data) != len(g_data):
            all_errors.append(f"Course Overview row count: agent={len(a_data)}, expected={len(g_data)}")

        a_lookup = {}
        for row in a_data:
            if row and row[0] is not None:
                a_lookup[str(row[0]).strip().lower()] = row
        for g_row in g_data:
            if not g_row or g_row[0] is None:
                continue
            key = str(g_row[0]).strip().lower()
            a_row = a_lookup.get(key)
            if a_row is None:
                all_errors.append(f"Missing course: {g_row[0]}")
                continue
            # Col 1: Course Code (string match)
            if len(a_row) > 1 and len(g_row) > 1:
                if not str_match(a_row[1], g_row[1]):
                    all_errors.append(f"{key}.Course Code: {a_row[1]} vs {g_row[1]}")
            # Col 2: Total Students
            if len(a_row) > 2 and len(g_row) > 2:
                if not num_close(a_row[2], g_row[2], 5):
                    all_errors.append(f"{key}.Total Students: {a_row[2]} vs {g_row[2]}")
            # Col 3: Enrollments
            if len(a_row) > 3 and len(g_row) > 3:
                if not num_close(a_row[3], g_row[3], 5):
                    all_errors.append(f"{key}.Enrollments: {a_row[3]} vs {g_row[3]}")
            # Col 4: Assignments
            if len(a_row) > 4 and len(g_row) > 4:
                if not num_close(a_row[4], g_row[4], 1):
                    all_errors.append(f"{key}.Assignments: {a_row[4]} vs {g_row[4]}")
            # Col 5: Avg Points Possible
            if len(a_row) > 5 and len(g_row) > 5:
                if not num_close(a_row[5], g_row[5], 2.0):
                    all_errors.append(f"{key}.Avg Points: {a_row[5]} vs {g_row[5]}")

        if not any("Course Overview" in e or "Missing course" in e for e in all_errors):
            print("    PASS")

    # Check Summary sheet
    print("  Checking Summary...")
    a_rows = load_sheet_rows(agent_wb, "Summary")
    g_rows = load_sheet_rows(gt_wb, "Summary")
    if a_rows is None:
        all_errors.append("Sheet 'Summary' not found in agent output")
    elif g_rows is None:
        all_errors.append("Sheet 'Summary' not found in groundtruth")
    else:
        a_data = a_rows[1:] if len(a_rows) > 1 else []
        g_data = g_rows[1:] if len(g_rows) > 1 else []

        a_lookup = {}
        for row in a_data:
            if row and row[0] is not None:
                a_lookup[str(row[0]).strip().lower()] = row
        for g_row in g_data:
            if not g_row or g_row[0] is None:
                continue
            key = str(g_row[0]).strip().lower()
            a_row = a_lookup.get(key)
            if a_row is None:
                all_errors.append(f"Missing metric: {g_row[0]}")
                continue
            if len(a_row) > 1 and len(g_row) > 1:
                if not num_close(a_row[1], g_row[1], 2.0):
                    all_errors.append(f"{key}.Value: {a_row[1]} vs {g_row[1]}")

        if not any("Summary" in e or "Missing metric" in e for e in all_errors):
            print("    PASS")

    # ---- Check PPT ----
    agent_ppt = os.path.join(args.agent_workspace, "Semester_Summary.pptx")
    if not os.path.exists(agent_ppt):
        all_errors.append("Agent output Semester_Summary.pptx not found")
    else:
        print("  Checking Semester_Summary.pptx...")
        prs = Presentation(agent_ppt)
        slides = list(prs.slides)

        # Minimum: title + overview + 7 courses + takeaways = 10
        if len(slides) < 10:
            all_errors.append(f"PPT has {len(slides)} slides, expected at least 10")
        else:
            # Check title slide
            title_text = ""
            for shape in slides[0].shapes:
                if shape.has_text_frame:
                    title_text += shape.text_frame.text.lower() + " "
            if "fall 2014" not in title_text:
                all_errors.append(f"Title slide missing 'Fall 2014'. Found: {title_text[:100]}")

            # Check last slide mentions key takeaways
            last_text = ""
            for shape in slides[-1].shapes:
                if shape.has_text_frame:
                    last_text += shape.text_frame.text.lower() + " "
            if "takeaway" not in last_text and "key" not in last_text and "summary" not in last_text and "conclusion" not in last_text:
                all_errors.append("Last slide missing takeaways/summary content")

            # Check that course names appear in presentation
            all_ppt_text = ""
            for slide in slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        all_ppt_text += shape.text_frame.text.lower() + " "

            expected_courses = ["applied analytics", "biochemistry", "creative computing",
                                "data-driven design", "environmental economics",
                                "foundations of finance", "global governance"]
            for course in expected_courses:
                if course not in all_ppt_text:
                    all_errors.append(f"PPT missing course: {course}")

            # Check largest/smallest mentioned in last slide
            if "creative computing" not in last_text and "2498" not in last_text:
                all_errors.append("Last slide should mention Creative Computing (largest by students)")
            if "applied analytics" not in last_text and "365" not in last_text:
                all_errors.append("Last slide should mention Applied Analytics (smallest by students)")

        if not any("PPT" in e or "ppt" in e.lower() or "slide" in e.lower() for e in all_errors):
            print("    PASS")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:15]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
