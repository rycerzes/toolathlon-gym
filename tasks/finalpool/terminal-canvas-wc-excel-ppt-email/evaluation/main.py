import argparse
import json
import os
import sys
import openpyxl
import psycopg2
from pptx import Presentation

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def num_close(a, b, tol=2.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except Exception:
        return False


def load_groundtruth():
    """Compute expected values dynamically from Canvas and WC databases."""
    try:
        return _compute_groundtruth_from_db()
    except Exception:
        # Fallback to static JSON if DB query fails
        gt_path = os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
            "groundtruth_workspace",
            "groundtruth_data.json",
        )
        with open(gt_path) as f:
            return json.load(f)


def _compute_groundtruth_from_db():
    """Query canvas and wc schemas to compute expected values dynamically."""
    import csv as _csv

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        # Read student_bookstore_registry.csv to get mappings
        registry_path = os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
            "initial_workspace",
            "student_bookstore_registry.csv",
        )
        mappings = []
        with open(registry_path) as f:
            reader = _csv.DictReader(f)
            for row in reader:
                mappings.append({
                    "student_id": int(row["student_id"]),
                    "bookstore_customer_id": int(row["bookstore_customer_id"]),
                    "department": row.get("department", ""),
                })

        student_ids = [m["student_id"] for m in mappings]
        sid_to_bcid = {m["student_id"]: m["bookstore_customer_id"] for m in mappings}

        # Get enrollments for courses 1 and 2
        cur.execute("""
            SELECT DISTINCT user_id FROM canvas.enrollments
            WHERE course_id IN (1, 2)
        """)
        enrolled_ids = {r[0] for r in cur.fetchall()}

        # Filter to enrolled students in registry
        matched_sids = [sid for sid in student_ids if sid in enrolled_ids]

        # Get student names
        if matched_sids:
            cur.execute(
                "SELECT id, name FROM canvas.users WHERE id = ANY(%s)",
                (matched_sids,),
            )
            sid_to_name = dict(cur.fetchall())
        else:
            sid_to_name = {}

        # Get average scores per student for courses 1,2
        cur.execute("""
            SELECT s.user_id, AVG(s.score)
            FROM canvas.submissions s
            JOIN canvas.assignments a ON s.assignment_id = a.id
            WHERE a.course_id IN (1, 2)
              AND s.score IS NOT NULL
              AND s.user_id = ANY(%s)
            GROUP BY s.user_id
        """, (matched_sids,))
        sid_to_avg_score = {r[0]: round(float(r[1]), 2) for r in cur.fetchall()}

        # Get submission counts
        cur.execute("""
            SELECT s.user_id, COUNT(*)
            FROM canvas.submissions s
            JOIN canvas.assignments a ON s.assignment_id = a.id
            WHERE a.course_id IN (1, 2)
              AND s.score IS NOT NULL
              AND s.user_id = ANY(%s)
            GROUP BY s.user_id
        """, (matched_sids,))
        sid_to_sub_count = {r[0]: int(r[1]) for r in cur.fetchall()}

        # Get electronics spending from WC
        # Find Electronics/Cameras category IDs
        cur.execute("""
            SELECT id, name FROM wc.product_categories
            WHERE lower(name) IN ('electronics', 'cameras')
        """)
        elec_cats = {r[0]: r[1] for r in cur.fetchall()}
        elec_cat_names = [name.lower() for name in elec_cats.values()]

        # Get orders for matched bookstore customer IDs and compute electronics spend
        bcids = [sid_to_bcid[sid] for sid in matched_sids if sid in sid_to_bcid]
        bcid_spend = {}
        bcid_order_count = {}
        bcid_item_count = {}

        if bcids:
            cur.execute("""
                SELECT id, customer_id, line_items FROM wc.orders
                WHERE customer_id = ANY(%s)
            """, (bcids,))
            for order_id, cust_id, line_items in cur.fetchall():
                if not line_items:
                    continue
                items = line_items if isinstance(line_items, list) else json.loads(line_items)
                has_electronics = False
                for item in items:
                    # Check if product is in electronics/cameras category
                    product_id = item.get("product_id")
                    # Check categories in line item
                    cats = []
                    if "categories" in item:
                        cats = [c.get("name", "").lower() if isinstance(c, dict) else str(c).lower()
                                for c in item["categories"]]
                    # Also check via product_categories join
                    is_elec = any(c in elec_cat_names for c in cats)
                    if not is_elec and product_id:
                        # Look up product categories
                        cur.execute("""
                            SELECT categories FROM wc.products WHERE id = %s
                        """, (product_id,))
                        prod_row = cur.fetchone()
                        if prod_row and prod_row[0]:
                            prod_cats = prod_row[0] if isinstance(prod_row[0], list) else json.loads(str(prod_row[0]))
                            for pc in prod_cats:
                                if isinstance(pc, dict) and pc.get("name", "").lower() in elec_cat_names:
                                    is_elec = True
                                    break
                    if is_elec:
                        has_electronics = True
                        total = float(item.get("total", 0))
                        qty = int(item.get("quantity", 0))
                        bcid_spend[cust_id] = bcid_spend.get(cust_id, 0) + total
                        bcid_item_count[cust_id] = bcid_item_count.get(cust_id, 0) + qty
                if has_electronics:
                    bcid_order_count[cust_id] = bcid_order_count.get(cust_id, 0) + 1

        # Build matched_data
        matched_data = []
        for sid in matched_sids:
            if sid not in sid_to_avg_score:
                continue
            bcid = sid_to_bcid.get(sid)
            if bcid is None or bcid not in bcid_spend:
                continue
            matched_data.append({
                "student_name": sid_to_name.get(sid, f"Student_{sid}"),
                "student_id": sid,
                "avg_score": sid_to_avg_score[sid],
                "submission_count": sid_to_sub_count.get(sid, 0),
                "bookstore_customer_id": bcid,
                "total_electronics_spend": round(bcid_spend.get(bcid, 0), 2),
                "electronics_order_count": bcid_order_count.get(bcid, 0),
                "electronics_item_count": bcid_item_count.get(bcid, 0),
            })

        num_matched = len(matched_data)

        # Compute correlation
        if num_matched >= 2:
            scores = [m["avg_score"] for m in matched_data]
            spends = [m["total_electronics_spend"] for m in matched_data]
            mean_s = sum(scores) / len(scores)
            mean_sp = sum(spends) / len(spends)
            cov = sum((s - mean_s) * (sp - mean_sp) for s, sp in zip(scores, spends))
            std_s = (sum((s - mean_s) ** 2 for s in scores)) ** 0.5
            std_sp = (sum((sp - mean_sp) ** 2 for sp in spends)) ** 0.5
            correlation = round(cov / (std_s * std_sp), 4) if std_s > 0 and std_sp > 0 else 0
        else:
            scores = [m["avg_score"] for m in matched_data]
            spends = [m["total_electronics_spend"] for m in matched_data]
            mean_s = sum(scores) / len(scores) if scores else 0
            mean_sp = sum(spends) / len(spends) if spends else 0
            correlation = 0

        return {
            "matched_data": matched_data,
            "num_matched": num_matched,
            "correlation": correlation,
            "mean_score": round(mean_s, 2),
            "mean_spend": round(mean_sp, 2),
        }
    finally:
        cur.close()
        conn.close()


def check_excel(workspace, gt):
    print("\n=== Excel Checks ===")
    xlsx_path = os.path.join(workspace, "Student_Purchase_Analysis.xlsx")
    check("Excel file exists", os.path.exists(xlsx_path), xlsx_path)
    if not os.path.exists(xlsx_path):
        return

    wb = openpyxl.load_workbook(xlsx_path)
    sheets = wb.sheetnames

    # Sheet existence
    check("Sheet Student_Performance exists", "Student_Performance" in sheets, sheets)
    check("Sheet Purchase_Summary exists", "Purchase_Summary" in sheets, sheets)
    check("Sheet Correlation_Analysis exists", "Correlation_Analysis" in sheets, sheets)
    check("Sheet Recommendations exists", "Recommendations" in sheets, sheets)

    # Student_Performance content
    if "Student_Performance" in sheets:
        ws = wb["Student_Performance"]
        headers = [ws.cell(1, c).value for c in range(1, ws.max_column + 1)]
        check(
            "Student_Performance has Student_Name column",
            "Student_Name" in headers,
            headers,
        )
        check(
            "Student_Performance has Average_Score column",
            "Average_Score" in headers,
            headers,
        )
        data_rows = ws.max_row - 1
        check(
            "Student_Performance has correct row count",
            data_rows == gt["num_matched"],
            f"got {data_rows}, expected {gt['num_matched']}",
        )

        # Spot check a few scores
        if "Average_Score" in headers:
            score_col = headers.index("Average_Score") + 1
            name_col = headers.index("Student_Name") + 1
            gt_by_name = {m["student_name"]: m["avg_score"] for m in gt["matched_data"]}
            checked = 0
            for row in range(2, ws.max_row + 1):
                name = ws.cell(row, name_col).value
                score = ws.cell(row, score_col).value
                if name in gt_by_name:
                    check(
                        f"Score for {name}",
                        num_close(score, gt_by_name[name], 1.0),
                        f"got {score}, expected {gt_by_name[name]}",
                    )
                    checked += 1
                    if checked >= 3:
                        break

    # Purchase_Summary content
    if "Purchase_Summary" in sheets:
        ws = wb["Purchase_Summary"]
        headers = [ws.cell(1, c).value for c in range(1, ws.max_column + 1)]
        check(
            "Purchase_Summary has Total_Electronics_Spend",
            "Total_Electronics_Spend" in headers,
            headers,
        )
        check(
            "Purchase_Summary has Electronics_Order_Count",
            "Electronics_Order_Count" in headers,
            headers,
        )

        # Spot check spending
        if "Total_Electronics_Spend" in headers and "Student_Name" in headers:
            spend_col = headers.index("Total_Electronics_Spend") + 1
            name_col = headers.index("Student_Name") + 1
            gt_by_name = {
                m["student_name"]: m["total_electronics_spend"]
                for m in gt["matched_data"]
            }
            checked = 0
            for row in range(2, ws.max_row + 1):
                name = ws.cell(row, name_col).value
                spend = ws.cell(row, spend_col).value
                if name in gt_by_name:
                    check(
                        f"Spend for {name}",
                        num_close(spend, gt_by_name[name], 5.0),
                        f"got {spend}, expected {gt_by_name[name]}",
                    )
                    checked += 1
                    if checked >= 3:
                        break

    # Correlation_Analysis content
    if "Correlation_Analysis" in sheets:
        ws = wb["Correlation_Analysis"]
        metrics = {}
        for row in range(2, ws.max_row + 1):
            key = ws.cell(row, 1).value
            val = ws.cell(row, 2).value
            if key:
                metrics[key] = val

        check(
            "Correlation has Pearson_Correlation",
            "Pearson_Correlation" in metrics,
            list(metrics.keys()),
        )
        if "Pearson_Correlation" in metrics:
            check(
                "Pearson_Correlation value",
                num_close(metrics["Pearson_Correlation"], gt["correlation"], 0.1),
                f"got {metrics['Pearson_Correlation']}, expected {gt['correlation']}",
            )
        check(
            "Correlation has Mean_Academic_Score",
            "Mean_Academic_Score" in metrics,
            list(metrics.keys()),
        )
        if "Mean_Academic_Score" in metrics:
            check(
                "Mean_Academic_Score value",
                num_close(metrics["Mean_Academic_Score"], gt["mean_score"], 2.0),
                f"got {metrics['Mean_Academic_Score']}, expected {gt['mean_score']}",
            )
        check(
            "Correlation has Mean_Electronics_Spend",
            "Mean_Electronics_Spend" in metrics,
            list(metrics.keys()),
        )
        if "Mean_Electronics_Spend" in metrics:
            check(
                "Mean_Electronics_Spend value",
                num_close(metrics["Mean_Electronics_Spend"], gt["mean_spend"], 20.0),
                f"got {metrics['Mean_Electronics_Spend']}, expected {gt['mean_spend']}",
            )
        check(
            "Correlation has Recommendation",
            "Recommendation" in metrics,
            list(metrics.keys()),
        )
        if "Recommendation" in metrics:
            rec = str(metrics["Recommendation"]).lower()
            expected_keyword = "discount" if gt["correlation"] < 0 else "expand"
            check(
                "Recommendation matches correlation direction",
                expected_keyword in rec,
                f"got '{metrics['Recommendation'][:80]}', expected keyword '{expected_keyword}'",
            )

    # Recommendations sheet
    if "Recommendations" in sheets:
        ws = wb["Recommendations"]
        check(
            "Recommendations has at least 3 data rows",
            ws.max_row >= 4,
            f"rows={ws.max_row}",
        )


def check_pptx(workspace, gt):
    print("\n=== PowerPoint Checks ===")
    pptx_path = os.path.join(workspace, "Purchase_Behavior_Presentation.pptx")
    check("PPTX file exists", os.path.exists(pptx_path), pptx_path)
    if not os.path.exists(pptx_path):
        return

    prs = Presentation(pptx_path)
    slides = prs.slides
    check("PPTX has 5 slides", len(slides) == 5, f"got {len(slides)}")

    if len(slides) >= 1:
        title = slides[0].shapes.title
        if title:
            check(
                "Slide 1 title contains 'Purchase Behavior' or 'Electronics'",
                "purchase" in title.text.lower() or "electronics" in title.text.lower(),
                title.text,
            )

    expected_titles = [
        "Study Overview",
        "Student Performance Summary",
        "Correlation Findings",
        "Recommendations",
    ]
    for i, exp in enumerate(expected_titles):
        if i + 1 < len(slides):
            slide = slides[i + 1]
            t = slide.shapes.title
            if t:
                check(
                    f"Slide {i+2} title contains '{exp.split()[0]}'",
                    exp.split()[0].lower() in t.text.lower(),
                    t.text,
                )

    # Check slide content mentions key data
    all_text = ""
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + " "

    check(
        "PPTX mentions number of matched students",
        str(gt["num_matched"]) in all_text,
        f"looking for {gt['num_matched']}",
    )
    check(
        "PPTX mentions correlation coefficient",
        str(gt["correlation"]) in all_text
        or str(abs(gt["correlation"])) in all_text,
        f"looking for {gt['correlation']}",
    )


def check_emails(gt):
    print("\n=== Email Checks ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        sent_folder = "(SELECT id FROM email.folders WHERE name = 'Sent' LIMIT 1)"

        # Check bookstore manager email
        cur.execute(
            f"""SELECT subject, body_text, to_addr FROM email.messages
            WHERE folder_id = {sent_folder}
            AND to_addr::text LIKE '%bookstore_manager@university.edu%'"""
        )
        mgr_emails = cur.fetchall()
        check(
            "Email sent to bookstore_manager@university.edu",
            len(mgr_emails) >= 1,
            f"found {len(mgr_emails)}",
        )
        if mgr_emails:
            subj = mgr_emails[0][0] or ""
            body = mgr_emails[0][1] or ""
            check(
                "Manager email subject mentions purchase/behavior/report",
                any(
                    kw in subj.lower()
                    for kw in ["purchase", "behavior", "report", "electronics"]
                ),
                subj,
            )
            check(
                "Manager email body mentions correlation or matched students",
                "correlation" in body.lower()
                or "matched" in body.lower()
                or str(gt["num_matched"]) in body,
                body[:100],
            )

        # Check academic affairs email
        cur.execute(
            f"""SELECT subject, body_text, to_addr FROM email.messages
            WHERE folder_id = {sent_folder}
            AND to_addr::text LIKE '%academic_affairs@university.edu%'"""
        )
        acad_emails = cur.fetchall()
        check(
            "Email sent to academic_affairs@university.edu",
            len(acad_emails) >= 1,
            f"found {len(acad_emails)}",
        )
        if acad_emails:
            subj = acad_emails[0][0] or ""
            body = acad_emails[0][1] or ""
            check(
                "Academic email subject mentions correlation or academic",
                any(
                    kw in subj.lower()
                    for kw in ["correlation", "academic", "performance"]
                ),
                subj,
            )
            check(
                "Academic email body mentions mean score or correlation",
                "correlation" in body.lower()
                or "mean" in body.lower()
                or str(gt["mean_score"]) in body,
                body[:100],
            )
    finally:
        cur.close()
        conn.close()


def check_reverse_validation(workspace):
    print("\n=== Reverse Validation ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        # Check no emails sent to noise recipients
        noise_recipients = [
            "newsletter@university.edu",
            "all-staff@university.edu",
            "alumni@university.edu",
            "admissions@university.edu",
        ]
        for addr in noise_recipients:
            cur.execute(
                "SELECT COUNT(*) FROM email.messages WHERE to_addr::text ILIKE %s",
                (f"%{addr}%",),
            )
            cnt = cur.fetchone()[0]
            check(f"No email sent to noise recipient {addr}", cnt == 0,
                  f"Found {cnt} emails to {addr}")

        # Check Excel does not include non-Electronics categories in Purchase_Summary
        xlsx_path = os.path.join(workspace, "Student_Purchase_Analysis.xlsx")
        if os.path.exists(xlsx_path):
            wb = openpyxl.load_workbook(xlsx_path)
            if "Purchase_Summary" in wb.sheetnames:
                ws = wb["Purchase_Summary"]
                all_text = " ".join(
                    str(ws.cell(r, c).value) for r in range(2, ws.max_row + 1)
                    for c in range(1, ws.max_column + 1) if ws.cell(r, c).value
                ).lower()
                check("Purchase_Summary does not contain Clothing category data",
                      "clothing" not in all_text,
                      "Found 'clothing' in Purchase_Summary")
    except Exception as e:
        check("Reverse validation", False, str(e))
    finally:
        cur.close()
        conn.close()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    gt = load_groundtruth()

    check_excel(args.agent_workspace, gt)
    check_pptx(args.agent_workspace, gt)
    check_emails(gt)
    check_reverse_validation(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total > 0 else 0
    print(f"\nOverall: {PASS_COUNT}/{total} ({accuracy:.1f}%)")
    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
    }
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)
    sys.exit(0 if accuracy >= 70 else 1)


if __name__ == "__main__":
    main()
