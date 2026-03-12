"""
Evaluation script for canvas-enrollment-overview-gsheet-ppt task.

Checks:
1. Google Sheet "Fall 2013 Enrollment Overview" exists with correct enrollment data
2. PowerPoint file Enrollment_Overview_F2013.pptx exists with correct structure
3. Word document Enrollment_Report_F2013.docx exists with correct content
4. Email sent to academic.office@university.edu
"""

import argparse
import json
import os
import sys

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


# ============================================================================
# Check 1: Google Sheet
# ============================================================================

def check_gsheet():
    print("\n=== Checking Google Sheet ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("SELECT id, title FROM gsheet.spreadsheets")
    sheets = cur.fetchall()
    print(f"[check_gsheet] Found {len(sheets)} spreadsheets.")
    record("At least 1 spreadsheet created", len(sheets) >= 1)

    found_sheet = False
    for ss_id, title in sheets:
        if title and ("enrollment" in title.lower() or "fall 2013" in title.lower()):
            found_sheet = True
            record(f"Spreadsheet '{title}' found", True)

            # Get cell data
            cur.execute("""
                SELECT c.row_index, c.col_index, c.value
                FROM gsheet.cells c
                JOIN gsheet.sheets s ON c.sheet_id = s.id
                WHERE c.spreadsheet_id = %s
                ORDER BY c.row_index, c.col_index
            """, (ss_id,))
            cells = cur.fetchall()

            # Build a grid
            grid = {}
            for row_i, col_i, val in cells:
                if row_i not in grid:
                    grid[row_i] = {}
                grid[row_i][col_i] = val

            max_row = max(grid.keys()) if grid else 0
            record("Spreadsheet has at least 7 rows (header + 6 courses)",
                   max_row >= 6,
                   f"Max row index: {max_row}")

            # Check for expected courses in data
            all_vals = [str(v).lower() for row in grid.values() for v in row.values() if v]
            for code in ["fff-2013j", "bbb-2013j", "ddd-2013j"]:
                found_code = any(code in v for v in all_vals)
                record(f"Spreadsheet contains {code.upper()}", found_code)

            # Check for numeric student counts
            numeric_vals = []
            for row in grid.values():
                for v in row.values():
                    try:
                        numeric_vals.append(float(v))
                    except (TypeError, ValueError):
                        pass
            record("Spreadsheet has numeric enrollment data",
                   any(v >= 100 for v in numeric_vals),
                   f"Numeric values: {numeric_vals[:5]}")
            break

    if not found_sheet:
        record("Fall 2013 Enrollment spreadsheet found", False,
               f"Spreadsheets: {[(s[0], s[1]) for s in sheets]}")

    cur.close()
    conn.close()
    return found_sheet


# ============================================================================
# Check 2: PowerPoint
# ============================================================================

def check_pptx(agent_workspace):
    print("\n=== Checking Enrollment_Overview_F2013.pptx ===")

    pptx_path = os.path.join(agent_workspace, "Enrollment_Overview_F2013.pptx")
    if not os.path.isfile(pptx_path):
        record("PPT file exists", False, f"Not found: {pptx_path}")
        return False
    record("PPT file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        record("PPT has at least 8 slides (title + 6 courses + summary)",
               slide_count >= 8,
               f"Found {slide_count} slides")

        # Collect all text from slides
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        record("PPT mentions Fall 2013",
               "fall 2013" in all_text or "2013" in all_text,
               "Missing year 2013")
        record("PPT mentions enrollment",
               "enrollment" in all_text or "enrolled" in all_text or "students" in all_text,
               "Missing enrollment content")

        # Check for course codes
        for code in ["fff-2013j", "bbb-2013j"]:
            record(f"PPT mentions {code.upper()}",
                   code in all_text,
                   f"Missing {code} in slides")

        return slide_count >= 8

    except ImportError:
        size = os.path.getsize(pptx_path)
        record("PPT file has content (>5KB)", size > 5000, f"Size: {size} bytes")
        return size > 5000
    except Exception as e:
        record("PPT file readable", False, str(e))
        return False


# ============================================================================
# Check 3: Word document
# ============================================================================

def check_word(agent_workspace):
    print("\n=== Checking Enrollment_Report_F2013.docx ===")

    docx_path = os.path.join(agent_workspace, "Enrollment_Report_F2013.docx")
    if not os.path.isfile(docx_path):
        record("Word file exists", False, f"Not found: {docx_path}")
        return False
    record("Word file exists", True)

    try:
        from docx import Document
        doc = Document(docx_path)
        all_text = " ".join(p.text for p in doc.paragraphs).lower()
        headings = " ".join(p.text for p in doc.paragraphs
                           if p.style.name.startswith("Heading")).lower()

        record("Word doc has substantial content",
               len(all_text.strip()) >= 100,
               f"Content length: {len(all_text.strip())}")
        record("Word doc mentions enrollment",
               "enrollment" in all_text or "enroll" in all_text,
               "Missing 'enrollment' in document")
        record("Word doc mentions Fall 2013",
               "fall 2013" in all_text or "2013" in all_text,
               "Missing 2013")

        # Check for table
        tables = doc.tables
        record("Word doc has at least 1 table",
               len(tables) >= 1,
               f"Found {len(tables)} tables")

        # Check for total enrollment number
        record("Word doc mentions total enrollment",
               "8845" in all_text or "8,845" in all_text or "total" in all_text,
               "Missing total enrollment figure")

        return True

    except ImportError:
        size = os.path.getsize(docx_path)
        record("Word file has content (>3KB)", size > 3000, f"Size: {size} bytes")
        return size > 3000
    except Exception as e:
        record("Word file readable", False, str(e))
        return False


# ============================================================================
# Check 4: Email
# ============================================================================

def check_emails():
    print("\n=== Checking Emails ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
    """)
    all_emails = cur.fetchall()
    cur.close()
    conn.close()

    print(f"[check_emails] Found {len(all_emails)} total emails.")
    record("At least 1 email sent", len(all_emails) >= 1, f"Found {len(all_emails)}")

    found_email = False
    for subject, from_addr, to_addr, body_text in all_emails:
        to_str = str(to_addr or "").lower()
        subject_lower = (subject or "").lower()
        if ("academic.office@university.edu" in to_str or
                "enrollment" in subject_lower or "fall 2013" in subject_lower):
            found_email = True
            record("Email to academic.office@university.edu found", True)

            record("Email subject mentions enrollment or Fall 2013",
                   "enrollment" in subject_lower or "fall 2013" in subject_lower or "2013" in subject_lower,
                   f"Subject: {subject}")

            body_lower = (body_text or "").lower()
            record("Email body has enrollment stats",
                   any(term in body_lower for term in ["enrollment", "students", "courses", "2013"]),
                   "Body missing enrollment content")
            break

    if not found_email:
        record("Enrollment email found", False,
               f"Emails: {[(e[0], str(e[2])[:60]) for e in all_emails[:3]]}")

    return found_email


# ============================================================================
# Main
# ============================================================================

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    gsheet_ok = check_gsheet()
    pptx_ok = check_pptx(args.agent_workspace)
    word_ok = check_word(args.agent_workspace)
    email_ok = check_emails()

    all_passed = gsheet_ok and pptx_ok and word_ok and email_ok

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Overall: {'PASS' if all_passed else 'FAIL'}")

    if args.res_log_file:
        result = {
            "passed": PASS_COUNT,
            "failed": FAIL_COUNT,
            "success": all_passed,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if all_passed else 1)


if __name__ == "__main__":
    main()
