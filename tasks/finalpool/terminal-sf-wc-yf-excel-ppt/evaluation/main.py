"""Evaluation for terminal-sf-wc-yf-excel-ppt.

Checks:
1. Executive_Dashboard.xlsx with 4 sheets and correct data
2. Executive_Briefing.pptx with 8 slides
3. unified_metrics.json exists with expected keys
4. business_health.json with composite score
5. alerts.json with below-target KPIs
"""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432,
          dbname=os.environ.get("PGDATABASE", "toolathlon_gym"),
          user="eigent", password="camel")

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:300]}")


def num_close(a, b, tol=2.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def get_sheet(wb, name):
    for s in wb.sheetnames:
        if s.strip().lower() == name.strip().lower():
            return wb[s]
    return None


def get_expected_data():
    """Query DB to get expected values."""
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()

        # SF segments
        cur.execute("""
            SELECT c."SEGMENT",
                   ROUND(SUM(o."TOTAL_AMOUNT")::numeric, 2) as rev,
                   COUNT(DISTINCT c."CUSTOMER_ID") as cust,
                   COUNT(DISTINCT o."ORDER_ID") as orders
            FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
            JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
            GROUP BY c."SEGMENT" ORDER BY c."SEGMENT"
        """)
        segments = {}
        total_enterprise_rev = 0
        for row in cur.fetchall():
            seg, rev, cust, orders = row
            segments[seg.lower()] = {"revenue": float(rev), "customers": cust, "orders": orders}
            total_enterprise_rev += float(rev)

        # WC total revenue
        cur.execute("SELECT ROUND(SUM(total::numeric), 2) FROM wc.orders")
        wc_total_rev = float(cur.fetchone()[0])

        # WC avg order value
        cur.execute("SELECT ROUND(AVG(total::numeric), 2) FROM wc.orders")
        wc_avg_order = float(cur.fetchone()[0])

        # WC top category
        cur.execute("""
            SELECT p.categories->0->>'name' as cat,
                   ROUND(SUM((li.value->>'total')::numeric), 2) as rev
            FROM wc.orders o, jsonb_array_elements(o.line_items) li
            JOIN wc.products p ON p.id = (li.value->>'product_id')::int
            GROUP BY p.categories->0->>'name'
            ORDER BY rev DESC LIMIT 1
        """)
        top_ecom_cat = cur.fetchone()[0]

        # YF latest prices
        yf_prices = {}
        for sym in ['^DJI', 'AMZN', 'JPM']:
            cur.execute("SELECT close FROM yf.stock_prices WHERE symbol=%s ORDER BY date DESC LIMIT 1", (sym,))
            row = cur.fetchone()
            yf_prices[sym] = float(row[0]) if row else 0

        # Top segment
        top_seg = max(segments.keys(), key=lambda k: segments[k]["revenue"])

        cur.close()
        conn.close()

        return {
            "segments": segments,
            "total_enterprise_rev": total_enterprise_rev,
            "wc_total_rev": wc_total_rev,
            "wc_avg_order": wc_avg_order,
            "top_ecom_cat": top_ecom_cat,
            "yf_prices": yf_prices,
            "top_seg": top_seg,
        }
    except Exception as e:
        print(f"  [WARN] DB query failed: {e}")
        return None


def check_excel(agent_workspace, expected):
    print("\n=== Checking Executive_Dashboard.xlsx ===")
    path = os.path.join(agent_workspace, "Executive_Dashboard.xlsx")
    check("Excel file exists", os.path.isfile(path), path)
    if not os.path.isfile(path):
        return

    try:
        wb = openpyxl.load_workbook(path, data_only=True)
    except Exception as e:
        check("Excel readable", False, str(e))
        return

    # Check 4 sheets exist
    check("Has 4 sheets", len(wb.sheetnames) >= 4, f"Got {wb.sheetnames}")

    # Enterprise_Sales
    ws = get_sheet(wb, "Enterprise_Sales")
    check("Sheet Enterprise_Sales exists", ws is not None, f"Sheets: {wb.sheetnames}")
    if ws and expected:
        rows = list(ws.iter_rows(min_row=2, values_only=True))
        check("Enterprise_Sales has 4 rows", len(rows) == 4, f"Got {len(rows)}")
        lookup = {str(r[0]).strip().lower(): r for r in rows if r and r[0]}
        for seg_name, seg_data in expected["segments"].items():
            r = lookup.get(seg_name)
            if r is None:
                check(f"Segment '{seg_name}' present", False, "Missing")
                continue
            check(f"'{seg_name}' revenue",
                  num_close(r[1], seg_data["revenue"], 500),
                  f"Expected {seg_data['revenue']}, got {r[1]}")
            check(f"'{seg_name}' customer_count",
                  num_close(r[2], seg_data["customers"], 5),
                  f"Expected {seg_data['customers']}, got {r[2]}")

    # Ecommerce_Performance
    ws = get_sheet(wb, "Ecommerce_Performance")
    check("Sheet Ecommerce_Performance exists", ws is not None, f"Sheets: {wb.sheetnames}")
    if ws:
        rows = list(ws.iter_rows(min_row=2, values_only=True))
        check("Ecommerce_Performance has >= 5 rows", len(rows) >= 5, f"Got {len(rows)}")
        if rows and expected:
            # First row should be highest revenue category
            top_cat = str(rows[0][0]).strip().lower() if rows[0] and rows[0][0] else ""
            check("Top ecom category is TV & Home Theater",
                  "tv" in top_cat and "home" in top_cat,
                  f"Got '{top_cat}'")

    # Market_Context
    ws = get_sheet(wb, "Market_Context")
    check("Sheet Market_Context exists", ws is not None, f"Sheets: {wb.sheetnames}")
    if ws and expected:
        rows = list(ws.iter_rows(min_row=2, values_only=True))
        check("Market_Context has 3 rows", len(rows) == 3, f"Got {len(rows)}")
        lookup = {str(r[0]).strip(): r for r in rows if r and r[0]}
        for sym in ['^DJI', 'AMZN', 'JPM']:
            r = lookup.get(sym)
            if r is None:
                check(f"Symbol '{sym}' present", False, "Missing")
                continue
            check(f"'{sym}' price",
                  num_close(r[1], expected["yf_prices"][sym], 100),
                  f"Expected {expected['yf_prices'][sym]}, got {r[1]}")

    # Unified_Dashboard
    ws = get_sheet(wb, "Unified_Dashboard")
    check("Sheet Unified_Dashboard exists", ws is not None, f"Sheets: {wb.sheetnames}")
    if ws and expected:
        rows = list(ws.iter_rows(min_row=2, values_only=True))
        check("Unified_Dashboard has 5 rows", len(rows) == 5, f"Got {len(rows)}")
        lookup = {str(r[0]).strip().lower(): r for r in rows if r and r[0]}

        ent_row = lookup.get("enterprise_total_revenue")
        if ent_row:
            check("Enterprise total revenue value",
                  num_close(ent_row[1], expected["total_enterprise_rev"], 5000),
                  f"Expected {expected['total_enterprise_rev']}, got {ent_row[1]}")
            check("Enterprise target is 3200000",
                  num_close(ent_row[2], 3200000, 100),
                  f"Got {ent_row[2]}")
            status = str(ent_row[3]).strip().lower() if ent_row[3] else ""
            check("Enterprise status is Below Target",
                  "below" in status,
                  f"Got '{ent_row[3]}'")

        ecom_row = lookup.get("ecommerce_total_revenue")
        if ecom_row:
            check("Ecom total revenue value",
                  num_close(ecom_row[1], expected["wc_total_rev"], 2000),
                  f"Expected {expected['wc_total_rev']}, got {ecom_row[1]}")

        mkt_row = lookup.get("market_index")
        if mkt_row:
            check("Market index value",
                  num_close(mkt_row[1], expected["yf_prices"]["^DJI"], 200),
                  f"Expected {expected['yf_prices']['^DJI']}, got {mkt_row[1]}")
            status = str(mkt_row[3]).strip().lower() if mkt_row[3] else ""
            check("Market status is On Track",
                  "on track" in status,
                  f"Got '{mkt_row[3]}'")


def check_pptx(agent_workspace):
    print("\n=== Checking Executive_Briefing.pptx ===")
    path = os.path.join(agent_workspace, "Executive_Briefing.pptx")
    check("PPTX file exists", os.path.isfile(path), path)
    if not os.path.isfile(path):
        return
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slide_count = len(prs.slides)
        check("Has 8 slides", slide_count == 8, f"Got {slide_count}")
        check("Has >= 6 slides", slide_count >= 6, f"Got {slide_count}")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        check("Contains 'business intelligence' or 'dashboard'",
              "business intelligence" in all_text or "dashboard" in all_text,
              "Missing title reference")
        check("Contains 'health score' or 'composite'",
              "health score" in all_text or "composite" in all_text,
              "Missing health score slide")
        check("Contains 'enterprise' reference",
              "enterprise" in all_text)
        check("Contains 'e-commerce' or 'ecommerce'",
              "e-commerce" in all_text or "ecommerce" in all_text)
        check("Contains 'market' reference",
              "market" in all_text)
        check("Contains 'risk' reference",
              "risk" in all_text)
        check("Contains 'strategic' or 'priorities' or 'recommendation'",
              "strategic" in all_text or "priorities" in all_text or "recommend" in all_text)
    except ImportError:
        check("python-pptx available", False)
    except Exception as e:
        check("PPTX readable", False, str(e))


def check_json_files(agent_workspace, expected):
    print("\n=== Checking JSON Output Files ===")

    # unified_metrics.json
    um_path = os.path.join(agent_workspace, "unified_metrics.json")
    check("unified_metrics.json exists", os.path.isfile(um_path), um_path)
    if os.path.isfile(um_path):
        try:
            with open(um_path) as f:
                um = json.load(f)
            check("unified_metrics.json is valid JSON", True)
            # Should contain enterprise, ecommerce, and market data
            um_str = json.dumps(um).lower()
            check("Contains enterprise data", "enterprise" in um_str or "segment" in um_str)
            check("Contains ecommerce data", "ecommerce" in um_str or "e-commerce" in um_str or "category" in um_str or "woocommerce" in um_str)
            check("Contains market data", "dji" in um_str or "market" in um_str or "amzn" in um_str)
        except json.JSONDecodeError as e:
            check("unified_metrics.json valid JSON", False, str(e))

    # business_health.json
    bh_path = os.path.join(agent_workspace, "business_health.json")
    check("business_health.json exists", os.path.isfile(bh_path), bh_path)
    if os.path.isfile(bh_path) and expected:
        try:
            with open(bh_path) as f:
                bh = json.load(f)
            # enterprise_score = min(100, total_ent_rev / 3200000 * 100)
            exp_ent_score = min(100, expected["total_enterprise_rev"] / 3200000 * 100)
            exp_ecom_score = min(100, expected["wc_total_rev"] / 70000 * 100)
            exp_mkt_score = 80 if expected["yf_prices"]["^DJI"] > 40000 else 50
            exp_composite = 0.4 * exp_ent_score + 0.3 * exp_ecom_score + 0.3 * exp_mkt_score

            check("enterprise_score present", "enterprise_score" in bh, f"Keys: {list(bh.keys())}")
            check("composite_score present", "composite_score" in bh, f"Keys: {list(bh.keys())}")
            if "composite_score" in bh:
                check("composite_score value",
                      num_close(bh["composite_score"], exp_composite, 5),
                      f"Expected ~{exp_composite:.1f}, got {bh['composite_score']}")
            if "enterprise_score" in bh:
                check("enterprise_score value",
                      num_close(bh["enterprise_score"], exp_ent_score, 3),
                      f"Expected ~{exp_ent_score:.1f}, got {bh['enterprise_score']}")
        except json.JSONDecodeError as e:
            check("business_health.json valid JSON", False, str(e))

    # alerts.json
    al_path = os.path.join(agent_workspace, "alerts.json")
    check("alerts.json exists", os.path.isfile(al_path), al_path)
    if os.path.isfile(al_path):
        try:
            with open(al_path) as f:
                alerts = json.load(f)
            check("alerts.json is a list", isinstance(alerts, list))
            # Enterprise and Ecommerce should be below target
            check("At least 2 alerts (enterprise + ecommerce below target)",
                  len(alerts) >= 2, f"Got {len(alerts)} alerts")
            if alerts:
                alert_names = [str(a.get("kpi_name", "")).lower() for a in alerts]
                check("Enterprise below target alert present",
                      any("enterprise" in n for n in alert_names),
                      f"Alert names: {alert_names}")
                check("Ecommerce below target alert present",
                      any("ecommerce" in n or "e-commerce" in n for n in alert_names),
                      f"Alert names: {alert_names}")
        except json.JSONDecodeError as e:
            check("alerts.json valid JSON", False, str(e))


def check_scripts(agent_workspace):
    print("\n=== Checking Scripts ===")
    for script in ["merge_data.py", "health_score.py", "generate_alerts.py"]:
        path = os.path.join(agent_workspace, script)
        check(f"{script} exists", os.path.isfile(path), path)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    expected = get_expected_data()

    check_excel(args.agent_workspace, expected)
    check_pptx(args.agent_workspace)
    check_json_files(args.agent_workspace, expected)
    check_scripts(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total > 0 else 0
    print(f"\nOverall: {PASS_COUNT}/{total} ({accuracy:.1f}%)")
    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)
    sys.exit(0 if accuracy >= 70 else 1)


if __name__ == "__main__":
    main()
