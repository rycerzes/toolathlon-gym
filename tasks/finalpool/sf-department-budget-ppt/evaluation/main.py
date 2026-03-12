"""
Evaluation script for sf-department-budget-ppt task.

Checks:
1. Excel file (budget_vs_actuals.xlsx) - both sheets
2. PowerPoint file (budget_vs_actuals.pptx) - structure and content
"""

import argparse
import json
import os
import sys

import openpyxl

PASS_COUNT = 0
FAIL_COUNT = 0

DEPARTMENTS = ["Engineering", "Finance", "HR", "Operations", "R&D", "Sales", "Support"]

# Expected actual expenditures from Snowflake
EXPECTED_ACTUALS = {
    "Engineering": 418604451.00,
    "Finance": 413713333.00,
    "HR": 416980025.00,
    "Operations": 411598247.00,
    "R&D": 410147671.00,
    "Sales": 425710161.00,
    "Support": 423053067.00,
}

EXPECTED_BUDGETS = {
    "Engineering": 633555641.27,
    "Finance": 616980984.11,
    "HR": 697911374.65,
    "Operations": 572668189.63,
    "R&D": 715914847.42,
    "Sales": 642575220.92,
    "Support": 684741486.31,
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def str_match(a, b):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    return str(a).strip().lower() == str(b).strip().lower()


def num_close(a, b, rel_tol=0.02):
    try:
        a, b = float(a), float(b)
        if b == 0:
            return abs(a) < 1000
        return abs(a - b) / abs(b) <= rel_tol
    except (TypeError, ValueError):
        return False


# ============================================================================
# Check 1: Excel file
# ============================================================================

def check_excel(agent_workspace, groundtruth_workspace):
    print("\n=== Checking Excel Output ===")

    agent_file = os.path.join(agent_workspace, "budget_vs_actuals.xlsx")

    if not os.path.isfile(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False
    record("Excel file exists", True)

    wb = openpyxl.load_workbook(agent_file, data_only=True)

    def get_sheet(wb, target):
        for name in wb.sheetnames:
            if name.strip().lower() == target.strip().lower():
                return wb[name]
        return None

    # Sheet 1: Department Details
    ws1 = get_sheet(wb, "Department Details")
    if ws1 is None:
        record("Sheet 'Department Details' exists", False, f"Sheets: {wb.sheetnames}")
        return False
    record("Sheet 'Department Details' exists", True)

    headers = [str(c.value).strip() if c.value else "" for c in ws1[1]]
    expected = ["Department", "Approved_Budget", "Actual_Expenditure", "Variance",
                "Variance_Pct", "Headcount", "Avg_Salary"]
    headers_ok = all(str_match(h, e) for h, e in zip(headers, expected))
    record("Department Details headers match", headers_ok,
           f"Expected: {expected}, Got: {headers}")

    rows = list(ws1.iter_rows(min_row=2, values_only=True))
    record("Department Details has 7 rows", len(rows) == 7, f"Got {len(rows)}")

    for dept in DEPARTMENTS:
        agent_row = None
        for r in rows:
            if r and str_match(r[0], dept):
                agent_row = r
                break

        if not agent_row:
            record(f"Dept {dept} present", False, "Missing")
            continue

        record(f"Dept {dept} present", True)

        # Check actual expenditure
        record(f"Dept {dept}: Actual Expenditure",
               num_close(agent_row[2], EXPECTED_ACTUALS[dept]),
               f"Expected ~{EXPECTED_ACTUALS[dept]:.0f}, got {agent_row[2]}")

        # Check budget
        record(f"Dept {dept}: Approved Budget",
               num_close(agent_row[1], EXPECTED_BUDGETS[dept]),
               f"Expected ~{EXPECTED_BUDGETS[dept]:.0f}, got {agent_row[1]}")

        # Check variance = budget - actuals
        expected_var = EXPECTED_BUDGETS[dept] - EXPECTED_ACTUALS[dept]
        record(f"Dept {dept}: Variance",
               num_close(agent_row[3], expected_var),
               f"Expected ~{expected_var:.0f}, got {agent_row[3]}")

    # Sheet 2: Totals
    ws2 = get_sheet(wb, "Totals")
    if ws2 is None:
        record("Sheet 'Totals' exists", False, f"Sheets: {wb.sheetnames}")
        return False
    record("Sheet 'Totals' exists", True)

    summary = {}
    for row in ws2.iter_rows(min_row=1, values_only=True):
        if row and row[0]:
            summary[str(row[0]).strip().lower()] = row[1]

    record("Totals: Total_Departments = 7",
           str(summary.get("total_departments", "")).strip() == "7",
           f"Got {summary.get('total_departments')}")

    total_budget = sum(EXPECTED_BUDGETS.values())
    total_actuals = sum(EXPECTED_ACTUALS.values())

    record("Totals: Total_Budget",
           num_close(summary.get("total_budget", 0), total_budget),
           f"Expected ~{total_budget:.0f}, got {summary.get('total_budget')}")

    record("Totals: Total_Actuals",
           num_close(summary.get("total_actuals", 0), total_actuals),
           f"Expected ~{total_actuals:.0f}, got {summary.get('total_actuals')}")

    record("Totals: Most_Under_Budget = R&D",
           str_match(summary.get("most_under_budget", ""), "R&D"),
           f"Got {summary.get('most_under_budget')}")

    record("Totals: Most_Over_Budget = Operations",
           str_match(summary.get("most_over_budget", ""), "Operations"),
           f"Got {summary.get('most_over_budget')}")

    return True


# ============================================================================
# Check 2: PowerPoint file
# ============================================================================

def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")

    pptx_path = os.path.join(agent_workspace, "budget_vs_actuals.pptx")

    if not os.path.isfile(pptx_path):
        record("PPTX file exists", False, f"Not found: {pptx_path}")
        return False
    record("PPTX file exists", True)

    try:
        from pptx import Presentation
    except ImportError:
        record("python-pptx installed", False, "Cannot import pptx")
        return False

    prs = Presentation(pptx_path)
    slides = list(prs.slides)

    # Should have: 1 title + 1 overview + 7 departments + 1 summary = 10 slides
    record("PPTX has at least 10 slides", len(slides) >= 10,
           f"Found {len(slides)} slides")

    # Collect all text
    all_text = ""
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + "\n"

    all_lower = all_text.lower()

    record("PPTX: has title 'budget vs actuals'",
           "budget vs actuals" in all_lower or "fy2026" in all_lower,
           "Title not found")

    record("PPTX: mentions March 2026",
           "march 2026" in all_lower,
           "March 2026 not found")

    for dept in DEPARTMENTS:
        record(f"PPTX: mentions {dept}",
               dept.lower() in all_lower,
               f"{dept} not found")

    record("PPTX: has Summary slide",
           "summary" in all_lower,
           "Summary not found")

    record("PPTX: mentions 'variance' or 'budget'",
           "variance" in all_lower or "budget" in all_lower,
           "Budget/variance terms not found")

    return True


# ============================================================================
# Main
# ============================================================================

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    excel_ok = check_excel(args.agent_workspace, args.groundtruth_workspace)
    pptx_ok = check_pptx(args.agent_workspace)

    all_passed = excel_ok and pptx_ok

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Overall: {'PASS' if all_passed else 'FAIL'}")

    if args.res_log_file:
        result = {
            "passed": PASS_COUNT,
            "failed": FAIL_COUNT,
            "success": all_passed,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if all_passed else 1)


if __name__ == "__main__":
    main()
