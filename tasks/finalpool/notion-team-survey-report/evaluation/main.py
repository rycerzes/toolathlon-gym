"""
Evaluation script for notion-team-survey-report task.

Checks:
1. Excel file (Survey_Analysis.xlsx) exists with 2 sheets and correct data
2. PowerPoint file (Team_Report.pptx) exists with at least 5 slides
3. Email sent to management@company.com with survey results
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

EXPECTED_AVERAGES = {
    "leadership": 3.5,
    "workload": 3.5,
    "communication": 3.6,
    "growth": 3.5,
}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=0.3):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def check_excel(agent_workspace):
    print("\n=== Checking Excel Output ===")

    excel_path = os.path.join(agent_workspace, "Survey_Analysis.xlsx")
    check("Excel file exists", os.path.isfile(excel_path),
          f"Expected {excel_path}")
    if not os.path.isfile(excel_path):
        return

    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
    except Exception as e:
        check("Excel file readable", False, str(e))
        return
    check("Excel file readable", True)

    sheet_names_lower = [s.lower().replace("_", " ").strip() for s in wb.sheetnames]

    has_data_sheet = any("survey" in s and "data" in s for s in sheet_names_lower) or \
                     any("data" in s or "response" in s for s in sheet_names_lower)
    check("Has Survey Data sheet", has_data_sheet,
          f"Found sheets: {wb.sheetnames}")

    ws_data = None
    for s in wb.sheetnames:
        sl = s.lower().replace("_", " ")
        if ("survey" in sl and "data" in sl) or "response" in sl:
            ws_data = wb[s]
            break
    if ws_data is None and len(wb.sheetnames) >= 1:
        ws_data = wb[wb.sheetnames[0]]

    if ws_data:
        data_rows = list(ws_data.iter_rows(min_row=2, values_only=True))
        data_rows = [r for r in data_rows if r and r[0] is not None]
        check("Survey Data has 10 rows", len(data_rows) == 10,
              f"Found {len(data_rows)} data rows")

        all_names = " ".join(str(r[0]).lower() for r in data_rows if r and r[0])
        expected_names = ["alice", "bob", "carol", "david", "eva", "frank", "grace", "henry", "irene", "jack"]
        names_found = sum(1 for n in expected_names if n in all_names)
        check("At least 6 of 8 respondent names present",
              names_found >= 6,
              f"Found {names_found}/8 names")

    has_summary_sheet = any("summary" in s or "statistic" in s or "average" in s for s in sheet_names_lower)
    check("Has Summary Statistics sheet", has_summary_sheet,
          f"Found sheets: {wb.sheetnames}")

    ws_summary = None
    for s in wb.sheetnames:
        sl = s.lower().replace("_", " ")
        if "summary" in sl or "statistic" in sl or "average" in sl:
            ws_summary = wb[s]
            break

    if ws_summary:
        all_values = []
        for row in ws_summary.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    try:
                        all_values.append(float(cell))
                    except (TypeError, ValueError):
                        pass

        for dimension, expected_avg in EXPECTED_AVERAGES.items():
            found = any(num_close(v, expected_avg, 0.2) for v in all_values)
            check(f"Average {dimension} ~{expected_avg}",
                  found,
                  f"Expected {expected_avg}, numeric values: {[v for v in all_values if abs(v - expected_avg) < 1][:5]}")


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint Output ===")

    pptx_path = os.path.join(agent_workspace, "Team_Report.pptx")
    check("PowerPoint file exists", os.path.isfile(pptx_path),
          f"Expected {pptx_path}")
    if not os.path.isfile(pptx_path):
        return

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        check("PowerPoint file readable", False, str(e))
        return
    check("PowerPoint file readable", True)

    slide_count = len(prs.slides)
    check("PowerPoint has at least 5 slides",
          slide_count >= 5,
          f"Found {slide_count} slides")

    all_slide_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_slide_text += para.text + " "
    all_text_lower = all_slide_text.lower()

    if slide_count > 0:
        first_text = ""
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    first_text += para.text + " "
        first_lower = first_text.lower()
        check("Title slide mentions engineering or quarterly or report",
              "engineering" in first_lower or "quarterly" in first_lower or "report" in first_lower,
              f"First slide text: {first_text[:150]}")

    project_names = ["alpha", "beta", "gamma", "delta", "epsilon"]
    projects_found = sum(1 for p in project_names if p in all_text_lower)
    check("PowerPoint references at least 3 project names",
          projects_found >= 3,
          f"Found {projects_found}/5 project names")

    dimensions = ["leadership", "workload", "communication", "growth"]
    dims_found = sum(1 for d in dimensions if d in all_text_lower)
    check("PowerPoint mentions at least 3 survey dimensions",
          dims_found >= 3,
          f"Found {dims_found}/4 dimensions")

    check("PowerPoint has recommendation content",
          any(kw in all_text_lower for kw in
              ["recommend", "improvement", "concern", "action", "suggestion", "next step"]),
          "No recommendation keywords found")


def check_emails():
    print("\n=== Checking Emails ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
        ORDER BY created_at DESC
    """)
    all_messages = cur.fetchall()
    cur.close()
    conn.close()

    matching_email = None
    for subject, from_addr, to_addr, body_text in all_messages:
        subj_lower = (subject or "").lower()
        if "quarterly" in subj_lower or "report" in subj_lower or "engineering" in subj_lower:
            matching_email = (subject, from_addr, to_addr, body_text)
            break

    check("Email with report-related subject exists",
          matching_email is not None,
          f"Found {len(all_messages)} total emails, none matching")

    if matching_email:
        subject, from_addr, to_addr, body_text = matching_email

        to_str = ""
        if isinstance(to_addr, list):
            to_str = " ".join(str(r).lower() for r in to_addr)
        elif isinstance(to_addr, str):
            try:
                parsed = json.loads(to_addr)
                if isinstance(parsed, list):
                    to_str = " ".join(str(r).lower() for r in parsed)
                else:
                    to_str = str(to_addr).lower()
            except (json.JSONDecodeError, TypeError):
                to_str = str(to_addr).lower()

        check("Email sent to management@company.com",
              "management@company.com" in to_str,
              f"Recipient: {to_addr}")

        body_lower = (body_text or "").lower()

        dims_mentioned = sum(1 for d in ["leadership", "workload", "communication", "growth"]
                           if d in body_lower)
        check("Email mentions at least 3 survey dimensions",
              dims_mentioned >= 3,
              f"Found {dims_mentioned}/4 dimensions in email body")

        has_scores = any(str(v) in body_lower or f"{v:.1f}" in body_lower
                        for v in EXPECTED_AVERAGES.values())
        check("Email mentions average scores",
              has_scores or "3.5" in body_lower or "3.625" in body_lower or "3.63" in body_lower,
              "No average scores found in email body")

        project_mentions = sum(1 for p in ["alpha", "beta", "gamma", "delta", "epsilon"]
                             if p in body_lower)
        check("Email mentions at least 2 project names",
              project_mentions >= 2 or "project" in body_lower,
              f"Found {project_mentions} project name mentions")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    pass_rate = PASS_COUNT / total if total > 0 else 0

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Pass Rate: {pass_rate:.1%}")

    result = {
        "passed": PASS_COUNT,
        "failed": FAIL_COUNT,
        "pass_rate": round(pass_rate, 3),
        "success": pass_rate >= 0.7,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if pass_rate >= 0.7 else 1)


if __name__ == "__main__":
    main()
