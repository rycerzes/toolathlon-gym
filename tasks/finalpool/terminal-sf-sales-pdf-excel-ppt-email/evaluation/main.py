"""Evaluation for terminal-sf-sales-pdf-excel-ppt-email.

Checks:
1. Quarterly_Sales_Review.xlsx with 4 sheets
2. Sales_Review_Presentation.pptx with >= 5 slides
3. Email sent to regional_managers@company.com
"""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432,
          dbname=os.environ.get("PGDATABASE", "toolathlon_gym"),
          user="eigent", password="camel")

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:300]}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def get_sheet(wb, name):
    for s in wb.sheetnames:
        if s.strip().lower() == name.strip().lower():
            return wb[s]
    return None


def check_excel(agent_workspace, groundtruth_workspace):
    print("\n=== Checking Quarterly_Sales_Review.xlsx ===")
    agent_file = os.path.join(agent_workspace, "Quarterly_Sales_Review.xlsx")
    gt_file = os.path.join(groundtruth_workspace, "Quarterly_Sales_Review.xlsx")

    check("Excel file exists", os.path.isfile(agent_file), agent_file)
    if not os.path.isfile(agent_file):
        return

    try:
        agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_file, data_only=True)
    except Exception as e:
        check("Excel readable", False, str(e))
        return

    # Regional_Performance
    print("  Checking Regional_Performance...")
    a_sheet = get_sheet(agent_wb, "Regional_Performance")
    g_sheet = get_sheet(gt_wb, "Regional_Performance")
    check("Sheet 'Regional_Performance' exists", a_sheet is not None,
          f"Sheets: {agent_wb.sheetnames}")
    if a_sheet and g_sheet:
        a_rows = list(a_sheet.iter_rows(min_row=2, values_only=True))
        g_rows = list(g_sheet.iter_rows(min_row=2, values_only=True))
        check("Regional_Performance has 5 rows", len(a_rows) == 5, f"Got {len(a_rows)}")

        a_lookup = {str(r[0]).strip().lower(): r for r in a_rows if r and r[0]}
        for g_row in g_rows:
            if not g_row or not g_row[0]:
                continue
            key = str(g_row[0]).strip().lower()
            a_row = a_lookup.get(key)
            if a_row is None:
                check(f"Region '{g_row[0]}' present", False, "Missing")
                continue
            # Target_Revenue
            if len(a_row) > 1 and len(g_row) > 1:
                check(f"'{key}' Target",
                      num_close(a_row[1], g_row[1], 1000),
                      f"Expected {g_row[1]}, got {a_row[1]}")
            # Actual_Revenue
            if len(a_row) > 2 and len(g_row) > 2:
                check(f"'{key}' Actual",
                      num_close(a_row[2], g_row[2], 500),
                      f"Expected {g_row[2]}, got {a_row[2]}")
            # Variance_Pct
            if len(a_row) > 4 and len(g_row) > 4:
                check(f"'{key}' Variance_Pct",
                      num_close(a_row[4], g_row[4], 1.0),
                      f"Expected {g_row[4]}, got {a_row[4]}")
            # Status
            if len(a_row) > 5 and len(g_row) > 5:
                a_status = str(a_row[5]).strip().lower() if a_row[5] else ""
                g_status = str(g_row[5]).strip().lower() if g_row[5] else ""
                check(f"'{key}' Status",
                      a_status == g_status,
                      f"Expected {g_status}, got {a_status}")

    # Segment_Breakdown
    print("  Checking Segment_Breakdown...")
    a_sheet = get_sheet(agent_wb, "Segment_Breakdown")
    g_sheet = get_sheet(gt_wb, "Segment_Breakdown")
    check("Sheet 'Segment_Breakdown' exists", a_sheet is not None,
          f"Sheets: {agent_wb.sheetnames}")
    if a_sheet and g_sheet:
        a_rows = list(a_sheet.iter_rows(min_row=2, values_only=True))
        check("Segment_Breakdown has 4 rows", len(a_rows) == 4, f"Got {len(a_rows)}")

        a_lookup = {str(r[0]).strip().lower(): r for r in a_rows if r and r[0]}
        for g_row in g_sheet.iter_rows(min_row=2, values_only=True):
            if not g_row or not g_row[0]:
                continue
            key = str(g_row[0]).strip().lower()
            a_row = a_lookup.get(key)
            if a_row is None:
                check(f"Segment '{g_row[0]}' present", False, "Missing")
                continue
            if len(a_row) > 2 and len(g_row) > 2:
                check(f"'{key}' Order_Count",
                      num_close(a_row[2], g_row[2], 50),
                      f"Expected {g_row[2]}, got {a_row[2]}")
            if len(a_row) > 3 and len(g_row) > 3:
                check(f"'{key}' Total_Revenue",
                      num_close(a_row[3], g_row[3], 5000),
                      f"Expected {g_row[3]}, got {a_row[3]}")

    # Top_Products
    print("  Checking Top_Products...")
    a_sheet = get_sheet(agent_wb, "Top_Products")
    check("Sheet 'Top_Products' exists", a_sheet is not None,
          f"Sheets: {agent_wb.sheetnames}")
    if a_sheet:
        a_rows = list(a_sheet.iter_rows(min_row=2, values_only=True))
        check("Top_Products has 5 rows", len(a_rows) == 5, f"Got {len(a_rows)}")
        if a_rows:
            # Check top product has highest revenue
            revenues = [float(r[3]) for r in a_rows if r and len(r) > 3 and r[3]]
            if revenues:
                check("Top product revenue > 100000",
                      revenues[0] > 100000,
                      f"Got {revenues[0]}")

    # Summary
    print("  Checking Summary...")
    a_sheet = get_sheet(agent_wb, "Summary")
    g_sheet = get_sheet(gt_wb, "Summary")
    check("Sheet 'Summary' exists", a_sheet is not None,
          f"Sheets: {agent_wb.sheetnames}")
    if a_sheet and g_sheet:
        a_data = {}
        for row in a_sheet.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                a_data[str(row[0]).strip().lower()] = row[1]
        g_data = {}
        for row in g_sheet.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                g_data[str(row[0]).strip().lower()] = row[1]

        check("Total_Target",
              num_close(a_data.get("total_target"), g_data.get("total_target"), 1000),
              f"Expected {g_data.get('total_target')}, got {a_data.get('total_target')}")
        check("Total_Actual",
              num_close(a_data.get("total_actual"), g_data.get("total_actual"), 5000),
              f"Expected {g_data.get('total_actual')}, got {a_data.get('total_actual')}")
        check("Regions_Met_Target",
              num_close(a_data.get("regions_met_target"), g_data.get("regions_met_target"), 0),
              f"Expected {g_data.get('regions_met_target')}, got {a_data.get('regions_met_target')}")
        # Dynamically determine best region from DB
        expected_best_region = "europe"  # fallback
        try:
            conn_db = psycopg2.connect(**DB)
            cur_db = conn_db.cursor()
            cur_db.execute("""
                SELECT region FROM sf_data.orders
                GROUP BY region ORDER BY SUM(sales) DESC LIMIT 1
            """)
            result = cur_db.fetchone()
            if result and result[0]:
                expected_best_region = str(result[0]).strip().lower()
            cur_db.close()
            conn_db.close()
        except Exception:
            pass
        br = a_data.get("best_region")
        check("Best_Region correct",
              br is not None and expected_best_region in str(br).lower(),
              f"Got {br}, expected {expected_best_region}")


def check_pptx(agent_workspace):
    print("\n=== Checking Sales_Review_Presentation.pptx ===")
    pptx_path = os.path.join(agent_workspace, "Sales_Review_Presentation.pptx")
    check("PPTX file exists", os.path.isfile(pptx_path))
    if not os.path.isfile(pptx_path):
        return
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        check("Presentation has >= 5 slides", slide_count >= 5,
              f"Got {slide_count} slides")

        # Check slide content
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "
        check("Contains regional/territory reference",
              "region" in all_text or "territory" in all_text)
        check("Contains sales/revenue reference",
              "sales" in all_text or "revenue" in all_text)
        check("Contains recommendation",
              "recommend" in all_text or "next step" in all_text or "action" in all_text)
    except ImportError:
        check("python-pptx available", False)
    except Exception as e:
        check("PPTX readable", False, str(e))


def check_email():
    print("\n=== Checking Email ===")
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT m.subject, m.to_addr, m.body_text
            FROM email.sent_log sl
            JOIN email.messages m ON sl.message_id = m.id
            WHERE lower(m.subject) LIKE '%%sales%%territory%%review%%'
               OR lower(m.subject) LIKE '%%q1%%sales%%'
        """)
        rows = cur.fetchall()
        if not rows:
            cur.execute("""
                SELECT subject, to_addr, body_text FROM email.messages
                WHERE lower(subject) LIKE '%%sales%%territory%%review%%'
                   OR lower(subject) LIKE '%%q1%%sales%%'
            """)
            rows = cur.fetchall()
        check("Sales review email sent", len(rows) > 0, f"Found {len(rows)}")
        if rows:
            to_str = str(rows[0][1]).lower() if rows[0][1] else ""
            check("Email to regional_managers",
                  "regional_manager" in to_str,
                  f"To: {rows[0][1]}")
        cur.close()
        conn.close()
    except Exception as e:
        check("Email check", False, str(e))


def check_reverse_validation(workspace):
    """Verify things that should NOT exist in output."""
    print("\n=== Reverse Validation ===")

    # Excel: no negative revenue or sales values
    path = os.path.join(workspace, "Quarterly_Sales_Review.xlsx")
    if os.path.isfile(path):
        wb = openpyxl.load_workbook(path, data_only=True)
        has_negative = False
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for row in ws.iter_rows(min_row=2, values_only=True):
                for cell in row:
                    if isinstance(cell, (int, float)) and cell < -100000:
                        has_negative = True
                        break
                if has_negative:
                    break
            if has_negative:
                break
        check("No large negative values in sales Excel", not has_negative,
              "Found unexpectedly large negative revenue value")

    # Email: no sales review emails to competitor addresses
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM email.messages
            WHERE (lower(subject) LIKE '%%sales%%' OR lower(subject) LIKE '%%revenue%%')
              AND to_addr::text ILIKE '%%competitor%%'
        """)
        bad_count = cur.fetchone()[0]
        check("No sales emails to competitor addresses", bad_count == 0,
              f"Found {bad_count}")
        cur.close()
        conn.close()
    except Exception:
        pass


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    check_excel(args.agent_workspace, gt_dir)
    check_pptx(args.agent_workspace)
    check_email()
    check_reverse_validation(args.agent_workspace)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    overall = FAIL_COUNT == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
