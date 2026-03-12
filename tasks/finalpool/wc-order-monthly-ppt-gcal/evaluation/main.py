"""Evaluation for wc-order-monthly-ppt-gcal."""
import argparse
import json
import os
import sys

import psycopg2

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0

MONTH_NAMES = [
    "january", "february", "march", "april", "may", "june",
    "july", "august", "september", "october", "november", "december"
]

# Actual revenue figures from DB
EXPECTED_REVENUES = {
    "january": 5938.09, "february": 4055.26, "march": 3649.45,
    "april": 4349.71, "may": 6023.29, "june": 9555.69,
    "july": 3921.19, "august": 3873.37, "september": 5425.60,
    "october": 4864.59, "november": 1061.13, "december": 8994.67,
}


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def check_ppt(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    pptx_path = os.path.join(agent_workspace, "Monthly_Sales_Review.pptx")
    if not os.path.isfile(pptx_path):
        check("Monthly_Sales_Review.pptx exists", False, f"Not found: {pptx_path}")
        return
    check("Monthly_Sales_Review.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        check("PPT has at least 4 slides", slide_count >= 4, f"Found {slide_count} slides")

        # Extract all text from all slides
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        all_text += para.text.lower() + " "

        check("PPT contains title 'Monthly Sales'",
              "monthly sales" in all_text,
              f"Sample: {all_text[:300]}")

        # Check for month names
        months_found = [m for m in MONTH_NAMES if m in all_text]
        check("PPT contains at least 6 month names",
              len(months_found) >= 6,
              f"Found months: {months_found}")

        # Check for revenue figures (at least peak month's revenue)
        check("PPT contains June revenue or total revenue figures",
              "9555" in all_text or "9,555" in all_text or "61712" in all_text or "61,712" in all_text,
              f"Revenue sample: {all_text[:400]}")

        # Check for key insights
        check("PPT contains key insights",
              "insight" in all_text or "peak" in all_text or "revenue" in all_text,
              f"Sample: {all_text[:300]}")

    except ImportError:
        check("PPT file has content", os.path.getsize(pptx_path) > 5000,
              f"Size: {os.path.getsize(pptx_path)}")
    except Exception as e:
        check("PPT readable", False, str(e))


def check_calendar():
    print("\n=== Checking Google Calendar ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT summary, start_datetime, end_datetime
            FROM gcal.events
            WHERE LOWER(summary) LIKE '%monthly sales%'
               OR LOWER(summary) LIKE '%sales review%'
        """)
        events = cur.fetchall()
        check("'Monthly Sales Review' calendar event exists",
              len(events) >= 1,
              f"Found {len(events)} matching events")

        if events:
            event = events[0]
            start_time = str(event[1]) if event[1] else ""
            check("Event is on 2026-04-15",
                  "2026-04-15" in start_time,
                  f"Start time: {start_time}")

        cur.close()
        conn.close()
    except Exception as e:
        check("Calendar check", False, str(e))


def check_email():
    print("\n=== Checking Email ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT id, subject, to_addr, body_text
            FROM email.messages
            WHERE to_addr::text ILIKE '%ecommerce_manager@store.com%'
               OR subject ILIKE '%monthly sales%'
               OR subject ILIKE '%sales performance%'
        """)
        emails = cur.fetchall()
        check("Email sent to ecommerce_manager@store.com", len(emails) >= 1,
              "No matching email found")
        if emails:
            email = emails[0]
            subject = str(email[1]).lower() if email[1] else ""
            check("Email subject contains 'sales'",
                  "sales" in subject,
                  f"Subject: {email[1]}")
            body = str(email[3]) if email[3] else ""
            check("Email body has content", len(body) > 30,
                  f"Body length: {len(body)}")
        cur.close()
        conn.close()
    except Exception as e:
        check("Email check", False, str(e))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_ppt(args.agent_workspace)
    check_calendar()
    check_email()

    total = PASS_COUNT + FAIL_COUNT
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ===")
    if FAIL_COUNT > 0:
        print(f"{FAIL_COUNT} checks failed")
        sys.exit(1)
    else:
        print("All checks passed!")
        sys.exit(0)


if __name__ == "__main__":
    main()
