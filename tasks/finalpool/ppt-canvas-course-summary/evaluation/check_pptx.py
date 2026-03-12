"""
Check the agent's Fall2014_Teaching_Review.pptx for correct structure and content.
"""

import os

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def check_pptx(agent_workspace):
    """
    Validate Fall2014_Teaching_Review.pptx from agent workspace.
    Returns (passed_count, failed_count, error_details).
    """
    agent_file = os.path.join(agent_workspace, "Fall2014_Teaching_Review.pptx")

    if Presentation is None:
        return 0, 1, ["python-pptx not installed; cannot check PPTX"]

    if not os.path.isfile(agent_file):
        return 0, 1, [f"Agent workspace file does not exist: {agent_file}"]

    passed = 0
    failed = 0
    errors = []

    try:
        prs = Presentation(agent_file)
    except Exception as e:
        return 0, 1, [f"Error reading PPTX file: {e}"]

    slides = prs.slides
    slide_count = len(slides)

    # Check minimum slide count: 1 title + 7 courses + 1 summary = 9
    if slide_count >= 9:
        passed += 1
    else:
        failed += 1
        errors.append(f"Expected at least 9 slides, got {slide_count}")

    # Collect all text from all slides
    all_texts = []
    for slide in slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            slide_text.append(text)
        all_texts.append(slide_text)

    all_text_flat = " ".join(" ".join(sl) for sl in all_texts).lower()

    # Check title slide contains "Fall 2014"
    first_slide_text = " ".join(all_texts[0]).lower() if all_texts else ""
    if "fall 2014" in first_slide_text:
        passed += 1
    else:
        failed += 1
        errors.append(f"Title slide does not contain 'Fall 2014'. Text: {first_slide_text[:200]}")

    # Check title slide contains "teaching review" or "review"
    if "review" in first_slide_text or "teaching" in first_slide_text:
        passed += 1
    else:
        failed += 1
        errors.append(f"Title slide does not contain 'review' or 'teaching'")

    # Check that all 7 course codes or course names appear somewhere in the presentation
    course_identifiers = [
        ("aaa-2014j", "applied analytics"),
        ("bbb-2014j", "biochemistry"),
        ("ccc-2014j", "creative computing"),
        ("ddd-2014j", "data-driven"),
        ("eee-2014j", "environmental economics"),
        ("fff-2014j", "foundations of finance"),
        ("ggg-2014j", "global governance"),
    ]

    courses_found = 0
    for code, name_fragment in course_identifiers:
        if code in all_text_flat or name_fragment in all_text_flat:
            courses_found += 1
        else:
            errors.append(f"PPTX: neither '{code}' nor '{name_fragment}' found in slides")

    if courses_found == 7:
        passed += 1
    else:
        failed += 1
        errors.append(f"Only {courses_found}/7 courses found in presentation")

    # Check that summary/total info exists somewhere
    summary_keywords = ["total", "summary", "overall"]
    summary_found = any(kw in all_text_flat for kw in summary_keywords)
    if summary_found:
        passed += 1
    else:
        failed += 1
        errors.append("No summary/total information found in presentation")

    return passed, failed, errors
