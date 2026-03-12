"""
Evaluation for 12306-yf-investor-roadshow-excel-ppt-email task.

Checks:
1. Roadshow_Analysis.xlsx exists with Travel_Plan, Stock_Summary, Financial_Highlights sheets
2. Travel_Plan has G11 and G105 rows
3. Stock_Summary has AAPL and MSFT rows
4. Financial_Highlights has >= 2 rows
5. Investor_Roadshow.pptx exists with >= 4 slides and AI/finance content
6. Emails sent to both investors@fundmanager.com and shanghai_partners@finance.com
"""
import json
import os
import sys
from argparse import ArgumentParser

import psycopg2
import openpyxl

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def check_excel(agent_workspace, groundtruth_workspace="."):
    print("\n=== Check 1: Excel Roadshow_Analysis.xlsx ===")

    xlsx_path = os.path.join(agent_workspace, "Roadshow_Analysis.xlsx")
    if not os.path.exists(xlsx_path):
        record("Roadshow_Analysis.xlsx exists", False, f"Not found at {xlsx_path}")
        return
    record("Roadshow_Analysis.xlsx exists", True)

    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    except Exception as e:
        record("Excel readable", False, str(e))
        return
    record("Excel readable", True)

    sheet_names_lower = [s.lower() for s in wb.sheetnames]
    has_travel = any("travel" in s for s in sheet_names_lower)
    has_stock = any("stock" in s or "summary" in s for s in sheet_names_lower)
    has_financial = any("financial" in s or "highlight" in s for s in sheet_names_lower)

    record("Excel has Travel_Plan sheet", has_travel, f"Sheets: {wb.sheetnames}")
    record("Excel has Stock_Summary sheet", has_stock, f"Sheets: {wb.sheetnames}")
    record("Excel has Financial_Highlights sheet", has_financial, f"Sheets: {wb.sheetnames}")

    if has_travel:
        ws_name = wb.sheetnames[next(i for i, s in enumerate(sheet_names_lower) if "travel" in s)]
        ws = wb[ws_name]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if any(c for c in r)]
        record("Travel_Plan has >= 2 rows", len(data_rows) >= 2, f"Found {len(data_rows)} rows")

        all_text = " ".join(str(c) for row in rows for c in row if c).lower()
        record("Travel_Plan contains G11", "g11" in all_text, f"Content sample: {all_text[:200]}")
        record("Travel_Plan contains G105", "g105" in all_text, f"Content sample: {all_text[:200]}")

    if has_stock:
        ws_name = wb.sheetnames[next(i for i, s in enumerate(sheet_names_lower) if "stock" in s or "summary" in s)]
        ws = wb[ws_name]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if any(c for c in r)]
        record("Stock_Summary has >= 2 rows", len(data_rows) >= 2, f"Found {len(data_rows)} rows")

        all_text = " ".join(str(c) for row in rows for c in row if c).upper()
        record("Stock_Summary contains AAPL", "AAPL" in all_text, f"Content: {all_text[:200]}")
        record("Stock_Summary contains MSFT", "MSFT" in all_text, f"Content: {all_text[:200]}")

    if has_financial:
        ws_name = wb.sheetnames[next(i for i, s in enumerate(sheet_names_lower) if "financial" in s or "highlight" in s)]
        ws = wb[ws_name]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if any(c for c in r)]
        record("Financial_Highlights has >= 2 rows", len(data_rows) >= 2, f"Found {len(data_rows)} rows")

    # --- Groundtruth value comparison ---
    gt_path = os.path.join(groundtruth_workspace, "Roadshow_Analysis.xlsx")
    if not os.path.isfile(gt_path):
        record("Groundtruth xlsx exists", False, gt_path)
        return

    gt_wb = openpyxl.load_workbook(gt_path, data_only=True)
    for gt_sheet_name in gt_wb.sheetnames:
        gt_ws = gt_wb[gt_sheet_name]
        agent_ws = None
        for asn in wb.sheetnames:
            if asn.strip().lower() == gt_sheet_name.strip().lower():
                agent_ws = wb[asn]
                break
        if agent_ws is None:
            record(f"GT sheet '{gt_sheet_name}' exists in agent", False, f"Available: {wb.sheetnames}")
            continue

        gt_rows = [r for r in gt_ws.iter_rows(min_row=2, values_only=True) if any(c is not None for c in r)]
        agent_rows = [r for r in agent_ws.iter_rows(min_row=2, values_only=True) if any(c is not None for c in r)]

        record(f"GT '{gt_sheet_name}' row count", len(agent_rows) == len(gt_rows),
               f"Expected {len(gt_rows)}, got {len(agent_rows)}")

        check_indices = list(range(min(3, len(gt_rows))))
        if len(gt_rows) > 3:
            check_indices.append(len(gt_rows) - 1)
        for idx in check_indices:
            gt_row = gt_rows[idx]
            if idx < len(agent_rows):
                a_row = agent_rows[idx]
                row_ok = True
                for col_idx in range(min(len(gt_row), len(a_row) if a_row else 0)):
                    gt_val = gt_row[col_idx]
                    a_val = a_row[col_idx]
                    if gt_val is None:
                        continue
                    if isinstance(gt_val, (int, float)):
                        ok = num_close(a_val, gt_val, max(abs(gt_val) * 0.1, 1.0))
                    else:
                        ok = str_match(a_val, gt_val)
                    if not ok:
                        record(f"GT '{gt_sheet_name}' row {idx+1} col {col_idx+1}",
                               False, f"Expected {gt_val}, got {a_val}")
                        row_ok = False
                        break
                if row_ok:
                    record(f"GT '{gt_sheet_name}' row {idx+1} values match", True)
            else:
                record(f"GT '{gt_sheet_name}' row {idx+1} exists", False, "Row missing in agent")
    gt_wb.close()


def check_pptx(agent_workspace):
    print("\n=== Check 2: PPT Investor_Roadshow.pptx ===")

    pptx_path = os.path.join(agent_workspace, "Investor_Roadshow.pptx")
    if not os.path.exists(pptx_path):
        record("Investor_Roadshow.pptx exists", False, f"Not found at {pptx_path}")
        return
    record("Investor_Roadshow.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        record("PPT has >= 4 slides", slide_count >= 4, f"Found {slide_count} slides")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += " " + shape.text
        all_text_lower = all_text.lower()

        has_roadshow = any(kw in all_text_lower for kw in ["roadshow", "investor", "q1 2026", "financial"])
        record("PPT contains roadshow/investor/financial content", has_roadshow,
               f"Text sample: {all_text[:200]}")

        has_train = any(kw in all_text_lower for kw in ["g11", "g105", "shanghai", "guangzhou", "train", "travel"])
        record("PPT mentions travel itinerary", has_train,
               f"Text sample: {all_text[:200]}")

        has_finance = any(kw in all_text_lower for kw in ["aapl", "revenue", "earnings", "stock", "price"])
        record("PPT mentions financial data", has_finance,
               f"Text sample: {all_text[:200]}")

    except ImportError:
        record("python-pptx available", False, "python-pptx not installed")
    except Exception as e:
        record("PPT readable", False, str(e))


def check_emails():
    print("\n=== Check 3: Emails sent ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
    messages = cur.fetchall()
    cur.close()
    conn.close()

    def to_addresses(to_addr):
        if isinstance(to_addr, list):
            return " ".join(str(r).lower() for r in to_addr)
        elif to_addr:
            try:
                parsed = json.loads(str(to_addr))
                return " ".join(str(r).lower() for r in parsed) if isinstance(parsed, list) else str(to_addr).lower()
            except Exception:
                return str(to_addr).lower()
        return ""

    to_fundmanager = [m for m in messages if "investors@fundmanager.com" in to_addresses(m[2])]
    to_shanghai = [m for m in messages if "shanghai_partners@finance.com" in to_addresses(m[2])]

    record("Email sent to investors@fundmanager.com", len(to_fundmanager) >= 1,
           f"Total messages: {len(messages)}")
    record("Email sent to shanghai_partners@finance.com", len(to_shanghai) >= 1,
           f"Total messages: {len(messages)}")

    if to_fundmanager:
        subj, _, _, body = to_fundmanager[0]
        content = ((subj or "") + " " + (body or "")).lower()
        has_finance = any(kw in content for kw in ["roadshow", "financial", "presentation", "schedule"])
        record("Fundmanager email mentions roadshow/financial content", has_finance,
               f"Subject: {subj}")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace, args.groundtruth_workspace)
    check_pptx(args.agent_workspace)
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
