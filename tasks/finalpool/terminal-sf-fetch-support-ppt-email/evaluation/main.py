"""Evaluation for terminal-sf-fetch-support-ppt-email."""
import argparse
import json
import os
import sys

import psycopg2

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432,
          dbname=os.environ.get("PGDATABASE", "toolathlon_gym"),
          user="eigent", password="camel")

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        d = f": {str(detail)[:200]}" if detail else ""
        print(f"  [FAIL] {name}{d}")


def check_pptx(ws_path):
    """Check Support_Performance_Review.pptx."""
    print("\n=== Checking PowerPoint ===")
    path = os.path.join(ws_path, "Support_Performance_Review.pptx")
    if not os.path.isfile(path):
        check("PPT file exists", False, f"Not found: {path}")
        return
    check("PPT file exists", True)

    if Presentation is None:
        check("python-pptx available", False, "Cannot import pptx")
        return

    try:
        prs = Presentation(path)
        slides = prs.slides
        check("PPT has >= 7 slides", len(slides) >= 7, f"Found {len(slides)} slides")

        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text.lower() + " "

        check("PPT mentions benchmark", "benchmark" in all_text or "industry" in all_text)
        check("PPT mentions priority levels",
              "high" in all_text and ("medium" in all_text or "low" in all_text))
        check("PPT mentions satisfaction or CSAT",
              "satisfaction" in all_text or "csat" in all_text)
        check("PPT mentions compliance or SLA",
              "compliance" in all_text or "sla" in all_text or "compliant" in all_text)
        check("PPT mentions agents or performance",
              "agent" in all_text or "performance" in all_text)
        check("PPT mentions recommendations or takeaways",
              "recommend" in all_text or "takeaway" in all_text or "improvement" in all_text)
    except Exception as e:
        check("PPT readable", False, str(e))


def check_email():
    """Check email sent to managers."""
    print("\n=== Checking Email ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
    all_emails = cur.fetchall()

    target_email = None
    for subj, from_addr, to_addr, body in all_emails:
        if to_addr:
            recipients = []
            if isinstance(to_addr, list):
                recipients = [str(r).strip().lower() for r in to_addr]
            elif isinstance(to_addr, str):
                try:
                    parsed = json.loads(to_addr)
                    if isinstance(parsed, list):
                        recipients = [str(r).strip().lower() for r in parsed]
                    else:
                        recipients = [str(to_addr).strip().lower()]
                except (json.JSONDecodeError, TypeError):
                    recipients = [str(to_addr).strip().lower()]
            if "managers@support-team.example.com" in recipients:
                target_email = (subj, from_addr, to_addr, body)
                break

    check("Email sent to managers@support-team.example.com", target_email is not None,
          f"Total emails: {len(all_emails)}")

    if target_email:
        subj, from_addr, to_addr, body = target_email
        check("Email subject mentions benchmark or performance",
              "benchmark" in (subj or "").lower() or "performance" in (subj or "").lower(),
              f"Subject: {subj}")
        check("Email from analytics@support-team.example.com",
              "analytics@support-team.example.com" in (from_addr or "").lower(),
              f"From: {from_addr}")
        body_lower = (body or "").lower()
        check("Email body mentions satisfaction or CSAT",
              "satisfaction" in body_lower or "csat" in body_lower,
              "Expected satisfaction/CSAT in body")
        check("Email body mentions compliance or priority",
              "compliant" in body_lower or "compliance" in body_lower or "priority" in body_lower,
              "Expected compliance/priority in body")

    conn.close()


def check_reverse_validation(workspace):
    """Verify things that should NOT exist in output."""
    print("\n=== Reverse Validation ===")

    # PPT: should not have empty slides
    path = os.path.join(workspace, "Support_Performance_Review.pptx")
    if os.path.isfile(path) and Presentation is not None:
        try:
            prs = Presentation(path)
            empty_slides = 0
            for slide in prs.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text += shape.text_frame.text.strip()
                if not slide_text:
                    empty_slides += 1
            check("No empty slides in PPT", empty_slides == 0,
                  f"Found {empty_slides} empty slides")
        except Exception:
            pass

    # Email: no emails sent to wrong recipients about support performance
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM email.messages
            WHERE (lower(subject) LIKE '%%benchmark%%' OR lower(subject) LIKE '%%performance%%')
              AND to_addr::text ILIKE '%%competitor%%'
        """)
        bad_count = cur.fetchone()[0]
        check("No benchmark emails to competitor addresses", bad_count == 0,
              f"Found {bad_count}")
        cur.close()
        conn.close()
    except Exception:
        pass


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print("=" * 70)
    print("TERMINAL-SF-FETCH-SUPPORT-PPT-EMAIL - EVALUATION")
    print("=" * 70)

    check_pptx(args.agent_workspace)
    check_email()
    check_reverse_validation(args.agent_workspace)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    overall = FAIL_COUNT == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
