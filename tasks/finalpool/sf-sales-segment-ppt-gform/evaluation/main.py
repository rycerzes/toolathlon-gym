"""Evaluation for sf-sales-segment-ppt-gform."""
import argparse
import os
import sys

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    agent_ws = args.agent_workspace or task_root

    all_errors = []

    # --- Check 1: PowerPoint ---
    print("Checking PowerPoint presentation...")
    pptx_path = os.path.join(agent_ws, "Segment_Performance.pptx")
    if not os.path.exists(pptx_path):
        all_errors.append("Segment_Performance.pptx not found in agent workspace")
    else:
        from pptx import Presentation
        prs = Presentation(pptx_path)

        if len(prs.slides) < 4:
            all_errors.append(f"PPT has only {len(prs.slides)} slides, expected at least 4")

        # Collect all text from slides
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += " " + shape.text

        # Check for segment data
        if "Consumer" not in all_text:
            all_errors.append("PPT missing 'Consumer' segment")
        if "Enterprise" not in all_text:
            all_errors.append("PPT missing 'Enterprise' segment")

        # Check for revenue figure - accept various formats
        revenue_present = any(x in all_text for x in ["839609", "839,609", "839,609.20"])
        if not revenue_present:
            all_errors.append("PPT missing Consumer revenue figure (839609.20)")

        # Check title slide
        if prs.slides:
            first_slide_text = ""
            for shape in prs.slides[0].shapes:
                if hasattr(shape, "text"):
                    first_slide_text += shape.text
            if "Customer Segment" not in first_slide_text and "Segment Performance" not in first_slide_text:
                all_errors.append("First slide missing 'Customer Segment Performance Report' title")

        print(f"    PPT checks done ({len(prs.slides)} slides)")

    # --- Check 2: GForm exists ---
    print("Checking Google Form...")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT COUNT(*) FROM gform.forms WHERE LOWER(title) LIKE '%sales strategy%'")
        count = cur.fetchone()[0]
        cur.close()
        conn.close()
        if count == 0:
            all_errors.append("Google Form 'Sales Strategy Feedback' not found in gform.forms")
        else:
            print(f"    GForm found ({count} matching forms)")
    except Exception as e:
        all_errors.append(f"Error checking GForm: {e}")

    # --- Check 3: Email sent ---
    print("Checking email...")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM email.messages
            WHERE to_addr::text ILIKE '%sales_director@company.com%'
        """)
        count = cur.fetchone()[0]
        cur.close()
        conn.close()
        if count == 0:
            all_errors.append("No email sent to sales_director@company.com")
        else:
            print(f"    Email found ({count} messages)")
    except Exception as e:
        all_errors.append(f"Error checking email: {e}")

    # --- Final result ---
    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
