"""
Evaluation for wc-category-performance-ppt.
Checks:
1. Excel file Category_Data.xlsx with correct category metrics
2. PPT file Category_Review.pptx with correct slide count and content
"""
import argparse
import json
import os
import sys
from collections import defaultdict

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def get_expected_data():
    """Query WC DB for expected category performance."""
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("SELECT id, name, categories, total_sales, price, average_rating FROM wc.products")
    rows = cur.fetchall()

    cat_data = defaultdict(lambda: {'count': 0, 'prices': [], 'total_units_sold': 0, 'ratings': []})

    for r in rows:
        cats = r[2]
        if isinstance(cats, str):
            cats = json.loads(cats)
        if isinstance(cats, list):
            for c in cats:
                cat_name = c.get('name', '') if isinstance(c, dict) else str(c)
                cat_data[cat_name]['count'] += 1
                try:
                    cat_data[cat_name]['prices'].append(float(r[4]) if r[4] else 0)
                except (TypeError, ValueError):
                    pass
                cat_data[cat_name]['total_units_sold'] += int(r[3]) if r[3] else 0
                try:
                    if r[5] and float(r[5]) > 0:
                        cat_data[cat_name]['ratings'].append(float(r[5]))
                except (TypeError, ValueError):
                    pass

    results = []
    for cat in sorted(cat_data.keys()):
        d = cat_data[cat]
        avg_price = round(sum(d['prices']) / len(d['prices']), 2) if d['prices'] else 0
        avg_rating = round(sum(d['ratings']) / len(d['ratings']), 2) if d['ratings'] else 0
        results.append({
            'category': cat,
            'product_count': d['count'],
            'avg_price': avg_price,
            'total_units_sold': d['total_units_sold'],
            'avg_rating': avg_rating,
        })

    cur.close()
    conn.close()
    return results


def check_excel(agent_workspace, gt_workspace, expected):
    """Check Category_Data.xlsx."""
    print("\n=== Checking Excel ===")
    agent_file = os.path.join(agent_workspace, "Category_Data.xlsx")
    gt_file = os.path.join(gt_workspace, "Category_Data.xlsx")

    if not os.path.exists(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False
    record("Excel file exists", True)

    agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
    gt_wb = openpyxl.load_workbook(gt_file, data_only=True)

    a_rows = load_sheet_rows(agent_wb, "Category Performance")
    g_rows = load_sheet_rows(gt_wb, "Category Performance")

    if a_rows is None:
        record("Sheet 'Category Performance' exists", False, f"Sheets: {agent_wb.sheetnames}")
        return False
    record("Sheet 'Category Performance' exists", True)

    a_data = a_rows[1:] if len(a_rows) > 1 else []
    g_data = g_rows[1:] if len(g_rows) > 1 else []

    record("Category row count", len(a_data) == len(g_data),
           f"Expected {len(g_data)}, got {len(a_data)}")

    a_lookup = {}
    for row in a_data:
        if row and row[0]:
            a_lookup[str(row[0]).strip().lower()] = row

    for g_row in g_data:
        if not g_row or not g_row[0]:
            continue
        cat = str(g_row[0]).strip()
        key = cat.lower()
        a_row = a_lookup.get(key)
        if a_row is None:
            record(f"Category '{cat}' found", False, "Missing")
            continue
        record(f"Category '{cat}' found", True)

        if len(a_row) > 1 and len(g_row) > 1:
            record(f"  {cat} Product_Count",
                   num_close(a_row[1], g_row[1], 1),
                   f"Agent={a_row[1]}, GT={g_row[1]}")
        if len(a_row) > 2 and len(g_row) > 2:
            record(f"  {cat} Avg_Price",
                   num_close(a_row[2], g_row[2], 5.0),
                   f"Agent={a_row[2]}, GT={g_row[2]}")
        if len(a_row) > 3 and len(g_row) > 3:
            record(f"  {cat} Total_Units_Sold",
                   num_close(a_row[3], g_row[3], 5),
                   f"Agent={a_row[3]}, GT={g_row[3]}")
        if len(a_row) > 4 and len(g_row) > 4:
            record(f"  {cat} Avg_Rating",
                   num_close(a_row[4], g_row[4], 0.1),
                   f"Agent={a_row[4]}, GT={g_row[4]}")

    return True


def check_pptx(agent_workspace, expected):
    """Check Category_Review.pptx."""
    print("\n=== Checking PowerPoint ===")
    pptx_path = os.path.join(agent_workspace, "Category_Review.pptx")

    if not os.path.exists(pptx_path):
        record("PPTX file exists", False, f"Not found: {pptx_path}")
        return False
    record("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        num_slides = len(prs.slides)

        expected_slides = len(expected) + 2  # title + per-category + summary
        record(f"Slide count (expected ~{expected_slides})",
               abs(num_slides - expected_slides) <= 2,
               f"Got {num_slides}")

        # Check for title slide
        first_slide = prs.slides[0]
        first_text = " ".join(
            shape.text for shape in first_slide.shapes if shape.has_text_frame
        ).lower()
        record("Title slide mentions category/performance/review",
               any(w in first_text for w in ["category", "performance", "review"]),
               f"Title text: {first_text[:100]}")

        # Check that category names appear in slides
        all_slide_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_slide_text += " " + shape.text
        all_slide_lower = all_slide_text.lower()

        for cat_info in expected:
            cat = cat_info['category']
            record(f"PPT mentions '{cat}'",
                   cat.lower() in all_slide_lower,
                   f"'{cat}' not found in slides")

        # Check summary slide
        last_slide = prs.slides[-1]
        last_text = " ".join(
            shape.text for shape in last_slide.shapes if shape.has_text_frame
        ).lower()
        record("Last slide is summary",
               "summary" in last_text or "overall" in last_text,
               f"Last slide text: {last_text[:100]}")

    except ImportError:
        record("python-pptx available", False, "Cannot import pptx module")
        return False
    except Exception as e:
        record("PPTX readable", False, str(e))
        return False

    return True


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    expected = get_expected_data()
    print(f"[eval] {len(expected)} categories found")

    excel_ok = check_excel(args.agent_workspace, gt_dir, expected)
    pptx_ok = check_pptx(args.agent_workspace, expected)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    overall = PASS_COUNT > 0 and FAIL_COUNT == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")

    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
