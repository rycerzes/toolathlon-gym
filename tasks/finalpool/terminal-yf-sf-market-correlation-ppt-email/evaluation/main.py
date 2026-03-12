"""Evaluation for terminal-yf-sf-market-correlation-ppt-email.
Checks:
1. Market_Correlation_Report.pptx with 6+ slides
2. Email sent to cfo@company.com
3. correlation_analysis.py script exists
"""
import argparse
import json
import os
import sys
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def check_pptx(workspace):
    print("\n=== Check 1: Market_Correlation_Report.pptx ===")
    path = os.path.join(workspace, "Market_Correlation_Report.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(path)
        slide_count = len(prs.slides)
        check("Has at least 6 slides", slide_count >= 6, f"Found {slide_count} slides")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += " " + shape.text
        all_lower = all_text.lower()

        check("Mentions market or stock", "market" in all_lower or "stock" in all_lower,
              f"Text snippet: {all_lower[:100]}")
        check("Mentions correlation or analysis", "correlation" in all_lower or "analysis" in all_lower,
              f"Text snippet: {all_lower[:100]}")
        check("Mentions AMZN or Amazon", "amzn" in all_lower or "amazon" in all_lower,
              f"Text snippet: {all_lower[:100]}")
        check("Mentions revenue", "revenue" in all_lower,
              f"Text snippet: {all_lower[:100]}")
        check("Mentions segment or consumer", "segment" in all_lower or "consumer" in all_lower,
              f"Text snippet: {all_lower[:100]}")
        check("Has recommendations", "recommend" in all_lower,
              f"Text snippet: {all_lower[:200]}")
    except Exception as e:
        check("PPTX readable", False, str(e))


def check_email():
    print("\n=== Check 2: Email to cfo@company.com ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
    messages = cur.fetchall()

    matching = None
    for subject, from_addr, to_addr, body_text in messages:
        to_str = ""
        if isinstance(to_addr, list):
            to_str = " ".join(str(r).lower() for r in to_addr)
        elif isinstance(to_addr, str):
            try:
                parsed = json.loads(to_addr)
                to_str = " ".join(str(r).lower() for r in parsed) if isinstance(parsed, list) else str(to_addr).lower()
            except Exception:
                to_str = str(to_addr).lower()
        if "cfo@company.com" in to_str:
            matching = (subject, from_addr, to_addr, body_text)
            break

    check("Email sent to cfo@company.com", matching is not None,
          f"Messages found: {len(messages)}")

    if matching:
        subject, _, _, body_text = matching
        all_text = ((subject or "") + " " + (body_text or "")).lower()
        check("Email mentions correlation or market", "correlation" in all_text or "market" in all_text,
              f"Subject: {subject}")
        check("Email mentions revenue", "revenue" in all_text,
              f"Body snippet: {all_text[:100]}")
        check("Email mentions presentation or report", "presentation" in all_text or "report" in all_text or "powerpoint" in all_text,
              f"Body snippet: {all_text[:100]}")

    cur.close()
    conn.close()


def check_script(workspace):
    print("\n=== Check 3: correlation_analysis.py ===")
    path = os.path.join(workspace, "correlation_analysis.py")
    check("correlation_analysis.py exists", os.path.exists(path))


def check_reverse_validation(workspace):
    """Verify things that should NOT exist in the output."""
    print("\n=== Reverse Validation ===")
    # Email: no emails about correlation sent to wrong recipients
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM email.messages
            WHERE (lower(subject) LIKE '%%correlation%%' OR lower(subject) LIKE '%%market%%')
              AND to_addr::text NOT ILIKE '%%cfo%%'
              AND to_addr::text NOT ILIKE '%%company.com%%'
        """)
        wrong_emails = cur.fetchone()[0]
        check("No correlation emails to wrong recipients", wrong_emails == 0,
              f"Found {wrong_emails} misrouted emails")
        cur.close()
        conn.close()
    except Exception:
        pass

    # PPTX: should not have empty slides
    path = os.path.join(workspace, "Market_Correlation_Report.pptx")
    if os.path.exists(path):
        try:
            from pptx import Presentation
            prs = Presentation(path)
            empty_slides = 0
            for slide in prs.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text
                if not slide_text.strip():
                    empty_slides += 1
            check("No empty slides in PPTX", empty_slides == 0,
                  f"Found {empty_slides} empty slides")
        except Exception:
            pass


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_email()
    check_script(args.agent_workspace)
    check_reverse_validation(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
