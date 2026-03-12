"""Evaluation for wc-product-launch-dashboard."""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0

MARKET_DATA = {
    "TV & Home Theater": {"growth": 12.5, "size": 48200, "trend": "Strong Growth"},
    "Electronics": {"growth": 8.3, "size": 125400, "trend": "Moderate Growth"},
    "Audio": {"growth": 15.2, "size": 22800, "trend": "Strong Growth"},
    "Cameras": {"growth": 3.1, "size": 8900, "trend": "Slow Growth"},
    "Watches": {"growth": 6.7, "size": 15600, "trend": "Moderate Growth"},
    "Home Appliances": {"growth": -2.4, "size": 35100, "trend": "Declining"},
}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=5.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def safe_str(v):
    return str(v).strip() if v is not None else ""


def get_expected_performance():
    """Compute expected category performance from WC database."""
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("""
        WITH items AS (
            SELECT (elem->>'product_id')::int as product_id,
                (elem->>'quantity')::int as quantity,
                (elem->>'price')::numeric as price
            FROM wc.orders o, jsonb_array_elements(o.line_items::jsonb) elem
        )
        SELECT p.categories::jsonb->0->>'name' as category,
            ROUND(SUM(i.quantity * i.price), 2) as revenue,
            SUM(i.quantity) as units_sold,
            ROUND(SUM(i.quantity * i.price) / SUM(i.quantity), 2) as avg_price
        FROM items i JOIN wc.products p ON i.product_id = p.id
        GROUP BY p.categories::jsonb->0->>'name'
        ORDER BY SUM(i.quantity * i.price) DESC
    """)
    perf_rows = cur.fetchall()

    cur.execute("""
        SELECT p.categories::jsonb->0->>'name' as category,
            ROUND(AVG(p.average_rating::numeric), 2) as avg_rating
        FROM wc.products p
        GROUP BY p.categories::jsonb->0->>'name'
    """)
    rating_rows = cur.fetchall()
    ratings = {r[0]: float(r[1]) for r in rating_rows}

    cur.close()
    conn.close()

    result = {}
    for cat, revenue, units, avg_price in perf_rows:
        result[cat] = {
            "revenue": float(revenue),
            "units": int(units),
            "avg_price": float(avg_price),
            "avg_rating": ratings.get(cat, 0),
        }
    return result


def check_excel(agent_workspace):
    """Check Product_Review.xlsx."""
    print("\n=== Checking Product_Review.xlsx ===")

    excel_path = os.path.join(agent_workspace, "Product_Review.xlsx")
    if not os.path.isfile(excel_path):
        record("Excel file exists", False, f"Not found: {excel_path}")
        return False
    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
    except Exception as e:
        record("Excel readable", False, str(e))
        return False

    expected = get_expected_performance()
    all_ok = True

    # Check Category Performance sheet
    cp_sheet = None
    for name in wb.sheetnames:
        if "category" in name.lower() and "perform" in name.lower():
            cp_sheet = wb[name]
            break
    if cp_sheet is None:
        record("Sheet 'Category Performance' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Category Performance' exists", True)
        headers = [safe_str(cp_sheet.cell(1, c).value).lower() for c in range(1, 8)]
        record("Has Category column", any("categ" in h for h in headers))
        record("Has Revenue column", any("revenue" in h for h in headers))
        record("Has Units_Sold column", any("unit" in h for h in headers))
        record("Has Avg_Rating column", any("rating" in h for h in headers))

        rows = list(cp_sheet.iter_rows(min_row=2, values_only=True))
        record("Category Performance has 6 rows", len(rows) >= 6,
               f"Got {len(rows)} rows")

        for cat, exp in expected.items():
            found = False
            for r in rows:
                if r and r[0] and cat.lower() in safe_str(r[0]).lower():
                    found = True
                    ok_rev = num_close(r[1], exp["revenue"], 50.0)
                    record(f"{cat} Revenue ~{exp['revenue']}", ok_rev, f"Got {r[1]}")
                    if not ok_rev:
                        all_ok = False
                    ok_units = num_close(r[2], exp["units"], 5)
                    record(f"{cat} Units ~{exp['units']}", ok_units, f"Got {r[2]}")
                    if not ok_units:
                        all_ok = False
                    break
            if not found:
                record(f"{cat} found in Category Performance", False)
                all_ok = False

    # Check Market Comparison sheet
    mc_sheet = None
    for name in wb.sheetnames:
        if "market" in name.lower() and "compar" in name.lower():
            mc_sheet = wb[name]
            break
    if mc_sheet is None:
        record("Sheet 'Market Comparison' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Market Comparison' exists", True)
        mc_rows = list(mc_sheet.iter_rows(min_row=2, values_only=True))
        record("Market Comparison has data", len(mc_rows) >= 6, f"Got {len(mc_rows)}")

        for cat, mdata in MARKET_DATA.items():
            found = False
            for r in mc_rows:
                if r and r[0] and cat.lower() in safe_str(r[0]).lower():
                    found = True
                    # Check growth rate
                    growth_col = None
                    for ci, h in enumerate([safe_str(mc_sheet.cell(1, c).value).lower() for c in range(1, 7)]):
                        if "growth" in h:
                            growth_col = ci
                            break
                    if growth_col is not None:
                        ok = num_close(r[growth_col], mdata["growth"], 0.5)
                        record(f"{cat} Market_Growth_Rate={mdata['growth']}", ok,
                               f"Got {r[growth_col]}")
                        if not ok:
                            all_ok = False
                    break
            if not found:
                record(f"{cat} found in Market Comparison", False)
                all_ok = False

    # Check Opportunities sheet
    opp_sheet = None
    for name in wb.sheetnames:
        if "opportunit" in name.lower():
            opp_sheet = wb[name]
            break
    if opp_sheet is None:
        record("Sheet 'Opportunities' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Opportunities' exists", True)
        opp_rows = list(opp_sheet.iter_rows(min_row=2, values_only=True))
        record("Opportunities has data", len(opp_rows) >= 6, f"Got {len(opp_rows)}")

        headers = [safe_str(opp_sheet.cell(1, c).value).lower() for c in range(1, 7)]
        score_col = None
        for ci, h in enumerate(headers):
            if "score" in h or "opportun" in h:
                score_col = ci
                break

        if score_col is not None and opp_rows:
            # Check top opportunity (should be Audio)
            top_cat = safe_str(opp_rows[0][0])
            record("Top opportunity is Audio",
                   "audio" in top_cat.lower(),
                   f"Got {top_cat}")
            # Check score for Audio
            exp_audio_score = round(15.2 * expected.get("Audio", {}).get("avg_rating", 4.54), 2)
            ok = num_close(opp_rows[0][score_col], exp_audio_score, 2.0)
            record(f"Audio Opportunity_Score ~{exp_audio_score}", ok,
                   f"Got {opp_rows[0][score_col]}")
            if not ok:
                all_ok = False

    return all_ok


def check_pptx(agent_workspace):
    """Check Q1_Product_Review.pptx."""
    print("\n=== Checking Q1_Product_Review.pptx ===")
    from pptx import Presentation

    pptx_file = os.path.join(agent_workspace, "Q1_Product_Review.pptx")
    if not os.path.isfile(pptx_file):
        record("PPTX file exists", False, f"Not found: {pptx_file}")
        return False
    record("PPTX file exists", True)

    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        record("PPTX readable", False, str(e))
        return False

    slide_count = len(prs.slides)
    record("PPTX has >= 5 slides", slide_count >= 5, f"Got {slide_count}")

    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text.lower() + " "

    record("PPTX mentions 'product' or 'performance'",
           "product" in all_text or "performance" in all_text)
    record("PPTX mentions 'revenue'", "revenue" in all_text)
    record("PPTX mentions 'audio'", "audio" in all_text)
    record("PPTX mentions 'market' or 'trend'",
           "market" in all_text or "trend" in all_text)
    record("PPTX mentions 'opportunit'", "opportunit" in all_text)

    return True


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    excel_ok = check_excel(args.agent_workspace)
    pptx_ok = check_pptx(args.agent_workspace)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    overall = excel_ok and pptx_ok
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
