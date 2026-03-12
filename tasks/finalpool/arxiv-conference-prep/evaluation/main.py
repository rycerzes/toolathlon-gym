"""
Evaluation for arxiv-conference-prep task.

Checks:
1. PPTX file exists and has at least 7 slides
2. Title slide mentions RLHF
3. Paper slides contain the 5 target RLHF paper keywords
4. Summary/conclusion slide present
5. Calendar event with RLHF on April 10, 2026
6. Email sent to collaborators@rlhf-lab.org
"""

import os
import sys
import json
from argparse import ArgumentParser

import psycopg2
from pptx import Presentation

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

TARGET_PAPER_KEYWORDS = [
    "instructgpt",
    "follow instructions",
    "summarize from human feedback",
    "constitutional ai",
    "direct preference optimization",
    "proximal policy optimization",
]

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                texts.append(para.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
    return " ".join(texts)


def slide_title(slide):
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip()
    return ""


def check_pptx(agent_workspace):
    """Check the PowerPoint presentation."""
    print("\n=== Check 1: PowerPoint Presentation ===")

    pptx_path = os.path.join(agent_workspace, "RLHF_Conference_Report.pptx")
    if not os.path.exists(pptx_path):
        record("PPTX file exists", False, f"Not found at {pptx_path}")
        return
    record("PPTX file exists", True)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPTX file readable", False, str(e))
        return
    record("PPTX file readable", True)

    slides = list(prs.slides)
    record("At least 7 slides", len(slides) >= 7, f"Found {len(slides)} slides")

    # Check title slide mentions RLHF
    if slides:
        first_text = slide_text(slides[0]).lower()
        has_rlhf = "rlhf" in first_text or "reinforcement learning from human feedback" in first_text
        record("Title slide mentions RLHF", has_rlhf, f"First slide text: {first_text[:100]}")

    # Check last slide has summary/conclusion
    if slides:
        last_text = slide_text(slides[-1]).lower()
        has_summary = any(kw in last_text for kw in ["summary", "conclusion", "themes", "synthesis"])
        record("Last slide has summary/conclusion", has_summary, f"Last slide: {last_text[:100]}")

    # Check paper keywords appear across all slides
    all_text = " ".join(slide_text(s) for s in slides).lower()
    papers_found = sum(1 for kw in TARGET_PAPER_KEYWORDS if kw.lower() in all_text)
    record(
        "At least 3 RLHF paper keywords in slides",
        papers_found >= 3,
        f"Found {papers_found}/{len(TARGET_PAPER_KEYWORDS)} keywords",
    )

    # Check for specific paper titles
    has_instructgpt = "instruct" in all_text and "human feedback" in all_text
    has_summarize = "summarize" in all_text and "human feedback" in all_text
    has_constitutional = "constitutional" in all_text
    has_dpo = "direct preference" in all_text or "dpo" in all_text
    has_ppo = "proximal policy" in all_text or "ppo" in all_text

    papers_present = sum([has_instructgpt, has_summarize, has_constitutional, has_dpo, has_ppo])
    record(
        "At least 4 of 5 target papers discussed",
        papers_present >= 4,
        f"Found {papers_present}/5 papers",
    )


def check_calendar():
    """Verify the RLHF Summit 2026 event was created."""
    print("\n=== Check 2: Google Calendar Event ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT summary, description, start_datetime, end_datetime, location
        FROM gcal.events
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    cur.close()
    conn.close()

    # Look for RLHF event
    rlhf_event = None
    for summary, description, start_dt, end_dt, *rest in events:
        summary_lower = (summary or "").lower()
        desc_lower = (description or "").lower() if description else ""
        if "rlhf" in summary_lower or "reinforcement learning" in summary_lower or "rlhf" in desc_lower:
            location = rest[0] if rest else None
            rlhf_event = (summary, description, start_dt, end_dt, location)
            break

    record("Calendar event with RLHF exists", rlhf_event is not None,
           "No event found with RLHF in summary/description")

    if rlhf_event:
        summary, description, start_dt, end_dt, location = rlhf_event
        if start_dt is not None:
            start_date_str = start_dt.strftime("%Y-%m-%d")
            record("Calendar event on 2026-04-10", start_date_str == "2026-04-10",
                   f"Event date is {start_date_str}")

            # Check duration is roughly a full day event (8 hours)
            if end_dt:
                duration_hours = (end_dt - start_dt).total_seconds() / 3600
                record("Event duration roughly 8 hours", 6.0 <= duration_hours <= 10.0,
                       f"Got {duration_hours} hours")
        else:
            record("Calendar event on 2026-04-10", False, "start_datetime is NULL")


def check_email():
    """Verify the email was sent to collaborators."""
    print("\n=== Check 3: Email to Collaborators ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
    """)
    all_messages = cur.fetchall()
    cur.close()
    conn.close()

    # Find email with RLHF in subject
    matching_email = None
    for subject, from_addr, to_addr, body_text in all_messages:
        subject_lower = (subject or "").lower()
        if "rlhf" in subject_lower or "reinforcement learning" in subject_lower:
            matching_email = (subject, from_addr, to_addr, body_text)
            break

    record("Email with RLHF in subject exists", matching_email is not None,
           "No email found with RLHF in subject")

    if matching_email:
        subject, from_addr, to_addr, body_text = matching_email
        # Check recipient
        to_str = ""
        if isinstance(to_addr, list):
            to_str = " ".join(str(r).lower() for r in to_addr)
        elif isinstance(to_addr, str):
            try:
                parsed = json.loads(to_addr)
                if isinstance(parsed, list):
                    to_str = " ".join(str(r).lower() for r in parsed)
                else:
                    to_str = str(to_addr).lower()
            except (json.JSONDecodeError, TypeError):
                to_str = str(to_addr).lower()

        record("Email to collaborators@rlhf-lab.org",
               "collaborators@rlhf-lab.org" in to_str,
               f"Recipient: {to_addr}")

        # Check body mentions conference date
        body_lower = (body_text or "").lower()
        has_date = "april 10" in body_lower or "2026-04-10" in body_lower or "april 10, 2026" in body_lower
        record("Email body mentions conference date", has_date,
               "Date not found in email body")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_calendar()
    check_email()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if accuracy >= 80:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
