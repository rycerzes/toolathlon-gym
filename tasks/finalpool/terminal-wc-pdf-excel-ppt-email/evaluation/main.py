"""Evaluation for terminal-wc-pdf-excel-ppt-email.
Checks:
1. Recall_Impact_Assessment.xlsx with 4 sheets and correct data
2. Recall_Briefing.pptx with 5 slides
3. Emails sent to operations, finance, legal
4. recall_analysis.py and customer_impact.py scripts exist
"""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0

# Recalled product IDs (task-designed constants from recall_notice.pdf)
RECALLED_PRODUCT_IDS = [16, 72, 74]


def get_expected_from_db():
    """Query WC schema dynamically for recall impact values."""
    defaults = {
        "max_product_revenue": 9781.38,
        "total_revenue": 12050.42,
        "total_units": 51,
        "total_orders": 22,
        "unique_customers": 15,
    }
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        # Per-product revenue (to find max)
        cur.execute("""
            SELECT SUM((li->>'total')::numeric) as rev
            FROM wc.orders o, jsonb_array_elements(o.line_items) li
            WHERE (li->>'product_id')::int IN %s
              AND o.status NOT IN ('cancelled','refunded','failed')
            GROUP BY (li->>'product_id')::int
            ORDER BY rev DESC LIMIT 1
        """, (tuple(RECALLED_PRODUCT_IDS),))
        row = cur.fetchone()
        if row and row[0]:
            defaults["max_product_revenue"] = float(row[0])

        # Aggregate totals
        cur.execute("""
            SELECT
                SUM((li->>'total')::numeric) as total_revenue,
                SUM((li->>'quantity')::int) as total_units,
                COUNT(DISTINCT o.id) as total_orders,
                COUNT(DISTINCT o.customer_id) as unique_customers
            FROM wc.orders o, jsonb_array_elements(o.line_items) li
            WHERE (li->>'product_id')::int IN %s
              AND o.status NOT IN ('cancelled','refunded','failed')
        """, (tuple(RECALLED_PRODUCT_IDS),))
        row = cur.fetchone()
        if row:
            if row[0]: defaults["total_revenue"] = float(row[0])
            if row[1]: defaults["total_units"] = int(row[1])
            if row[2]: defaults["total_orders"] = int(row[2])
            if row[3]: defaults["unique_customers"] = int(row[3])

        cur.close()
        conn.close()
    except Exception as e:
        print(f"  [WARN] DB query for expected values failed, using defaults: {e}")
    return defaults


EXPECTED = get_expected_from_db()


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def num_close(a, b, tol=2.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except:
        return False


def check_excel(workspace):
    print("\n=== Check 1: Recall_Impact_Assessment.xlsx ===")
    path = os.path.join(workspace, "Recall_Impact_Assessment.xlsx")
    if not os.path.exists(path):
        check("Excel file exists", False, f"Not found at {path}")
        return
    check("Excel file exists", True)

    wb = openpyxl.load_workbook(path, data_only=True)
    sheets = wb.sheetnames
    check("Has at least 4 sheets", len(sheets) >= 4, f"Found {len(sheets)}: {sheets}")

    sheets_lower = [s.lower().replace(" ", "_") for s in sheets]

    # Affected_Products
    ap_idx = next((i for i, s in enumerate(sheets_lower) if "affected" in s or "product" in s), 0)
    ws1 = wb[sheets[ap_idx]]
    rows1 = list(ws1.iter_rows(values_only=True))
    data1 = [r for r in rows1[1:] if any(c for c in r)]
    check("Affected_Products has 3 products", len(data1) == 3, f"Found {len(data1)}")

    if rows1:
        headers = [str(c).lower() if c else "" for c in rows1[0]]
        check("Has revenue_at_risk column", any("revenue" in h for h in headers), f"Headers: {headers}")
        check("Has sku column", any("sku" in h for h in headers), f"Headers: {headers}")

    # Check product names present
    all_text1 = " ".join(str(c) for r in rows1 for c in r if c).lower()
    check("Canon M50 in Affected_Products", "canon" in all_text1 and "m50" in all_text1)
    check("JBL Flip in Affected_Products", "jbl" in all_text1 and "flip" in all_text1)
    check("ZETONES in Affected_Products", "zetones" in all_text1)

    # Check revenue values
    if data1:
        # Find revenue column index
        rev_idx = next((i for i, h in enumerate(headers) if "revenue" in h), -1)
        if rev_idx >= 0:
            revenues = [float(r[rev_idx]) for r in data1 if r[rev_idx] is not None]
            if revenues:
                check(f"Highest revenue ~{EXPECTED['max_product_revenue']:.0f}",
                      num_close(max(revenues), EXPECTED["max_product_revenue"], 50),
                      f"Got {max(revenues)}")
                total_rev = sum(revenues)
                check(f"Total revenue ~{EXPECTED['total_revenue']:.0f}",
                      num_close(total_rev, EXPECTED["total_revenue"], 100),
                      f"Got {total_rev}")

    # Customer_Impact
    ci_idx = next((i for i, s in enumerate(sheets_lower) if "customer" in s), 1)
    if ci_idx < len(sheets):
        ws2 = wb[sheets[ci_idx]]
        rows2 = list(ws2.iter_rows(values_only=True))
        data2 = [r for r in rows2[1:] if any(c for c in r)]
        check(f"Customer_Impact has {EXPECTED['unique_customers']} customers",
              num_close(len(data2), EXPECTED["unique_customers"], 2), f"Found {len(data2)}")

    # Financial_Summary
    fs_idx = next((i for i, s in enumerate(sheets_lower) if "financial" in s or "summary" in s), 2)
    if fs_idx < len(sheets):
        ws3 = wb[sheets[fs_idx]]
        rows3 = list(ws3.iter_rows(values_only=True))
        all_text3 = " ".join(str(c) for r in rows3 for c in r if c).lower()
        check("Has Total Units metric", "total" in all_text3 and "unit" in all_text3)
        check("Has Revenue at Risk metric", "revenue" in all_text3 and "risk" in all_text3)

        # Check values
        val_map = {}
        for r in rows3[1:]:
            if r and r[0] and r[1] is not None:
                val_map[str(r[0]).lower()] = r[1]

        for key, expected in [("total units", EXPECTED["total_units"]),
                              ("total orders", EXPECTED["total_orders"]),
                              ("unique customers", EXPECTED["unique_customers"])]:
            matched = [(k, v) for k, v in val_map.items() if key.split()[0] in k and key.split()[1] in k]
            if matched:
                check(f"Financial_Summary {key} ~{expected}",
                      num_close(matched[0][1], expected, 3),
                      f"Got {matched[0][1]}")
            else:
                check(f"Financial_Summary has {key}", False, f"Keys: {list(val_map.keys())}")

        rev_matched = [(k, v) for k, v in val_map.items() if "revenue" in k and "risk" in k]
        if rev_matched:
            check(f"Financial_Summary revenue ~{EXPECTED['total_revenue']:.0f}",
                  num_close(rev_matched[0][1], EXPECTED["total_revenue"], 200),
                  f"Got {rev_matched[0][1]}")

    # Timeline
    tl_idx = next((i for i, s in enumerate(sheets_lower) if "timeline" in s), 3)
    if tl_idx < len(sheets):
        ws4 = wb[sheets[tl_idx]]
        rows4 = list(ws4.iter_rows(values_only=True))
        data4 = [r for r in rows4[1:] if any(c for c in r)]
        check("Timeline has at least 4 rows", len(data4) >= 4, f"Found {len(data4)}")
        all_text4 = " ".join(str(c) for r in rows4 for c in r if c).lower()
        check("Timeline has Pending status", "pending" in all_text4)
        check("Timeline has Complete status", "complete" in all_text4)


def check_pptx(workspace):
    print("\n=== Check 2: Recall_Briefing.pptx ===")
    path = os.path.join(workspace, "Recall_Briefing.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = prs.slides
        check("Has 5 slides", len(slides) >= 5, f"Found {len(slides)}")

        # Check title slide for Major/Minor
        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text.lower() + " "

        check("Title contains 'Major Recall'", "major" in all_text and "recall" in all_text,
              f"Text snippet: {all_text[:200]}")
        total_rev_str = f"{EXPECTED['total_revenue']:.0f}"
        total_rev_comma = f"{EXPECTED['total_revenue']:,.0f}"
        check("Contains financial impact info",
              "revenue" in all_text or total_rev_str in all_text or total_rev_comma in all_text,
              f"Text snippet: {all_text[:300]}")
        check("Contains customer communication", "customer" in all_text and ("communication" in all_text or "notification" in all_text or "notify" in all_text))
        check("Contains timeline/remediation", "timeline" in all_text or "remediation" in all_text or "milestone" in all_text)
    except ImportError:
        check("python-pptx available", False, "Cannot import pptx")


def check_emails():
    print("\n=== Check 3: Emails ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        # Check email to operations
        cur.execute("""
            SELECT subject, to_addr, body_text FROM email.messages
            WHERE to_addr::text ILIKE '%%operations%%'
               OR subject ILIKE '%%recall%%alert%%'
               OR subject ILIKE '%%recall%%impact%%'
        """)
        ops_emails = cur.fetchall()
        check("Email to operations sent", len(ops_emails) >= 1, "No matching email found")
        if ops_emails:
            body = str(ops_emails[0][2]).lower() if ops_emails[0][2] else ""
            subj = str(ops_emails[0][0]).lower() if ops_emails[0][0] else ""
            check("Operations email about recall",
                  "recall" in subj or "recall" in body)

        # Check email to finance
        cur.execute("""
            SELECT subject, to_addr, body_text FROM email.messages
            WHERE to_addr::text ILIKE '%%finance%%'
               OR subject ILIKE '%%financial%%impact%%'
               OR subject ILIKE '%%recall%%financial%%'
        """)
        fin_emails = cur.fetchall()
        check("Email to finance sent", len(fin_emails) >= 1, "No matching email found")
        if fin_emails:
            combined = (str(fin_emails[0][0]) + " " + str(fin_emails[0][2])).lower()
            check("Finance email mentions revenue or financial",
                  "revenue" in combined or "financial" in combined or "impact" in combined)

        # Check email to legal
        cur.execute("""
            SELECT subject, to_addr, body_text FROM email.messages
            WHERE to_addr::text ILIKE '%%legal%%'
               OR subject ILIKE '%%customer%%communication%%'
               OR subject ILIKE '%%recall%%customer%%'
        """)
        legal_emails = cur.fetchall()
        check("Email to legal sent", len(legal_emails) >= 1, "No matching email found")
        if legal_emails:
            combined = (str(legal_emails[0][0]) + " " + str(legal_emails[0][2])).lower()
            check("Legal email mentions customers or communication",
                  "customer" in combined or "communication" in combined)

    except Exception as e:
        check("Email check", False, str(e))
    finally:
        cur.close()
        conn.close()


def check_scripts(workspace):
    print("\n=== Check 4: Python scripts ===")
    check("recall_analysis.py exists",
          os.path.exists(os.path.join(workspace, "recall_analysis.py")))
    check("customer_impact.py exists",
          os.path.exists(os.path.join(workspace, "customer_impact.py")))


def check_json_outputs(workspace):
    print("\n=== Check 5: JSON outputs ===")
    for fname in ["recall_impact.json", "customer_impact.json"]:
        path = os.path.join(workspace, fname)
        if os.path.exists(path):
            try:
                with open(path) as f:
                    data = json.load(f)
                check(f"{fname} is valid JSON", True)
            except json.JSONDecodeError:
                check(f"{fname} is valid JSON", False, "Invalid JSON")
        else:
            check(f"{fname} exists", False)


def check_reverse_validation():
    """Check that emails were NOT sent to wrong/noise recipients."""
    print("\n=== Reverse Validation ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        # Noise recipients that should NOT have received recall-related emails
        noise_recipients = ["team@company.com", "all@company.com", "accounts@supplier.com",
                            "hr@company.com"]

        cur.execute("""
            SELECT subject, to_addr FROM email.messages
            WHERE subject ILIKE '%%recall%%' OR subject ILIKE '%%impact%%'
        """)
        recall_emails = cur.fetchall()

        for subj, to_addr in recall_emails:
            to_str = json.dumps(to_addr).lower() if to_addr else ""
            for noise_addr in noise_recipients:
                if noise_addr.lower() in to_str:
                    check(f"No recall email sent to noise address {noise_addr}", False,
                          f"Subject '{subj}' sent to {to_addr}")
                    break
            else:
                continue

        # Overall check: no recall emails to noise addresses
        all_to = []
        for subj, to_addr in recall_emails:
            to_str = json.dumps(to_addr).lower() if to_addr else ""
            all_to.append(to_str)
        combined_to = " ".join(all_to)
        no_noise = not any(na.lower() in combined_to for na in noise_recipients)
        check("No recall emails to noise recipients (team@, all@, accounts@, hr@)",
              no_noise, f"Recipients found: {combined_to[:200]}")

    except Exception as e:
        check("Reverse validation (email noise)", False, str(e))
    finally:
        cur.close()
        conn.close()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_emails()
    check_scripts(args.agent_workspace)
    check_json_outputs(args.agent_workspace)
    check_reverse_validation()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
