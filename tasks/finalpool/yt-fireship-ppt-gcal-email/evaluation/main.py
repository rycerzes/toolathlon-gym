"""
Evaluation for yt-fireship-ppt-gcal-email task.

Checks:
1. Fireship_2025_Digest.pptx exists and has 7 slides (title + 6 monthly)
2. Title slide contains 'Fireship 2025 Monthly Tech Digest'
3. Monthly slides contain month names and video data
4. GCal has 6 monthly review events in Jan-Jun 2025
5. GCal event titles match pattern 'Monthly Tech Digest Review - [Month Year]'
6. Email sent to team@company.com with subject 'Fireship 2025 Monthly Digest Ready'
7. Email body lists top videos from multiple months
"""
import json
import os
import sys
from argparse import ArgumentParser

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def check_pptx(agent_workspace):
    print("\n=== Check 1: Fireship_2025_Digest.pptx ===")
    pptx_path = os.path.join(agent_workspace, "Fireship_2025_Digest.pptx")
    if not os.path.exists(pptx_path):
        record("Fireship_2025_Digest.pptx exists", False, f"Not found at {pptx_path}")
        return
    record("Fireship_2025_Digest.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPTX file readable", False, str(e))
        return
    record("PPTX file readable", True)

    slide_count = len(prs.slides)
    record("PPTX has at least 7 slides (title + 6 months)", slide_count >= 7,
           f"Found {slide_count} slides")

    # Collect all text from slides
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text)
    full_text = " ".join(all_text).lower()

    record("Title slide contains 'Fireship 2025 Monthly Tech Digest'",
           "fireship 2025 monthly tech digest" in full_text,
           "Title text not found")

    # Check months present
    months_found = sum(1 for m in ["january 2025", "february 2025", "march 2025",
                                   "april 2025", "may 2025", "june 2025"]
                       if m in full_text)
    record("Slides contain at least 5 of 6 month names", months_found >= 5,
           f"Found {months_found}/6 months")

    # Check for some video title keywords
    has_deepseek = "deepseek" in full_text or "bubble" in full_text
    has_microsoft = "microsoft" in full_text or "chip" in full_text
    has_vibe = "vibe" in full_text or "coding" in full_text
    keywords_found = sum([has_deepseek, has_microsoft, has_vibe])
    record("Slides reference top video titles (at least 2 keywords)", keywords_found >= 2,
           f"DeepSeek:{has_deepseek}, Microsoft:{has_microsoft}, Vibe:{has_vibe}")


def check_gcal():
    print("\n=== Check 2: Google Calendar events ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT summary, description, start_datetime, end_datetime
        FROM gcal.events
        WHERE start_datetime >= '2025-01-01' AND start_datetime < '2025-07-01'
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    cur.close()
    conn.close()

    digest_events = [
        e for e in events
        if "monthly tech digest review" in (e[0] or "").lower()
    ]

    record("At least 6 Monthly Tech Digest Review events in Jan-Jun 2025",
           len(digest_events) >= 6,
           f"Found {len(digest_events)} matching events out of {len(events)} total")

    if digest_events:
        # Check duration (09:00-09:30 = 30 minutes)
        summary, desc, start_dt, end_dt = digest_events[0]
        if start_dt and end_dt:
            duration_min = (end_dt - start_dt).total_seconds() / 60
            record("Events are 30 minutes (09:00-09:30)", 25 <= duration_min <= 35,
                   f"Duration: {duration_min:.0f} minutes")

        # Check event descriptions mention video titles
        all_desc = " ".join((e[1] or "") for e in digest_events).lower()
        has_video_mention = any(kw in all_desc for kw in
                                ["deepseek", "microsoft", "vibe", "programming", "google", "cli"])
        record("Event descriptions mention top video titles", has_video_mention,
               "No video title keywords found in event descriptions")

        # Check months covered - should have Jan-Jun
        months_in_titles = set()
        for ev in digest_events:
            title_lower = (ev[0] or "").lower()
            for m in ["january", "february", "march", "april", "may", "june"]:
                if m in title_lower:
                    months_in_titles.add(m)
        record("Events cover at least 5 different months",
               len(months_in_titles) >= 5, f"Months found: {months_in_titles}")


def check_email():
    print("\n=== Check 3: Email to team@company.com ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT to_addr, subject, body_text FROM email.messages
        WHERE to_addr::text ILIKE '%team@company.com%'
        AND subject ILIKE '%fireship%digest%'
        ORDER BY id DESC LIMIT 5
    """)
    emails = cur.fetchall()
    cur.close()
    conn.close()

    record("Email to team@company.com with Fireship digest subject exists",
           len(emails) > 0,
           "No matching email found")

    if emails:
        to_addr, subject, body = emails[0]
        record("Email subject contains 'Fireship 2025 Monthly Digest Ready'",
               "fireship 2025 monthly digest ready" in subject.lower(),
               f"Subject: {subject}")

        body_lower = (body or "").lower()
        # Check that top video titles are mentioned
        has_deepseek = "deepseek" in body_lower or "bubble" in body_lower
        has_microsoft = "microsoft" in body_lower
        has_vibe = "vibe" in body_lower
        has_programming = "programming" in body_lower
        has_google = "google" in body_lower
        mentions = sum([has_deepseek, has_microsoft, has_vibe, has_programming, has_google])
        record("Email body mentions at least 4 top video topic keywords",
               mentions >= 4,
               f"Keywords found: deepseek={has_deepseek}, microsoft={has_microsoft}, "
               f"vibe={has_vibe}, programming={has_programming}, google={has_google}")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_gcal()
    check_email()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
