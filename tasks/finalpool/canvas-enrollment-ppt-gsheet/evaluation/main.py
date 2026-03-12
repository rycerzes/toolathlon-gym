"""Evaluation for canvas-enrollment-ppt-gsheet.

Blocking checks: Enrollment_Overview.xlsx and Enrollment_Overview.pptx.
Non-blocking: Google Sheet DB check.
"""
import argparse
import os
import sys
import openpyxl
from pptx import Presentation


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")
    all_errors = []

    # ---- Check Excel ----
    agent_excel = os.path.join(args.agent_workspace, "Enrollment_Overview.xlsx")
    gt_excel = os.path.join(gt_dir, "Enrollment_Overview.xlsx")

    if not os.path.exists(agent_excel):
        all_errors.append("Agent output Enrollment_Overview.xlsx not found")
    elif not os.path.exists(gt_excel):
        all_errors.append("Groundtruth Enrollment_Overview.xlsx not found")
    else:
        agent_wb = openpyxl.load_workbook(agent_excel, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_excel, data_only=True)

        # Check Enrollment Details
        print("  Checking Enrollment Details...")
        a_rows = load_sheet_rows(agent_wb, "Enrollment Details")
        g_rows = load_sheet_rows(gt_wb, "Enrollment Details")
        if a_rows is None:
            all_errors.append("Sheet 'Enrollment Details' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Enrollment Details' not found in groundtruth")
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []
            if abs(len(a_data) - len(g_data)) > 2:
                all_errors.append(f"Enrollment Details row count: agent={len(a_data)}, expected={len(g_data)}")

            # Match by course code (col 1)
            a_lookup = {}
            for row in a_data:
                if row and len(row) > 1 and row[1] is not None:
                    a_lookup[str(row[1]).strip().lower()] = row
            for g_row in g_data:
                if not g_row or len(g_row) < 2 or g_row[1] is None:
                    continue
                key = str(g_row[1]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    all_errors.append(f"Missing course: {g_row[0]} ({g_row[1]})")
                    continue
                # Col 2: Total_Enrollments
                if len(a_row) > 2 and len(g_row) > 2:
                    if not num_close(a_row[2], g_row[2], 10):
                        all_errors.append(f"{key}.Total_Enrollments: {a_row[2]} vs {g_row[2]} (tol=10)")
                # Col 3: Students
                if len(a_row) > 3 and len(g_row) > 3:
                    if not num_close(a_row[3], g_row[3], 10):
                        all_errors.append(f"{key}.Students: {a_row[3]} vs {g_row[3]} (tol=10)")
            if not all_errors:
                print("    PASS")

        # Check Summary sheet
        print("  Checking Summary...")
        a_rows = load_sheet_rows(agent_wb, "Summary")
        g_rows = load_sheet_rows(gt_wb, "Summary")
        prev_errors = len(all_errors)
        if a_rows is None:
            all_errors.append("Sheet 'Summary' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Summary' not found in groundtruth")
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []

            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None:
                    a_lookup[str(row[0]).strip().lower()] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = str(g_row[0]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    all_errors.append(f"Missing summary metric: {g_row[0]}")
                    continue
                if len(a_row) > 1 and len(g_row) > 1:
                    # For string values, use str_match; for numbers, use num_close
                    try:
                        float(g_row[1])
                        if not num_close(a_row[1], g_row[1], 50):
                            all_errors.append(f"Summary.{key}: {a_row[1]} vs {g_row[1]} (tol=50)")
                    except (TypeError, ValueError):
                        if not str_match(a_row[1], g_row[1]):
                            all_errors.append(f"Summary.{key}: {a_row[1]} vs {g_row[1]}")
            new_errors = len(all_errors) - prev_errors
            if new_errors == 0:
                print("    PASS")

    # ---- Check PowerPoint ----
    agent_ppt = os.path.join(args.agent_workspace, "Enrollment_Overview.pptx")
    if not os.path.exists(agent_ppt):
        all_errors.append("Agent output Enrollment_Overview.pptx not found")
    else:
        print("  Checking Enrollment_Overview.pptx...")
        prs = Presentation(agent_ppt)
        slides = list(prs.slides)
        if len(slides) < 4:
            all_errors.append(f"PPT has {len(slides)} slides, expected at least 4")
        else:
            # Check title slide
            title_text = ""
            for shape in slides[0].shapes:
                if shape.has_text_frame:
                    title_text += shape.text_frame.text.lower() + " "
            if "enrollment" not in title_text:
                all_errors.append(f"Title slide missing 'enrollment'. Found: {title_text[:100]}")

            # Check all PPT text for key content
            all_ppt_text = ""
            for slide in slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        all_ppt_text += shape.text_frame.text.lower() + " "

            if "top 5" not in all_ppt_text and "top five" not in all_ppt_text:
                all_errors.append("PPT missing 'Top 5' courses slide")

            if "distribution" not in all_ppt_text:
                all_errors.append("PPT missing 'Distribution' slide")

            # Check for actual course names
            if "creative computing" not in all_ppt_text:
                all_errors.append("PPT missing top course: Creative Computing")

        if not any("ppt" in e.lower() or "slide" in e.lower() for e in all_errors):
            print("    PASS")

    # ---- Non-blocking GSheet check ----
    print("  Non-blocking: Google Sheet DB check...")
    try:
        import psycopg2
        conn = psycopg2.connect(host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym",
                                user="eigent", password="camel")
        cur = conn.cursor()
        cur.execute("SELECT COUNT(*) FROM gsheet.spreadsheets")
        count = cur.fetchone()[0]
        print(f"    [INFO] Found {count} spreadsheet(s) (non-blocking)")
        cur.close()
        conn.close()
    except Exception as e:
        print(f"    [INFO] GSheet check skipped: {e} (non-blocking)")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:15]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
