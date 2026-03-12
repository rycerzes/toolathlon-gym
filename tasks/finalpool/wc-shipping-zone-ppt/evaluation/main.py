"""Evaluation for wc-shipping-zone-ppt."""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB = {"host": os.environ.get("PGHOST", "localhost"), "port": 5432, "dbname": "toolathlon_gym", "user": "eigent", "password": "camel"}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=0.5):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def get_expected_zone_data():
    """Compute expected zone performance from read-only DB."""
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    cur.execute("""
        SELECT
          CASE
            WHEN shipping->>'state' = 'CA' AND shipping->>'country' = 'US' THEN 'California'
            WHEN shipping->>'country' = 'US' THEN 'Domestic US'
            ELSE 'International'
          END as zone_name,
          COUNT(*) as order_count,
          ROUND(SUM(shipping_total)::numeric, 2) as total_shipping,
          ROUND(AVG(shipping_total)::numeric, 2) as avg_shipping
        FROM wc.orders
        GROUP BY 1
        ORDER BY COUNT(*) DESC
    """)
    rows = cur.fetchall()
    cur.close()
    conn.close()
    return rows


def check_excel(agent_workspace, groundtruth_workspace):
    """Check Shipping_Zone_Report.xlsx."""
    print("\n=== Checking Shipping_Zone_Report.xlsx ===")

    agent_file = os.path.join(agent_workspace, "Shipping_Zone_Report.xlsx")
    if not os.path.isfile(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False
    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(agent_file, data_only=True)
    except Exception as e:
        record("Excel readable", False, str(e))
        return False

    all_ok = True
    expected_zones = get_expected_zone_data()

    # Check Zone Performance sheet
    zp_sheet = None
    for name in wb.sheetnames:
        if "zone" in name.lower() and "perform" in name.lower():
            zp_sheet = wb[name]
            break
    if zp_sheet is None:
        record("Sheet 'Zone Performance' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Zone Performance' exists", True)
        rows = list(zp_sheet.iter_rows(min_row=2, values_only=True))
        record("Zone Performance has data rows", len(rows) >= 1, f"Got {len(rows)} rows")

        for ez in expected_zones:
            zone_name, exp_count, exp_total, exp_avg = ez
            found = False
            for r in rows:
                if r and r[0] and zone_name.lower() in str(r[0]).lower():
                    found = True
                    ok_count = num_close(r[1], exp_count, 2)
                    record(f"'{zone_name}' Order_Count", ok_count,
                           f"Expected {exp_count}, got {r[1]}")
                    if not ok_count:
                        all_ok = False
                    ok_total = num_close(r[2], exp_total, 5.0)
                    record(f"'{zone_name}' Total_Shipping_Cost", ok_total,
                           f"Expected {exp_total}, got {r[2]}")
                    if not ok_total:
                        all_ok = False
                    ok_avg = num_close(r[3], exp_avg, 0.5)
                    record(f"'{zone_name}' Avg_Shipping_Cost", ok_avg,
                           f"Expected {exp_avg}, got {r[3]}")
                    if not ok_avg:
                        all_ok = False
                    break
            if not found:
                record(f"Zone '{zone_name}' found in sheet", False, "Missing")
                all_ok = False

    # Check Summary sheet
    sum_sheet = None
    for name in wb.sheetnames:
        if "summary" in name.lower():
            sum_sheet = wb[name]
            break
    if sum_sheet is None:
        record("Sheet 'Summary' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Summary' exists", True)
        summary = {}
        for row in sum_sheet.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                summary[str(row[0]).strip().lower()] = row[1]

        total_orders = sum(ez[1] for ez in expected_zones)
        total_revenue = sum(ez[2] for ez in expected_zones)
        zones_count = len(expected_zones)

        for key in summary:
            if "total_order" in key or "total order" in key:
                ok = num_close(summary[key], total_orders, 2)
                record("Summary Total_Orders", ok,
                       f"Expected {total_orders}, got {summary[key]}")
                if not ok:
                    all_ok = False
            elif "total_shipping" in key or "shipping_revenue" in key or "total shipping" in key:
                ok = num_close(summary[key], total_revenue, 5.0)
                record("Summary Total_Shipping_Revenue", ok,
                       f"Expected {total_revenue}, got {summary[key]}")
                if not ok:
                    all_ok = False
            elif "zones_count" in key or "zone" in key and "count" in key:
                ok = num_close(summary[key], zones_count, 0)
                record("Summary Zones_Count", ok,
                       f"Expected {zones_count}, got {summary[key]}")
                if not ok:
                    all_ok = False

    return all_ok


def check_pptx(agent_workspace):
    """Check Shipping_Review.pptx."""
    print("\n=== Checking Shipping_Review.pptx ===")
    from pptx import Presentation

    pptx_file = os.path.join(agent_workspace, "Shipping_Review.pptx")
    if not os.path.isfile(pptx_file):
        record("PPTX file exists", False, f"Not found: {pptx_file}")
        return False
    record("PPTX file exists", True)

    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        record("PPTX readable", False, str(e))
        return False

    slide_count = len(prs.slides)
    record("PPTX has >= 4 slides", slide_count >= 4, f"Got {slide_count}")

    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text.lower() + " "

    record("PPTX mentions 'shipping'", "shipping" in all_text,
           "No mention of 'shipping' in slides")
    record("PPTX mentions 'performance' or 'review'",
           "performance" in all_text or "review" in all_text,
           "No mention of 'performance' or 'review'")

    return True


def check_gsheet():
    """Check Google Sheet with shipping data."""
    print("\n=== Checking Google Sheet ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("SELECT id, title FROM gsheet.spreadsheets WHERE title ILIKE '%shipping%'")
    rows = cur.fetchall()
    if not rows:
        record("GSheet with 'shipping' in title", False, "No matching spreadsheet found")
        cur.close()
        conn.close()
        return False
    record("GSheet with 'shipping' in title", True)

    ss_id = rows[0][0]
    cur.execute("SELECT id FROM gsheet.sheets WHERE spreadsheet_id = %s", (ss_id,))
    sheets = cur.fetchall()
    record("GSheet has at least one sheet", len(sheets) >= 1, f"Got {len(sheets)}")

    if sheets:
        sheet_id = sheets[0][0]
        cur.execute("SELECT COUNT(DISTINCT row_index) FROM gsheet.cells WHERE spreadsheet_id = %s AND sheet_id = %s",
                    (ss_id, sheet_id))
        row_count = cur.fetchone()[0]
        record("GSheet has data rows", row_count >= 2, f"Got {row_count} rows")

    cur.close()
    conn.close()
    return True


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    excel_ok = check_excel(args.agent_workspace, gt_dir)
    pptx_ok = check_pptx(args.agent_workspace)

    db_fail_before = FAIL_COUNT
    gsheet_ok = check_gsheet()
    db_failures = FAIL_COUNT - db_fail_before

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    if db_failures > 0:
        print(f"  WARNING: {db_failures} DB checks failed (not blocking)")

    overall = excel_ok and pptx_ok
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
