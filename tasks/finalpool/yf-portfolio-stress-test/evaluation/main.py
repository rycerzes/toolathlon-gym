"""Evaluation for yf-portfolio-stress-test."""
import argparse, os, sys


def num_close(a, b, abs_tol=1.0, rel_tol=0.05):
    try:
        a_f, b_f = float(a), float(b)
        return abs(a_f - b_f) <= max(abs_tol, abs(b_f) * rel_tol)
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def check_excel(agent_workspace):
    errors = []
    import openpyxl
    path = os.path.join(agent_workspace, "Stress_Test.xlsx")
    if not os.path.exists(path):
        return ["Stress_Test.xlsx not found"]
    try:
        wb = openpyxl.load_workbook(path, data_only=True)

        # Scenario Analysis
        rows = load_sheet_rows(wb, "Scenario Analysis")
        if rows is None:
            errors.append("Sheet 'Scenario Analysis' not found")
        else:
            data_rows = [r for r in rows[1:] if r and r[0] is not None]
            if len(data_rows) < 4:
                errors.append(f"Scenario Analysis has {len(data_rows)} rows, expected 4")
            scenarios = {str(r[0]).strip().lower(): r for r in data_rows if r[0]}
            if "market crash" in scenarios:
                r = scenarios["market crash"]
                if len(r) > 3 and not num_close(r[3], -27.75):
                    errors.append(f"Market Crash loss={r[3]}, expected ~-27.75")
                if len(r) > 4 and str(r[4]).strip().lower() != "yes":
                    errors.append(f"Market Crash Exceeds_Limit={r[4]}, expected Yes")

        # Position Impact
        rows2 = load_sheet_rows(wb, "Position Impact")
        if rows2 is None:
            errors.append("Sheet 'Position Impact' not found")
        else:
            data_rows2 = [r for r in rows2[1:] if r and r[0] is not None]
            if len(data_rows2) < 24:
                errors.append(f"Position Impact has {len(data_rows2)} rows, expected 24")

        # Risk Summary
        rows3 = load_sheet_rows(wb, "Risk Summary")
        if rows3 is None:
            errors.append("Sheet 'Risk Summary' not found")
        else:
            data_rows3 = [r for r in rows3[1:] if r and r[0] is not None]
            lookup = {str(r[0]).strip().lower(): r[1] for r in data_rows3 if r[0]}
            if "worst_case_loss" in lookup:
                if not num_close(lookup["worst_case_loss"], 27.75):
                    errors.append(f"Worst_Case_Loss={lookup['worst_case_loss']}, expected ~27.75")
            else:
                errors.append("Worst_Case_Loss not found")
            if "scenarios_exceeding_limit" in lookup:
                if not num_close(lookup["scenarios_exceeding_limit"], 1, abs_tol=0):
                    errors.append(f"Scenarios_Exceeding_Limit={lookup['scenarios_exceeding_limit']}, expected 1")
            if "expected_loss" in lookup:
                if not num_close(lookup["expected_loss"], -5.63):
                    errors.append(f"Expected_Loss={lookup['expected_loss']}, expected ~-5.63")
            if "risk_rating" in lookup:
                if str(lookup["risk_rating"]).strip().lower() != "elevated":
                    errors.append(f"Risk_Rating={lookup['risk_rating']}, expected Elevated")

    except Exception as e:
        errors.append(f"Error reading Excel: {e}")
    return errors


def check_pptx(agent_workspace):
    errors = []
    path = os.path.join(agent_workspace, "Risk_Presentation.pptx")
    if not os.path.exists(path):
        return ["Risk_Presentation.pptx not found"]
    try:
        from pptx import Presentation
        prs = Presentation(path)
        if len(prs.slides) < 5:
            errors.append(f"Presentation has {len(prs.slides)} slides, expected at least 5")
        # Check title slide
        first_slide = prs.slides[0]
        title_text = ""
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                title_text += shape.text_frame.text.lower()
        if "stress test" not in title_text and "portfolio" not in title_text:
            errors.append("Title slide does not mention stress test or portfolio")
    except Exception as e:
        errors.append(f"Error reading PPTX: {e}")
    return errors


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    agent_ws = args.agent_workspace or os.path.join(os.path.dirname(__file__), "..", "groundtruth_workspace")

    all_errors = []

    print("  Checking Excel file...")
    errs = check_excel(agent_ws)
    if errs:
        all_errors.extend(errs)
        for e in errs[:5]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking PowerPoint...")
    errs = check_pptx(agent_ws)
    if errs:
        all_errors.extend(errs)
        for e in errs[:3]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
