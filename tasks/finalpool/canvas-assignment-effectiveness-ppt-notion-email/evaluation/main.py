"""
Evaluation script for canvas-assignment-effectiveness-ppt-notion-email task.

Checks:
1. Excel file Assessment_Effectiveness.xlsx with 3 sheets and correct data
2. PowerPoint Curriculum_Review.pptx with 6+ slides
3. Notion database "Assignment Improvement Tracker" with revision entries
4. Email to curriculum_committee@university.edu
"""
import argparse
import json
import os
import sys

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def di_close(a, b, tol=0.05):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


# ============================================================
# Check 1: Excel
# ============================================================
def check_excel(agent_workspace, gt_workspace):
    print("\n=== Checking Excel ===")
    xlsx_path = os.path.join(agent_workspace, "Assessment_Effectiveness.xlsx")
    if not os.path.isfile(xlsx_path):
        check("Assessment_Effectiveness.xlsx exists", False, f"Not found: {xlsx_path}")
        return

    check("Assessment_Effectiveness.xlsx exists", True)

    try:
        import openpyxl
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        sheet_names = wb.sheetnames

        # Check sheets exist
        check("Sheet 'Assignment Metrics' exists",
              any("assignment" in s.lower() and "metric" in s.lower() for s in sheet_names),
              f"Sheets: {sheet_names}")
        check("Sheet 'Course Summary' exists",
              any("course" in s.lower() and "summary" in s.lower() for s in sheet_names),
              f"Sheets: {sheet_names}")
        check("Sheet 'Revision Needed' exists",
              any("revision" in s.lower() for s in sheet_names),
              f"Sheets: {sheet_names}")

        # Check Assignment Metrics sheet
        metrics_ws = None
        for s in sheet_names:
            if "assignment" in s.lower() and "metric" in s.lower():
                metrics_ws = wb[s]
                break
        if not metrics_ws:
            metrics_ws = wb[sheet_names[0]]

        rows = list(metrics_ws.iter_rows(values_only=True))
        check("Assignment Metrics has header + data rows",
              len(rows) >= 40,
              f"Found {len(rows)} rows (expected ~52)")

        # Check data content
        if len(rows) > 1:
            all_text = " ".join(str(c) for row in rows for c in row if c).lower()
            check("Contains Fall 2014 course names",
                  "applied analytics" in all_text or "biochemistry" in all_text,
                  f"Sample: {all_text[:200]}")
            check("Contains effectiveness labels",
                  "good" in all_text or "acceptable" in all_text,
                  f"Sample: {all_text[:200]}")

        # Compare with groundtruth if available
        gt_xlsx = os.path.join(gt_workspace, "Assessment_Effectiveness.xlsx")
        if os.path.isfile(gt_xlsx):
            gt_wb = openpyxl.load_workbook(gt_xlsx, data_only=True)
            gt_metrics = None
            for s in gt_wb.sheetnames:
                if "assignment" in s.lower() and "metric" in s.lower():
                    gt_metrics = gt_wb[s]
                    break
            if not gt_metrics:
                gt_metrics = gt_wb[gt_wb.sheetnames[0]]

            gt_rows = list(gt_metrics.iter_rows(values_only=True))
            agent_rows = rows

            # Check row count matches (within tolerance)
            check("Assignment count matches groundtruth",
                  abs(len(agent_rows) - len(gt_rows)) <= 3,
                  f"Agent: {len(agent_rows)}, GT: {len(gt_rows)}")

            # Spot-check a few DI values
            gt_data = {}
            for row in gt_rows[1:]:
                if row[0] and row[1]:
                    key = (str(row[0]).strip(), str(row[1]).strip())
                    gt_data[key] = row

            matches = 0
            total_checked = 0
            for row in agent_rows[1:]:
                if row[0] and row[1]:
                    key = (str(row[0]).strip(), str(row[1]).strip())
                    if key in gt_data:
                        total_checked += 1
                        gt_row = gt_data[key]
                        # Check DI (col 7, 0-indexed)
                        if len(row) > 7 and len(gt_row) > 7:
                            if di_close(row[7], gt_row[7], tol=0.1):
                                matches += 1
            if total_checked > 0:
                match_rate = matches / total_checked
                check(f"DI values match groundtruth (>= 70%)",
                      match_rate >= 0.7,
                      f"{matches}/{total_checked} = {match_rate:.1%}")
            else:
                check("DI spot-check feasible", False, "No matching rows found")

        # Check Course Summary sheet
        summary_ws = None
        for s in sheet_names:
            if "course" in s.lower() and "summary" in s.lower():
                summary_ws = wb[s]
                break
        if summary_ws:
            summary_rows = list(summary_ws.iter_rows(values_only=True))
            check("Course Summary has 7 courses + header",
                  len(summary_rows) >= 7,
                  f"Found {len(summary_rows)} rows")

        # Check Revision Needed sheet
        revision_ws = None
        for s in sheet_names:
            if "revision" in s.lower():
                revision_ws = wb[s]
                break
        if revision_ws:
            revision_rows = list(revision_ws.iter_rows(values_only=True))
            check("Revision Needed has entries",
                  len(revision_rows) >= 10,
                  f"Found {len(revision_rows)} rows (expected ~18)")

    except ImportError:
        check("openpyxl available", False, "Cannot parse Excel without openpyxl")
    except Exception as e:
        check("Excel parsing", False, str(e))


# ============================================================
# Check 2: PowerPoint
# ============================================================
def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    pptx_path = os.path.join(agent_workspace, "Curriculum_Review.pptx")
    if not os.path.isfile(pptx_path):
        check("Curriculum_Review.pptx exists", False, f"Not found: {pptx_path}")
        return

    check("Curriculum_Review.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        check("PPT has at least 6 slides", slide_count >= 6,
              f"Found {slide_count} slides")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        check("PPT mentions discrimination index or DI",
              "discrimination" in all_text or " di " in all_text or "di=" in all_text,
              "Missing DI content")
        check("PPT mentions methodology or formula",
              "methodol" in all_text or "formula" in all_text or "27%" in all_text,
              "Missing methodology")
        check("PPT mentions revision or improvement",
              "revision" in all_text or "improv" in all_text or "needs" in all_text,
              "Missing revision content")
        check("PPT mentions recommendations",
              "recommend" in all_text or "suggest" in all_text or "action" in all_text,
              "Missing recommendations")
        check("PPT mentions Fall 2014 or course names",
              "fall 2014" in all_text or "2014" in all_text,
              "Missing term reference")

    except ImportError:
        size = os.path.getsize(pptx_path)
        check("PPT file has content (>5KB)", size > 5000, f"Size: {size}")
    except Exception as e:
        check("PPT parsing", False, str(e))


# ============================================================
# Check 3: Notion
# ============================================================
def check_notion():
    print("\n=== Checking Notion ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        # Check for database
        cur.execute("""
            SELECT id, title, properties FROM notion.databases
            WHERE title::text ILIKE '%%improvement%%tracker%%'
               OR title::text ILIKE '%%assignment%%tracker%%'
        """)
        dbs = cur.fetchall()
        check("Notion database 'Assignment Improvement Tracker' exists",
              len(dbs) >= 1,
              f"Found {len(dbs)} matching databases")

        if dbs:
            db_id = dbs[0][0]
            props_raw = dbs[0][2]
            if isinstance(props_raw, str):
                props = json.loads(props_raw)
            else:
                props = props_raw

            # Check properties
            prop_names = [k.lower() for k in props.keys()] if props else []
            check("Database has 'Assignment' property",
                  any("assignment" in p for p in prop_names),
                  f"Properties: {prop_names}")
            check("Database has 'Course' property",
                  any("course" in p for p in prop_names),
                  f"Properties: {prop_names}")
            check("Database has 'Status' property",
                  any("status" in p for p in prop_names),
                  f"Properties: {prop_names}")

            # Check pages in this database
            cur.execute("""
                SELECT id, properties FROM notion.pages
                WHERE parent::text LIKE %s
                  AND (archived IS NULL OR archived = false)
                  AND (in_trash IS NULL OR in_trash = false)
            """, (f'%{db_id}%',))
            pages = cur.fetchall()
            check("Notion has entries for revision-needed assignments (>= 10)",
                  len(pages) >= 10,
                  f"Found {len(pages)} pages (expected ~17)")

            # Check some page content
            if pages:
                all_props_text = " ".join(
                    json.dumps(p[1]) if isinstance(p[1], dict) else str(p[1])
                    for p in pages
                ).lower()
                check("Pages contain assignment names",
                      "cma" in all_props_text or "tma" in all_props_text or "final" in all_props_text,
                      f"Sample: {all_props_text[:200]}")
                check("Pages contain course references",
                      "foundation" in all_props_text or "creative" in all_props_text or "finance" in all_props_text,
                      f"Sample: {all_props_text[:200]}")

        cur.close()
        conn.close()
    except Exception as e:
        check("Notion check", False, str(e))


# ============================================================
# Check 4: Email
# ============================================================
def check_email():
    print("\n=== Checking Email ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        cur.execute("""
            SELECT id, subject, from_addr, to_addr, body_text
            FROM email.messages
            WHERE to_addr::text ILIKE '%%curriculum_committee@university.edu%%'
               OR to_addr::text ILIKE '%%curriculum%%committee%%'
               OR subject ILIKE '%%assignment%%effectiveness%%'
               OR subject ILIKE '%%fall 2014%%assignment%%'
        """)
        emails = cur.fetchall()
        check("Email to curriculum_committee@university.edu found",
              len(emails) >= 1,
              f"Found {len(emails)} matching emails")

        if emails:
            email = emails[0]
            subject = str(email[1] or "").lower()
            body = str(email[4] or "").lower()

            check("Email subject mentions assignment or effectiveness",
                  "assignment" in subject or "effectiveness" in subject or "fall 2014" in subject,
                  f"Subject: {email[1]}")
            check("Email body has substantive content",
                  len(body) > 50,
                  f"Body length: {len(body)}")
            check("Email body mentions key findings",
                  any(term in body for term in ["revision", "poor", "completion", "discrimination", "good", "acceptable"]),
                  f"Body sample: {body[:200]}")

        cur.close()
        conn.close()
    except Exception as e:
        check("Email check", False, str(e))


# ============================================================
# Main
# ============================================================
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False, default="")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    gt_ws = args.groundtruth_workspace or args.agent_workspace

    check_excel(args.agent_workspace, gt_ws)
    check_pptx(args.agent_workspace)
    check_notion()
    check_email()

    total = PASS_COUNT + FAIL_COUNT
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ===")

    if args.res_log_file:
        result = {"passed": PASS_COUNT, "failed": FAIL_COUNT, "total": total}
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if FAIL_COUNT > 0:
        print(f"{FAIL_COUNT} checks failed")
        sys.exit(1)
    else:
        print("All checks passed!")
        sys.exit(0)


if __name__ == "__main__":
    main()
