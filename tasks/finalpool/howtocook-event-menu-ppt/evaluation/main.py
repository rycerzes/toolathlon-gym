"""
Evaluation script for howtocook-event-menu-ppt task.

Checks:
1. PPT exists with >=5 slides, title slide has "dinner" or "menu", course types present
2. Excel exists with correct sheet structure
Since HowToCook data is dynamic, we evaluate structure not specific values.
"""

import argparse
import json
import os
import sys


def check_pptx(workspace):
    """Check PowerPoint file structure and content."""
    from pptx import Presentation

    pptx_path = os.path.join(workspace, "Event_Menu_Presentation.pptx")
    if not os.path.exists(pptx_path):
        return False, "Event_Menu_Presentation.pptx not found"

    prs = Presentation(pptx_path)
    slides = list(prs.slides)

    if len(slides) < 5:
        return False, f"PPT has {len(slides)} slides, expected at least 5"

    # Collect all text from all slides
    all_text = ""
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    all_text += paragraph.text.lower() + " "
            if hasattr(shape, "text"):
                all_text += shape.text.lower() + " "

    # Check title slide contains dinner or menu
    title_slide_text = ""
    for shape in slides[0].shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                title_slide_text += paragraph.text.lower() + " "
        if hasattr(shape, "text"):
            title_slide_text += shape.text.lower() + " "

    if "dinner" not in title_slide_text and "menu" not in title_slide_text:
        return False, "Title slide does not contain 'dinner' or 'menu'"

    # Check course types are present somewhere in the presentation
    course_types = ["appetizer", "main", "dessert"]
    missing = [c for c in course_types if c not in all_text]
    if missing:
        return False, f"Missing course types in PPT: {missing}"

    return True, "PPT check passed"


def check_excel(workspace):
    """Check Excel file structure."""
    from openpyxl import load_workbook

    xlsx_path = os.path.join(workspace, "Menu_Budget.xlsx")
    if not os.path.exists(xlsx_path):
        return False, "Menu_Budget.xlsx not found"

    wb = load_workbook(xlsx_path)
    sheet_names = [s.lower() for s in wb.sheetnames]

    if "menu items" not in sheet_names:
        return False, f"Missing 'Menu Items' sheet. Found: {wb.sheetnames}"

    if "summary" not in sheet_names:
        return False, f"Missing 'Summary' sheet. Found: {wb.sheetnames}"

    # Check Menu Items sheet has data
    menu_sheet = wb[wb.sheetnames[sheet_names.index("menu items")]]
    headers = [str(cell.value).lower() if cell.value else "" for cell in menu_sheet[1]]

    required_headers = ["course", "recipe_name", "ingredients_count", "difficulty"]
    for rh in required_headers:
        if not any(rh.replace("_", " ") in h or rh in h for h in headers):
            return False, f"Menu Items sheet missing header: {rh}. Found: {headers}"

    # Check at least 3 data rows (one per course)
    data_rows = sum(1 for row in menu_sheet.iter_rows(min_row=2) if row[0].value is not None)
    if data_rows < 3:
        return False, f"Menu Items sheet has {data_rows} data rows, expected at least 3"

    # Check Summary sheet has data
    summary_sheet = wb[wb.sheetnames[sheet_names.index("summary")]]
    summary_text = ""
    for row in summary_sheet.iter_rows():
        for cell in row:
            if cell.value:
                summary_text += str(cell.value).lower() + " "

    if "total_recipes" not in summary_text and "total recipes" not in summary_text:
        return False, "Summary sheet missing Total_Recipes metric"

    if "total_ingredients" not in summary_text and "total ingredients" not in summary_text:
        return False, "Summary sheet missing Total_Ingredients metric"

    return True, "Excel check passed"


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    errors = []

    # Check PPT
    print("\n=== Checking PowerPoint ===")
    pptx_ok, pptx_msg = check_pptx(args.agent_workspace)
    if pptx_ok:
        print(f"  [PASS] {pptx_msg}")
    else:
        print(f"  [FAIL] {pptx_msg}")
        errors.append(pptx_msg)

    # Check Excel
    print("\n=== Checking Excel ===")
    excel_ok, excel_msg = check_excel(args.agent_workspace)
    if excel_ok:
        print(f"  [PASS] {excel_msg}")
    else:
        print(f"  [FAIL] {excel_msg}")
        errors.append(excel_msg)

    # Summary
    print(f"\n=== SUMMARY ===")
    if errors:
        for e in errors:
            print(f"  [ERROR] {e}")
        print("  Overall: FAIL")
    else:
        print("  Overall: PASS")

    if args.res_log_file:
        result = {"errors": errors, "success": len(errors) == 0}
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(1 if errors else 0)


if __name__ == "__main__":
    main()
