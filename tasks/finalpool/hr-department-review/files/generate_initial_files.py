"""
Generate initial workspace files for hr-department-review task.

Creates:
- Department_Goals_2025.pptx (department goals and KPIs)
- HR_Policies.pdf (company HR policies)
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer


def create_pptx():
    """Create Department_Goals_2025.pptx."""
    prs = Presentation()

    departments = [
        {
            "name": "Engineering",
            "target_headcount": 7451,
            "budget_target": 95,
            "perf_target": 3.5,
            "initiative": "Adopt microservices architecture and improve CI/CD pipeline efficiency",
        },
        {
            "name": "Finance",
            "target_headcount": 7505,
            "budget_target": 95,
            "perf_target": 3.5,
            "initiative": "Implement automated financial reporting and reduce close cycle time",
        },
        {
            "name": "HR",
            "target_headcount": 7431,
            "budget_target": 95,
            "perf_target": 3.5,
            "initiative": "Launch employee wellness program and streamline onboarding process",
        },
        {
            "name": "Operations",
            "target_headcount": 7476,
            "budget_target": 95,
            "perf_target": 3.5,
            "initiative": "Optimize supply chain logistics and reduce operational overhead by 10%",
        },
        {
            "name": "R&D",
            "target_headcount": 7437,
            "budget_target": 95,
            "perf_target": 3.5,
            "initiative": "Accelerate product innovation cycle and file 20 new patents",
        },
        {
            "name": "Sales",
            "target_headcount": 7594,
            "budget_target": 95,
            "perf_target": 3.5,
            "initiative": "Expand into three new regional markets and grow revenue by 15%",
        },
        {
            "name": "Support",
            "target_headcount": 7606,
            "budget_target": 95,
            "perf_target": 3.5,
            "initiative": "Reduce average ticket resolution time by 20% and improve CSAT score",
        },
    ]

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "2025 Department Goals & KPIs"
    slide.placeholders[1].text = "Annual Performance Targets\nPrepared by the HR Director"

    # Slides 2-8: One per department
    for dept in departments:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"{dept['name']} Goals"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        lines = [
            f"Target headcount: {dept['target_headcount']}",
            f"Salary budget utilization target: {dept['budget_target']}%",
            f"Performance rating target: {dept['perf_target']}",
            f"Key initiative: {dept['initiative']}",
        ]
        for i, line in enumerate(lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line

    output_path = os.path.join(os.path.dirname(__file__), "Department_Goals_2025.pptx")
    prs.save(output_path)
    print(f"Created {output_path}")


def create_pdf():
    """Create HR_Policies.pdf."""
    output_path = os.path.join(os.path.dirname(__file__), "HR_Policies.pdf")
    doc = SimpleDocTemplate(output_path, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "CustomTitle", parent=styles["Title"], fontSize=20, spaceAfter=30
    )
    heading_style = ParagraphStyle(
        "CustomHeading", parent=styles["Heading2"], fontSize=14, spaceAfter=12,
        spaceBefore=20
    )
    body_style = ParagraphStyle(
        "CustomBody", parent=styles["Normal"], fontSize=11, spaceAfter=8,
        leading=16
    )

    elements = []

    # Title
    elements.append(Paragraph("Company HR Policies 2025", title_style))
    elements.append(Spacer(1, 20))

    # Section 1: Salary Ranges
    elements.append(Paragraph("Salary Ranges", heading_style))
    elements.append(Paragraph(
        "The company maintains the following salary band guidelines for all departments. "
        "Entry level positions fall within the range of $15,000 to $30,000 annually. "
        "Mid-level positions are compensated between $30,000 and $80,000 per year. "
        "Senior level positions command salaries of $80,000 and above, with no fixed "
        "upper cap but subject to departmental budget constraints and executive approval.",
        body_style
    ))
    elements.append(Spacer(1, 10))

    # Section 2: Performance Ratings
    elements.append(Paragraph("Performance Rating Criteria", heading_style))
    elements.append(Paragraph(
        "All employees are evaluated on a scale of 1 to 5. A rating of 1.0 indicates "
        "unsatisfactory performance requiring immediate improvement. A rating of 2.0 "
        "indicates performance that needs improvement. A rating of 3.0 means the "
        "employee meets expectations. A rating of 3.5 indicates the employee exceeds "
        "expectations. A rating of 4.0 or above is considered outstanding performance "
        "and qualifies the employee for accelerated promotion consideration and bonus "
        "eligibility.",
        body_style
    ))
    elements.append(Spacer(1, 10))

    # Section 3: Review Process
    elements.append(Paragraph("Review Process", heading_style))
    elements.append(Paragraph(
        "The company follows a structured review process. Quarterly performance reviews "
        "are conducted by each department manager for their direct reports. An annual "
        "company-wide performance review is prepared by the HR Director, which consolidates "
        "department-level data and presents an overall assessment to the executive team. "
        "Department managers are responsible for communicating review results to their teams "
        "and creating individual development plans based on the findings.",
        body_style
    ))

    doc.build(elements)
    print(f"Created {output_path}")


if __name__ == "__main__":
    create_pptx()
    create_pdf()
    print("Initial workspace files generated successfully.")
