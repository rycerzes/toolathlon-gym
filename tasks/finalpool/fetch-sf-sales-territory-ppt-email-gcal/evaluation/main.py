"""Evaluation for fetch-sf-sales-territory-ppt-email-gcal."""
import argparse
import os
import sys

import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432,
          dbname=os.environ.get("PGDATABASE", "toolathlon_gym"),
          user="eigent", password="camel")


def num_close(a, b, tol=0.5):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def check_pptx(agent_workspace):
    errors = []
    path = os.path.join(agent_workspace, "Territory_Scorecard.pptx")
    if not os.path.exists(path):
        return ["Territory_Scorecard.pptx not found"]
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = list(prs.slides)
        if len(slides) < 5:
            errors.append(f"Expected 5 slides, found {len(slides)}")

        # Gather all text from the presentation
        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + "\n"
        all_lower = all_text.lower()

        # Check title slide
        if "territory performance" not in all_lower and "q1 2026" not in all_lower:
            errors.append("Title slide missing expected title text")

        # Check revenue data for key regions
        if "asia pacific" not in all_lower:
            errors.append("Asia Pacific not mentioned in presentation")
        if "europe" not in all_lower:
            errors.append("Europe not mentioned in presentation")
        if "latin america" not in all_lower:
            errors.append("Latin America not mentioned in presentation")

        # Check quota attainment values
        if "103.7" not in all_text and "103.6" not in all_text and "103.8" not in all_text:
            errors.append("Asia Pacific attainment ~103.7% not found")
        if "94.7" not in all_text and "94.6" not in all_text and "94.8" not in all_text:
            errors.append("Latin America attainment ~94.7% not found")

        # Check segment mentions
        if "consumer" not in all_lower:
            errors.append("Consumer segment not mentioned")
        if "enterprise" not in all_lower:
            errors.append("Enterprise segment not mentioned")

        # Check pipeline coverage
        if "pipeline" not in all_lower and "coverage" not in all_lower:
            errors.append("Pipeline coverage not discussed")

        # Check recommendations slide
        if "recommendation" not in all_lower:
            errors.append("Recommendations slide not found")

    except Exception as e:
        errors.append(f"Error reading PPTX: {e}")
    return errors


def check_gcal():
    errors = []
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT summary, description FROM gcal.events
            WHERE start_datetime::date = '2026-03-28'
        """)
        rows = cur.fetchall()
        cur.close()
        conn.close()
        if not rows:
            errors.append("No GCal event found on 2026-03-28")
        else:
            summaries = [r[0].lower() if r[0] else "" for r in rows]
            if not any("territory" in s or "review" in s or "executive" in s for s in summaries):
                errors.append(f"No territory review event (found: {[r[0] for r in rows]})")
    except Exception as e:
        errors.append(f"Error checking GCal: {e}")
    return errors


def check_email():
    errors = []
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT subject, body_text FROM email.messages
            WHERE to_addr::text ILIKE '%%executive_team@company.com%%'
            ORDER BY id DESC LIMIT 5
        """)
        rows = cur.fetchall()
        cur.close()
        conn.close()
        if not rows:
            errors.append("No email found to executive_team@company.com")
        else:
            subjects = [r[0].lower() if r[0] else "" for r in rows]
            if not any("territory" in s or "q1" in s or "performance" in s for s in subjects):
                errors.append(f"Email subject doesn't match (found: {[r[0] for r in rows]})")
    except Exception as e:
        errors.append(f"Error checking email: {e}")
    return errors


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    agent_ws = args.agent_workspace or os.path.join(os.path.dirname(__file__), "..", "groundtruth_workspace")

    all_errors = []

    print("  Checking PowerPoint...")
    errs = check_pptx(agent_ws)
    if errs:
        all_errors.extend(errs)
        for e in errs[:5]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking GCal event...")
    errs = check_gcal()
    if errs:
        all_errors.extend(errs)
        for e in errs[:3]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking email...")
    errs = check_email()
    if errs:
        all_errors.extend(errs)
        for e in errs[:3]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
