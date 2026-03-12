"""
Evaluation script for notion-gform-gcal-onboarding task.

Checks:
1. PPT file exists with 4+ slides, title slide mentions "Welcome" and "2026"
2. Calendar events on 2026-03-16 (Orientation Session + Team Lunch)
3. Emails sent with "Welcome" and "Onboarding" in subject, at least 3 emails
4. Notion page updated (blocks with "March 2026" text)

Usage:
    python -m evaluation.main \
        --agent_workspace /path/to/workspace \
        --groundtruth_workspace /path/to/groundtruth \
        --launch_time "2026-03-06 10:00:00"
"""

import argparse
import json
import os
import sys
from datetime import timezone

import psycopg2
from pptx import Presentation

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def str_contains(haystack, needle):
    """Case-insensitive containment check."""
    if haystack is None or needle is None:
        return False
    return needle.lower() in str(haystack).lower()


# ============================================================================
# Check 1: PowerPoint file
# ============================================================================

def check_pptx(agent_workspace):
    """Verify Onboarding_Presentation.pptx exists and has correct content."""
    print("\n=== Checking PowerPoint ===")

    pptx_path = os.path.join(agent_workspace, "Onboarding_Presentation.pptx")

    if not os.path.isfile(pptx_path):
        record("PPT file exists", False, f"Not found: {pptx_path}")
        return False
    record("PPT file exists", True)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPT file readable", False, str(e))
        return False
    record("PPT file readable", True)

    slide_count = len(prs.slides)
    record("PPT has at least 4 slides", slide_count >= 4,
           f"Found {slide_count} slides")

    all_ok = True

    # Check title slide (first slide)
    if slide_count > 0:
        first_slide = prs.slides[0]
        slide_text = ""
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                slide_text += " " + shape.text_frame.text

        record("Title slide mentions 'Welcome'",
               str_contains(slide_text, "Welcome"),
               f"Slide text: {slide_text[:200]}")

        record("Title slide mentions '2026'",
               str_contains(slide_text, "2026"),
               f"Slide text: {slide_text[:200]}")
    else:
        record("Title slide exists", False, "No slides in presentation")
        all_ok = False

    # Check that at least one slide mentions new team members or hires
    all_slide_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_slide_text += " " + shape.text_frame.text

    all_slide_lower = all_slide_text.lower()

    # Check for checklist-related content (at least 3 checklist items mentioned)
    checklist_keywords = ["hr", "paperwork", "workstation", "accounts", "team",
                          "orientation", "handbook", "compliance", "training",
                          "office", "tour", "manager", "policies", "it accounts"]
    checklist_found = sum(1 for kw in checklist_keywords if kw in all_slide_lower)
    record("PPT contains checklist content (>=3 keywords)",
           checklist_found >= 3,
           f"Found {checklist_found} checklist keywords in slides")

    # Check for new team members info
    record("PPT mentions new team members",
           "team member" in all_slide_lower or "new hire" in all_slide_lower
           or "new team" in all_slide_lower or "member" in all_slide_lower,
           "No reference to team members found")

    return all_ok


# ============================================================================
# Check 2: Google Calendar events
# ============================================================================

def check_gcal():
    """Verify Orientation Session and Team Lunch events on 2026-03-16."""
    print("\n=== Checking Google Calendar ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT summary, description, start_datetime, end_datetime, attendees
        FROM gcal.events
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    cur.close()
    conn.close()

    print(f"[check_gcal] Found {len(events)} calendar events.")
    for ev in events:
        print(f"  Event: {ev[0]} | {ev[2]} - {ev[3]}")

    record("At least 2 calendar events created", len(events) >= 2,
           f"Found {len(events)}")

    all_ok = True

    # Check for Orientation Session
    orientation_found = False
    for summary, description, start_dt, end_dt, attendees in events:
        summary_lower = (summary or "").lower()
        if "orientation" in summary_lower:
            orientation_found = True

            # Check date is March 16, 2026
            if start_dt is not None:
                start_date_str = start_dt.strftime("%Y-%m-%d")
                record("Orientation on 2026-03-16",
                       start_date_str == "2026-03-16",
                       f"Start date: {start_date_str}")

                # Check time: 9:00 AM to 12:00 PM
                start_hour = start_dt.hour
                end_hour = end_dt.hour if end_dt else None
                record("Orientation 9 AM - 12 PM",
                       start_hour == 9 and end_hour == 12,
                       f"Start hour: {start_hour}, End hour: {end_hour}")
            break

    record("Orientation Session event exists", orientation_found,
           "No event with 'Orientation' in summary")

    # Check for Team Lunch
    lunch_found = False
    for summary, description, start_dt, end_dt, attendees in events:
        summary_lower = (summary or "").lower()
        if "lunch" in summary_lower:
            lunch_found = True

            if start_dt is not None:
                start_date_str = start_dt.strftime("%Y-%m-%d")
                record("Team Lunch on 2026-03-16",
                       start_date_str == "2026-03-16",
                       f"Start date: {start_date_str}")

                start_hour = start_dt.hour
                end_hour = end_dt.hour if end_dt else None
                record("Team Lunch 12 PM - 1 PM",
                       start_hour == 12 and end_hour == 13,
                       f"Start hour: {start_hour}, End hour: {end_hour}")
            break

    record("Team Lunch event exists", lunch_found,
           "No event with 'Lunch' in summary")

    return all_ok


# ============================================================================
# Check 3: Emails
# ============================================================================

def check_emails():
    """Verify 3 welcome emails sent to new hires."""
    print("\n=== Checking Emails ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
        WHERE folder_id = (SELECT id FROM email.folders WHERE name = 'Sent' LIMIT 1)
    """)
    sent_emails = cur.fetchall()
    cur.close()
    conn.close()

    print(f"[check_emails] Found {len(sent_emails)} sent emails.")

    record("At least 3 emails sent", len(sent_emails) >= 3,
           f"Found {len(sent_emails)}")

    all_ok = True

    # Check that emails have correct subject
    welcome_emails = []
    for subject, from_addr, to_addr, body_text in sent_emails:
        if str_contains(subject, "Welcome") and str_contains(subject, "Onboarding"):
            welcome_emails.append((subject, from_addr, to_addr, body_text))

    record("At least 3 welcome/onboarding emails",
           len(welcome_emails) >= 3,
           f"Found {len(welcome_emails)} emails with 'Welcome' and 'Onboarding' in subject")

    # Check that emails mention March 16, 2026 or orientation
    for i, (subject, from_addr, to_addr, body_text) in enumerate(welcome_emails[:3]):
        body_lower = (body_text or "").lower()

        record(f"Email {i+1}: body mentions start date or orientation",
               "march" in body_lower or "2026" in body_lower or "orientation" in body_lower,
               f"Body snippet: {(body_text or '')[:150]}")

    return all_ok


# ============================================================================
# Check 4: Notion page updated
# ============================================================================

def check_notion():
    """Verify Notion page was updated with March 2026 New Hires section."""
    print("\n=== Checking Notion ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Find the onboarding checklist page
    cur.execute("""
        SELECT id, properties
        FROM notion.pages
    """)
    pages = cur.fetchall()

    print(f"[check_notion] Found {len(pages)} Notion pages.")

    # Look for the onboarding checklist page
    checklist_page_id = None
    for page_id, props in pages:
        props_str = json.dumps(props).lower() if props else ""
        if "onboarding" in props_str and "checklist" in props_str:
            checklist_page_id = page_id
            break

    record("Onboarding checklist page exists", checklist_page_id is not None,
           f"No page with 'onboarding' and 'checklist' found among {len(pages)} pages")

    if not checklist_page_id:
        return False

    # Check blocks for "March 2026" content
    cur.execute("""
        SELECT block_data::text FROM notion.blocks
        WHERE parent_id = %s
    """, (checklist_page_id,))
    blocks = cur.fetchall()

    all_block_text = " ".join(str(b[0]).lower() for b in blocks)

    record("Notion blocks contain 'march 2026'",
           "march 2026" in all_block_text,
           f"'march 2026' not found in {len(blocks)} blocks")

    # Check for new hire names - look for any of the expected names
    # We check for names from the form responses injected by preprocess
    new_hire_names = ["sarah", "mike", "amy", "chen", "park", "rodriguez"]
    names_found = sum(1 for name in new_hire_names if name in all_block_text)
    record("Notion blocks mention new hire names (>=2)",
           names_found >= 2,
           f"Found {names_found} name references in blocks")

    cur.close()
    conn.close()

    return True


# ============================================================================
# Main
# ============================================================================

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    pptx_ok = check_pptx(args.agent_workspace)
    gcal_ok = check_gcal()
    email_ok = check_emails()
    notion_ok = check_notion()

    all_passed = pptx_ok and gcal_ok and email_ok and notion_ok

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Overall: {'PASS' if all_passed else 'FAIL'}")

    if args.res_log_file:
        result = {
            "passed": PASS_COUNT,
            "failed": FAIL_COUNT,
            "success": all_passed,
            "details": {
                "pptx": pptx_ok,
                "gcal": gcal_ok,
                "email": email_ok,
                "notion": notion_ok,
            },
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if all_passed else 1)


if __name__ == "__main__":
    main()
