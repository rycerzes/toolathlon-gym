"""
Evaluation script for sf-sales-product-ranking-ppt task.

Checks:
1. Excel file (Product_Rankings.xlsx) - Top Products and Category Summary sheets
2. PowerPoint (Product_Performance_Review.pptx) - >=4 slides with 'product' in text
3. Google Sheet with 'product' or 'ranking' in title
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

FILE_PASS = 0
FILE_FAIL = 0
DB_PASS = 0
DB_FAIL = 0


def check(name, condition, detail="", db=False):
    global FILE_PASS, FILE_FAIL, DB_PASS, DB_FAIL
    if condition:
        if db:
            DB_PASS += 1
        else:
            FILE_PASS += 1
        print(f"  [PASS] {name}")
    else:
        if db:
            DB_FAIL += 1
        else:
            FILE_FAIL += 1
        detail_str = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{detail_str}")


def get_expected_top20():
    """Get expected top 20 products from DB."""
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute('''
        SELECT p."PRODUCT_NAME", p."CATEGORY",
               SUM(o."TOTAL_AMOUNT") as revenue,
               SUM(o."QUANTITY") as units
        FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
        JOIN sf_data."SALES_DW__PUBLIC__PRODUCTS" p ON o."PRODUCT_ID" = p."PRODUCT_ID"
        GROUP BY p."PRODUCT_NAME", p."CATEGORY"
        ORDER BY revenue DESC
        LIMIT 20
    ''')
    rows = cur.fetchall()
    cur.close()
    conn.close()
    return rows


def check_excel(agent_workspace, expected_top20):
    print("\n=== Checking Excel Output ===")
    excel_path = os.path.join(agent_workspace, "Product_Rankings.xlsx")
    check("Excel file exists", os.path.isfile(excel_path), f"Expected {excel_path}")
    if not os.path.isfile(excel_path):
        return False

    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
    except Exception as e:
        check("Excel file readable", False, str(e))
        return False

    # Check Top Products sheet
    ws = None
    for s in wb.sheetnames:
        if "top" in s.lower() and "product" in s.lower():
            ws = wb[s]
            break
    if ws is None:
        for s in wb.sheetnames:
            if "top" in s.lower():
                ws = wb[s]
                break

    check("Sheet with 'Top Products' exists", ws is not None, f"Sheets: {wb.sheetnames}")
    if ws is None:
        return False

    rows = list(ws.iter_rows(min_row=2, values_only=True))
    check("Top Products has 20 rows", len(rows) == 20, f"Got {len(rows)}")

    if expected_top20 and len(rows) >= 1:
        # Check top product revenue is within 5% tolerance
        top_expected_rev = float(expected_top20[0][2])
        # Find revenue column (usually col index 3)
        for r in rows:
            if r and len(r) >= 4:
                try:
                    agent_rev = float(r[3])
                    rel_err = abs(agent_rev - top_expected_rev) / top_expected_rev if top_expected_rev else 1
                    check("Top product revenue within 5%", rel_err < 0.05,
                          f"Expected ~{top_expected_rev:.2f}, got {agent_rev:.2f}")
                    break
                except (TypeError, ValueError):
                    continue

    # Check Category Summary sheet
    ws2 = None
    for s in wb.sheetnames:
        if "category" in s.lower() or "summary" in s.lower():
            ws2 = wb[s]
            break
    check("Category Summary sheet exists", ws2 is not None, f"Sheets: {wb.sheetnames}")

    return True


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint Output ===")
    pptx_path = os.path.join(agent_workspace, "Product_Performance_Review.pptx")
    check("PPTX file exists", os.path.isfile(pptx_path), f"Expected {pptx_path}")
    if not os.path.isfile(pptx_path):
        return False

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        check("PPTX file readable", False, str(e))
        return False

    slide_count = len(prs.slides)
    check("PPTX has >= 4 slides", slide_count >= 4, f"Got {slide_count} slides")

    # Check for 'product' keyword in slide text
    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text.lower() + " "

    check("PPTX contains 'product' in text", "product" in all_text,
          f"Text sample: {all_text[:200]}")
    check("PPTX contains 'Q1 2026' or 'performance' or 'review'",
          "q1 2026" in all_text or "performance" in all_text or "review" in all_text,
          f"Text sample: {all_text[:200]}")

    return True


def check_gsheet():
    print("\n=== Checking Google Sheet ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("SELECT id, title FROM gsheet.spreadsheets")
    sheets = cur.fetchall()
    cur.close()
    conn.close()

    matching = [
        s for s in sheets
        if s[1] and ("product" in s[1].lower() or "ranking" in s[1].lower())
    ]
    check("Google Sheet with 'product' or 'ranking' in title exists",
          len(matching) > 0,
          f"Sheet titles: {[s[1] for s in sheets]}", db=True)
    return len(matching) > 0


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    expected_top20 = get_expected_top20()
    print(f"Expected top product: {expected_top20[0][0] if expected_top20 else 'N/A'}")

    check_excel(args.agent_workspace, expected_top20)
    check_pptx(args.agent_workspace)
    check_gsheet()

    total_pass = FILE_PASS + DB_PASS
    total_fail = FILE_FAIL + DB_FAIL
    file_ok = FILE_FAIL == 0

    print(f"\n=== SUMMARY ===")
    print(f"  File checks - Passed: {FILE_PASS}, Failed: {FILE_FAIL}")
    print(f"  DB checks   - Passed: {DB_PASS}, Failed: {DB_FAIL}")
    if DB_FAIL > 0:
        print(f"  WARNING: {DB_FAIL} DB checks failed (not blocking)")
    print(f"  Overall: {'PASS' if file_ok else 'FAIL'}")

    if args.res_log_file:
        result = {"passed": total_pass, "failed": total_fail, "success": file_ok}
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if file_ok else 1)


if __name__ == "__main__":
    main()
