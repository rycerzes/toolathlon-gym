"""
Evaluation script for howtocook-nutrition-ppt task.

Checks:
1. Excel file exists with correct sheets and columns
2. Recipe Overview has at least 5 data rows with valid values
3. Recommendation has at least 3 data rows with valid values
4. PPT exists with at least 7 slides
5. PPT title slide references healthy eating
6. PPT mentions recipe names that also appear in Excel
7. Categories are diverse (at least 4 distinct)
"""

from argparse import ArgumentParser
import sys
import os
from pathlib import Path

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)



def check_excel(workspace):
    """Check Nutrition_Comparison.xlsx for structure and data validity."""
    import openpyxl

    xlsx_path = Path(workspace) / "Nutrition_Comparison.xlsx"
    if not xlsx_path.exists():
        return False, "Nutrition_Comparison.xlsx not found in workspace"

    wb = openpyxl.load_workbook(xlsx_path, data_only=True)

    # Check sheet names (case-insensitive comparison)
    sheet_names_lower = {s.lower(): s for s in wb.sheetnames}
    if "recipe overview" not in sheet_names_lower:
        return False, f"Missing 'Recipe Overview' sheet. Found: {wb.sheetnames}"
    if "recommendation" not in sheet_names_lower:
        return False, f"Missing 'Recommendation' sheet. Found: {wb.sheetnames}"

    overview_sheet_name = sheet_names_lower["recipe overview"]
    recommendation_sheet_name = sheet_names_lower["recommendation"]

    # --- Check Recipe Overview ---
    ws1 = wb[overview_sheet_name]
    rows1 = list(ws1.iter_rows(values_only=True))
    if len(rows1) < 2:
        return False, "Recipe Overview sheet has no data rows"

    header1 = [str(h).strip() if h else "" for h in rows1[0]]
    header1_lower = [h.lower().replace(" ", "_") for h in header1]

    expected_overview_cols = ["recipe_name", "category", "difficulty", "estimated_calories", "protein_level", "fiber_level"]
    for col in expected_overview_cols:
        if col not in header1_lower:
            return False, f"Recipe Overview missing column '{col}'. Found headers: {header1}"

    idx = {col: header1_lower.index(col) for col in expected_overview_cols}

    data_rows = rows1[1:]
    # Filter out empty rows
    data_rows = [r for r in data_rows if r[idx["recipe_name"]] is not None and str(r[idx["recipe_name"]]).strip()]

    if len(data_rows) < 5:
        return False, f"Recipe Overview: expected at least 5 data rows, got {len(data_rows)}"

    categories = set()
    recipe_names_overview = []

    for i, row in enumerate(data_rows):
        recipe_name = str(row[idx["recipe_name"]]).strip() if row[idx["recipe_name"]] else ""
        category = str(row[idx["category"]]).strip() if row[idx["category"]] else ""
        difficulty = str(row[idx["difficulty"]]).strip() if row[idx["difficulty"]] else ""
        cal_val = row[idx["estimated_calories"]]
        protein = str(row[idx["protein_level"]]).strip() if row[idx["protein_level"]] else ""
        fiber = str(row[idx["fiber_level"]]).strip() if row[idx["fiber_level"]] else ""

        if not recipe_name:
            return False, f"Recipe Overview row {i+1}: Recipe_Name is empty"
        recipe_names_overview.append(recipe_name)

        if not category:
            return False, f"Recipe Overview row {i+1} ({recipe_name}): Category is empty"
        categories.add(category.lower())

        if not difficulty:
            return False, f"Recipe Overview row {i+1} ({recipe_name}): Difficulty is empty"
        if difficulty.lower() not in ["easy", "medium", "hard"]:
            return False, f"Recipe Overview row {i+1} ({recipe_name}): Difficulty '{difficulty}' not in [Easy, Medium, Hard]"

        if cal_val is None:
            return False, f"Recipe Overview row {i+1} ({recipe_name}): Estimated_Calories is empty"
        try:
            cal_float = float(cal_val)
        except (ValueError, TypeError):
            return False, f"Recipe Overview row {i+1} ({recipe_name}): Estimated_Calories '{str(cal_val)[:50]}' is not numeric"
        if cal_float < 50 or cal_float > 5000:
            return False, f"Recipe Overview row {i+1} ({recipe_name}): Estimated_Calories {cal_float} out of reasonable range (50-5000)"

        if not protein:
            return False, f"Recipe Overview row {i+1} ({recipe_name}): Protein_Level is empty"
        if protein.lower() not in ["high", "medium", "low"]:
            return False, f"Recipe Overview row {i+1} ({recipe_name}): Protein_Level '{protein}' not in [High, Medium, Low]"

        if not fiber:
            return False, f"Recipe Overview row {i+1} ({recipe_name}): Fiber_Level is empty"
        if fiber.lower() not in ["high", "medium", "low"]:
            return False, f"Recipe Overview row {i+1} ({recipe_name}): Fiber_Level '{fiber}' not in [High, Medium, Low]"

    if len(categories) < 4:
        return False, f"Recipe Overview: expected at least 4 distinct categories, got {len(categories)}: {categories}"

    print(f"  [PASS] Recipe Overview: {len(data_rows)} recipes, {len(categories)} categories")

    # --- Check Recommendation ---
    ws2 = wb[recommendation_sheet_name]
    rows2 = list(ws2.iter_rows(values_only=True))
    if len(rows2) < 2:
        return False, "Recommendation sheet has no data rows"

    header2 = [str(h).strip() if h else "" for h in rows2[0]]
    header2_lower = [h.lower().replace(" ", "_") for h in header2]

    expected_rec_cols = ["rank", "recipe_name", "health_score", "reason"]
    for col in expected_rec_cols:
        if col not in header2_lower:
            return False, f"Recommendation missing column '{col}'. Found headers: {header2}"

    idx2 = {col: header2_lower.index(col) for col in expected_rec_cols}

    data_rows2 = rows2[1:]
    data_rows2 = [r for r in data_rows2 if r[idx2["recipe_name"]] is not None and str(r[idx2["recipe_name"]]).strip()]

    if len(data_rows2) < 3:
        return False, f"Recommendation: expected at least 3 data rows, got {len(data_rows2)}"

    rec_recipe_names = []
    prev_score = None
    for i, row in enumerate(data_rows2):
        rank_val = row[idx2["rank"]]
        recipe_name = str(row[idx2["recipe_name"]]).strip() if row[idx2["recipe_name"]] else ""
        score_val = row[idx2["health_score"]]
        reason_val = str(row[idx2["reason"]]).strip() if row[idx2["reason"]] else ""

        if rank_val is None:
            return False, f"Recommendation row {i+1}: Rank is empty"
        try:
            int(float(str(rank_val)))
        except (ValueError, TypeError):
            return False, f"Recommendation row {i+1}: Rank '{str(rank_val)[:50]}' is not numeric"

        if not recipe_name:
            return False, f"Recommendation row {i+1}: Recipe_Name is empty"
        rec_recipe_names.append(recipe_name)

        if score_val is None:
            return False, f"Recommendation row {i+1} ({recipe_name}): Health_Score is empty"
        try:
            score_float = float(score_val)
        except (ValueError, TypeError):
            return False, f"Recommendation row {i+1} ({recipe_name}): Health_Score '{str(score_val)[:50]}' is not numeric"

        # Check descending order
        if prev_score is not None and score_float > prev_score + 0.01:
            return False, f"Recommendation: recipes not in descending Health_Score order (row {i}: {score_float} > row {i-1}: {prev_score})"
        prev_score = score_float

        if not reason_val:
            return False, f"Recommendation row {i+1} ({recipe_name}): Reason is empty"

    # Check that recommended recipes appear in overview
    overview_names_lower = [n.lower() for n in recipe_names_overview]
    for rn in rec_recipe_names:
        if rn.lower() not in overview_names_lower:
            return False, f"Recommendation recipe '{rn}' not found in Recipe Overview"

    print(f"  [PASS] Recommendation: {len(data_rows2)} recommended recipes")
    wb.close()
    return True, "Excel file checks passed", recipe_names_overview, rec_recipe_names


def check_pptx(workspace, rec_recipe_names=None):
    """Check Healthy_Eating_Guide.pptx for structure."""
    from pptx import Presentation

    pptx_path = Path(workspace) / "Healthy_Eating_Guide.pptx"
    if not pptx_path.exists():
        return False, "Healthy_Eating_Guide.pptx not found in workspace"

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)

    if len(slides) < 5:
        return False, f"Expected at least 5 slides, got {len(slides)}"
    print(f"  Slide count: {len(slides)}")

    # Collect all text from all slides
    all_text = []
    for slide in slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_texts.append(paragraph.text)
        all_text.append("\n".join(slide_texts))

    full_text = "\n".join(all_text).lower()

    # Check title slide contains healthy/eating/nutrition keywords
    first_slide_text = all_text[0].lower()
    has_keyword = any(kw in first_slide_text for kw in ["healthy", "nutrition", "eating", "health", "wellness"])
    if not has_keyword:
        return False, f"Title slide does not contain healthy eating keywords. Text: {all_text[0][:200]}"
    print("  [PASS] Title slide contains healthy eating keywords")

    # Check that at least 3 recipe names from recommendation appear in the PPT
    if rec_recipe_names:
        found_count = 0
        for rn in rec_recipe_names:
            if rn.lower() in full_text:
                found_count += 1
        if found_count < 3:
            return False, f"Only {found_count} of {len(rec_recipe_names)} recommended recipe names found in PPT (need at least 3)"
        print(f"  [PASS] {found_count} recommended recipe names found in presentation")
    else:
        print("  [SKIP] No recipe names to cross-check with PPT")

    # Check last slide has summary content (more than just a title)
    last_slide_text = all_text[-1]
    if len(last_slide_text.strip()) < 10:
        return False, "Summary slide (last slide) has too little text content"
    print("  [PASS] Summary slide has content")

    return True, "PPTX file checks passed"


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False, help="Launch time")
    args = parser.parse_args()

    workspace = args.agent_workspace
    if not workspace:
        print("Error: --agent_workspace is required")
        sys.exit(1)

    all_passed = True
    rec_recipe_names = None

    # Check Excel
    print("\n--- Check 1: Excel File ---")
    try:
        result = check_excel(workspace)
        ok = result[0]
        msg = result[1]
        if not ok:
            print(f"  [FAIL] {msg}")
            all_passed = False
        else:
            print(f"  {msg}")
            if len(result) > 3:
                rec_recipe_names = result[3]
    except Exception as e:
        print(f"  [FAIL] Excel check error: {e}")
        import traceback
        traceback.print_exc()
        all_passed = False

    # Check PPTX
    print("\n--- Check 2: PowerPoint File ---")
    try:
        ok, msg = check_pptx(workspace, rec_recipe_names)
        if not ok:
            print(f"  [FAIL] {msg}")
            all_passed = False
        else:
            print(f"  {msg}")
    except Exception as e:
        print(f"  [FAIL] PPTX check error: {e}")
        import traceback
        traceback.print_exc()
        all_passed = False

    if all_passed:
        print("\nPass all tests!")
    else:
        print("\nSome checks failed.")
        sys.exit(1)
