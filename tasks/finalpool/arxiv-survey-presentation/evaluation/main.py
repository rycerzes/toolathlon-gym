"""
Evaluation for arxiv-survey-presentation task.
Checks Excel (Survey_Data.xlsx) and PowerPoint (NLP_Survey.pptx).
"""
import argparse
import os
import sys
from pathlib import Path

import openpyxl

PASS_COUNT = 0
FAIL_COUNT = 0

EXPECTED_PAPER_IDS = {"2404.00001", "2404.00002", "2404.00003", "2404.00004", "2404.00005"}


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower().replace(" ", "_") == sheet_name.strip().lower().replace(" ", "_"):
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
        if name.strip().lower().replace("_", " ") == sheet_name.strip().lower().replace("_", " "):
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def find_col(header, names):
    if not header:
        return None
    for i, cell in enumerate(header):
        if cell is None:
            continue
        cl = str(cell).strip().lower().replace(" ", "_")
        for n in names:
            if n.lower().replace(" ", "_") == cl:
                return i
    return None


def check_excel(workspace):
    print("\n=== Checking Excel ===")
    path = os.path.join(workspace, "Survey_Data.xlsx")
    if not os.path.isfile(path):
        record("Excel exists", False, f"Not found: {path}")
        return False
    record("Excel exists", True)

    wb = openpyxl.load_workbook(path, data_only=True)

    # Paper Comparison
    pc_rows = load_sheet_rows(wb, "Paper Comparison") or load_sheet_rows(wb, "Paper_Comparison")
    if pc_rows is None:
        record("Sheet 'Paper Comparison' exists", False, f"Sheets: {wb.sheetnames}")
    else:
        record("Sheet 'Paper Comparison' exists", True)
        data = pc_rows[1:]
        record("Paper Comparison has 5 rows", len(data) == 5, f"Found {len(data)}")

        id_col = find_col(pc_rows[0], ["Paper_ID", "Paper ID", "ID"])
        if id_col is not None:
            found = {str(r[id_col]).strip() for r in data if id_col < len(r) and r[id_col]}
            for pid in EXPECTED_PAPER_IDS:
                record(f"Paper {pid} present", pid in found, f"Found: {found}")

        method_col = find_col(pc_rows[0], ["Method", "method"])
        record("Method column exists", method_col is not None, f"Header: {pc_rows[0]}")

        dataset_col = find_col(pc_rows[0], ["Dataset", "dataset", "Dataset_Used"])
        record("Dataset column exists", dataset_col is not None, f"Header: {pc_rows[0]}")

    # Taxonomy
    tax_rows = load_sheet_rows(wb, "Taxonomy")
    if tax_rows is None:
        record("Sheet 'Taxonomy' exists", False, f"Sheets: {wb.sheetnames}")
    else:
        record("Sheet 'Taxonomy' exists", True)
        data = tax_rows[1:]
        record("Taxonomy has >= 2 categories", len(data) >= 2, f"Found {len(data)}")

    # Summary Statistics
    ss_rows = load_sheet_rows(wb, "Summary Statistics") or load_sheet_rows(wb, "Summary_Statistics")
    if ss_rows is None:
        record("Sheet 'Summary Statistics' exists", False, f"Sheets: {wb.sheetnames}")
    else:
        record("Sheet 'Summary Statistics' exists", True)
        metrics = {}
        for row in ss_rows[1:]:
            if row and row[0]:
                metrics[str(row[0]).strip().lower().replace(" ", "_")] = row[1] if len(row) > 1 else None

        tp_key = next((k for k in metrics if "total" in k and "paper" in k), None)
        if tp_key:
            try:
                record("Total_Papers = 5", abs(float(metrics[tp_key]) - 5) < 1, f"Got {metrics[tp_key]}")
            except (TypeError, ValueError):
                record("Total_Papers is numeric", False, f"Got {metrics[tp_key]}")

    return True


def check_pptx(workspace):
    print("\n=== Checking PowerPoint ===")
    path = os.path.join(workspace, "NLP_Survey.pptx")
    if not os.path.isfile(path):
        record("PPTX exists", False, f"Not found: {path}")
        return False
    record("PPTX exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = list(prs.slides)

        record("Has >= 6 slides", len(slides) >= 6, f"Found {len(slides)}")

        all_text = []
        for slide in slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        texts.append(p.text)
            all_text.append("\n".join(texts))

        full = "\n".join(all_text).lower()

        # Title slide
        first = all_text[0].lower() if all_text else ""
        has_title = any(kw in first for kw in ["survey", "nlp", "language model", "reasoning", "advances"])
        record("Title slide has survey keywords", has_title, f"First slide: {first[:200]}")

        # Content checks
        has_papers = any(kw in full for kw in ["chain-of-thought", "instruction tuning", "retrieval", "self-consistency", "multimodal"])
        record("Mentions paper methods", has_papers)

        # Summary/conclusion slide
        last = all_text[-1] if all_text else ""
        record("Last slide has content", len(last.strip()) > 10, f"Last slide: {last[:100]}")

        return True
    except Exception as e:
        record("PPTX readable", False, str(e))
        return False


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}, Failed: {FAIL_COUNT}")
    sys.exit(0 if FAIL_COUNT == 0 else 1)


if __name__ == "__main__":
    main()
