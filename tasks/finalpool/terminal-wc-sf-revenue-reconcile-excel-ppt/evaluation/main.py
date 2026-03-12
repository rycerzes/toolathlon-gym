"""Evaluation for terminal-wc-sf-revenue-reconcile-excel-ppt."""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432,
          dbname=os.environ.get("PGDATABASE", "toolathlon_gym"),
          user="eigent", password="camel")

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:300]}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def get_sheet(wb, name):
    for s in wb.sheetnames:
        if s.strip().lower().replace(" ", "_") == name.strip().lower().replace(" ", "_"):
            return wb[s]
    return None


def check_excel(agent_ws, gt_dir):
    print("\n=== Checking Revenue_Reconciliation.xlsx ===")
    agent_file = os.path.join(agent_ws, "Revenue_Reconciliation.xlsx")
    gt_file = os.path.join(gt_dir, "Revenue_Reconciliation.xlsx")

    check("Excel file exists", os.path.isfile(agent_file), agent_file)
    if not os.path.isfile(agent_file):
        return

    try:
        awb = openpyxl.load_workbook(agent_file, data_only=True)
        gwb = openpyxl.load_workbook(gt_file, data_only=True)
    except Exception as e:
        check("Excel readable", False, str(e))
        return

    # Sheet 1: WC_Summary
    print("  Checking WC_Summary...")
    ws1 = get_sheet(awb, "WC_Summary")
    check("Sheet WC_Summary exists", ws1 is not None, f"Sheets: {awb.sheetnames}")
    if ws1:
        rows = list(ws1.iter_rows(min_row=2, values_only=True))
        # Should have 7 status rows + 1 total
        data_rows = [r for r in rows if r and r[0] and "total" not in str(r[0]).lower()]
        check("WC_Summary has 7 status rows", len(data_rows) == 7, f"Got {len(data_rows)}")

        # Query dynamic WC values from DB
        try:
            conn = psycopg2.connect(**DB)
            cur = conn.cursor()
            cur.execute("SELECT COUNT(*), COALESCE(SUM(total::numeric), 0) FROM wc.orders WHERE LOWER(status) = 'completed'")
            expected_completed_count, expected_completed_rev = cur.fetchone()
            cur.execute("""
                SELECT COUNT(*), COALESCE(SUM(total::numeric), 0) FROM wc.orders
                WHERE LOWER(status) IN ('completed', 'processing')
            """)
            expected_total_count, expected_total_rev = cur.fetchone()
            cur.close(); conn.close()
        except Exception:
            expected_completed_count, expected_completed_rev = 72, 30296.82
            expected_total_count, expected_total_rev = 100, 39418.17

        # Check completed row
        completed = [r for r in data_rows if r[0] and "completed" in str(r[0]).lower()]
        if completed:
            check(f"WC completed order_count = {expected_completed_count}",
                  num_close(completed[0][1], expected_completed_count, 1),
                  f"Got {completed[0][1]}")
            check(f"WC completed revenue ~ {expected_completed_rev:.2f}",
                  num_close(completed[0][2], expected_completed_rev, 50),
                  f"Got {completed[0][2]}")

        # Check total row
        total_rows = [r for r in rows if r and r[0] and "total" in str(r[0]).lower()]
        if total_rows:
            check(f"WC confirmed total orders = {expected_total_count}",
                  num_close(total_rows[0][1], expected_total_count, 2),
                  f"Got {total_rows[0][1]}")
            check(f"WC confirmed total revenue ~ {expected_total_rev:.0f}",
                  num_close(total_rows[0][2], expected_total_rev, 100),
                  f"Got {total_rows[0][2]}")

    # Sheet 2: SF_Summary
    print("  Checking SF_Summary...")
    ws2 = get_sheet(awb, "SF_Summary")
    check("Sheet SF_Summary exists", ws2 is not None, f"Sheets: {awb.sheetnames}")
    if ws2:
        all_rows = list(ws2.iter_rows(min_row=2, values_only=True))
        # Find region data (before blank row)
        region_rows = []
        for r in all_rows:
            if r and r[0] and str(r[0]).strip():
                name = str(r[0]).strip().lower()
                if name in ["europe", "asia pacific", "north america", "middle east", "latin america"]:
                    region_rows.append(r)
        check("SF_Summary has 5 regions", len(region_rows) == 5, f"Got {len(region_rows)}")

        # Query dynamic SF Europe values from DB
        try:
            conn = psycopg2.connect(**DB)
            cur = conn.cursor()
            cur.execute("SELECT COUNT(*), COALESCE(SUM(sales), 0) FROM sf_data.orders WHERE LOWER(region) = 'europe'")
            expected_eu_count, expected_eu_rev = cur.fetchone()
            cur.close(); conn.close()
        except Exception:
            expected_eu_count, expected_eu_rev = 4100, 648798.47

        europe = [r for r in region_rows if "europe" in str(r[0]).lower()]
        if europe:
            check(f"Europe order_count = {expected_eu_count}",
                  num_close(europe[0][1], expected_eu_count, 5),
                  f"Got {europe[0][1]}")
            check(f"Europe revenue ~ {expected_eu_rev:.0f}",
                  num_close(europe[0][2], expected_eu_rev, 500),
                  f"Got {europe[0][2]}")

    # Sheet 3: Cross_Audit
    print("  Checking Cross_Audit...")
    ws3 = get_sheet(awb, "Cross_Audit")
    check("Sheet Cross_Audit exists", ws3 is not None, f"Sheets: {awb.sheetnames}")
    if ws3:
        rows3 = list(ws3.iter_rows(min_row=2, values_only=True))
        check("Cross_Audit has 3 metric rows", len(rows3) == 3, f"Got {len(rows3)}")

        a_lookup = {}
        for r in rows3:
            if r and len(r) >= 6 and r[1]:
                a_lookup[str(r[1]).strip().lower()] = r

        # Query dynamic cross-audit values from DB
        try:
            conn = psycopg2.connect(**DB)
            cur = conn.cursor()
            cur.execute("""
                SELECT COUNT(*), COALESCE(SUM(total::numeric), 0)
                FROM wc.orders WHERE LOWER(status) IN ('completed', 'processing')
            """)
            wc_total_orders, wc_total_rev = cur.fetchone()
            wc_aov = float(wc_total_rev) / wc_total_orders if wc_total_orders else 0
            cur.execute("SELECT COUNT(*), COALESCE(SUM(sales), 0) FROM sf_data.orders")
            sf_total_orders, sf_total_rev = cur.fetchone()
            sf_aov = float(sf_total_rev) / sf_total_orders if sf_total_orders else 0
            cur.close(); conn.close()
        except Exception:
            wc_total_orders, sf_total_orders = 100, 20000
            wc_aov, sf_aov = 394.18, 152.45

        # Check Total_Orders variance
        to = a_lookup.get("total_orders")
        if to:
            check(f"Total_Orders WC={wc_total_orders}", num_close(to[2], wc_total_orders, 5), f"Got {to[2]}")
            check(f"Total_Orders SF={sf_total_orders}", num_close(to[3], sf_total_orders, 50), f"Got {to[3]}")
            check("Total_Orders Flag=REVIEW",
                  to[5] is not None and "review" in str(to[5]).lower(),
                  f"Got {to[5]}")

        # Check Avg_Order_Value
        aov = a_lookup.get("avg_order_value")
        if aov:
            check(f"AOV WC ~ {wc_aov:.0f}", num_close(aov[2], wc_aov, 20), f"Got {aov[2]}")
            check(f"AOV SF ~ {sf_aov:.0f}", num_close(aov[3], sf_aov, 10), f"Got {aov[3]}")

    # Sheet 4: Recommendations
    print("  Checking Recommendations...")
    ws4 = get_sheet(awb, "Recommendations")
    check("Sheet Recommendations exists", ws4 is not None, f"Sheets: {awb.sheetnames}")
    if ws4:
        rows4 = list(ws4.iter_rows(min_row=2, values_only=True))
        data_rows4 = [r for r in rows4 if r and r[0]]
        check("Recommendations has 5 findings", len(data_rows4) == 5, f"Got {len(data_rows4)}")


def check_pptx(agent_ws):
    print("\n=== Checking Reconciliation_Presentation.pptx ===")
    pptx_path = os.path.join(agent_ws, "Reconciliation_Presentation.pptx")
    check("PPTX file exists", os.path.isfile(pptx_path), pptx_path)
    if not os.path.isfile(pptx_path):
        return
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        check("PPTX has >= 5 slides", len(prs.slides) >= 5, f"Got {len(prs.slides)}")

        # Check title slide
        if len(prs.slides) > 0:
            title_text = ""
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame:
                    title_text += shape.text_frame.text.lower()
            check("Title slide mentions reconciliation",
                  "reconcil" in title_text,
                  f"Title: {title_text[:100]}")

        # Check content slides have data
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text.lower() + " "

        check("PPTX mentions variance", "variance" in all_text or "%" in all_text,
              "No variance/percentage content found")
        check("PPTX mentions revenue", "revenue" in all_text,
              "No revenue content found")
    except ImportError:
        check("python-pptx available", False, "Cannot verify PPTX content")
    except Exception as e:
        check("PPTX readable", False, str(e))


def check_reverse_validation(workspace):
    print("\n=== Reverse Validation ===")
    # Check no unexpected sheets in the Excel file
    excel_path = os.path.join(workspace, "Revenue_Reconciliation.xlsx")
    if os.path.isfile(excel_path):
        try:
            wb = openpyxl.load_workbook(excel_path, data_only=True)
            expected_sheets = {"wc_summary", "sf_summary", "cross_audit", "recommendations"}
            actual_sheets = {s.strip().lower().replace(" ", "_") for s in wb.sheetnames}
            unexpected = actual_sheets - expected_sheets
            check("No unexpected sheets in Excel",
                  len(unexpected) == 0,
                  f"Unexpected sheets: {unexpected}")

            # Check no negative revenue values in WC_Summary
            ws = get_sheet(wb, "WC_Summary")
            if ws:
                for row in ws.iter_rows(min_row=2, values_only=True):
                    if row and row[2] is not None:
                        try:
                            val = float(row[2])
                            if val < 0:
                                check("No negative revenue in WC_Summary", False,
                                      f"Found negative revenue: {val}")
                                break
                        except (TypeError, ValueError):
                            pass
                else:
                    check("No negative revenue in WC_Summary", True)
            wb.close()
        except Exception as e:
            check("Reverse validation readable", False, str(e))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    check_excel(args.agent_workspace, gt_dir)
    check_pptx(args.agent_workspace)
    check_reverse_validation(args.agent_workspace)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    overall = FAIL_COUNT == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
