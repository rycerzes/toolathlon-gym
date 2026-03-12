"""
Evaluation for train-yf-roadshow-excel-ppt-email task.

Checks:
1. Roadshow_Prep.xlsx with Travel_Details and Stock_Summary sheets
2. G1 business class travel details correct
3. AMZN stock data correct (sector, 5-day stats)
4. Roadshow_Agenda.pptx with 4 slides and correct titles
5. Email sent to roadshow@bank.com
"""
import json
import os
import sys
from argparse import ArgumentParser

import psycopg2
import openpyxl

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def check_excel(agent_workspace):
    print("\n=== Check 1: Roadshow_Prep.xlsx ===")

    xlsx_path = os.path.join(agent_workspace, "Roadshow_Prep.xlsx")
    if not os.path.exists(xlsx_path):
        record("Roadshow_Prep.xlsx exists", False, f"Not found at {xlsx_path}")
        return
    record("Roadshow_Prep.xlsx exists", True)

    try:
        wb = openpyxl.load_workbook(xlsx_path)
    except Exception as e:
        record("Excel readable", False, str(e))
        return
    record("Excel readable", True)

    sheet_names_lower = [s.lower() for s in wb.sheetnames]

    # Travel_Details sheet
    if "travel_details" not in sheet_names_lower:
        record("Travel_Details sheet exists", False, f"Sheets: {wb.sheetnames}")
    else:
        record("Travel_Details sheet exists", True)
        ws = wb[wb.sheetnames[sheet_names_lower.index("travel_details")]]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if any(c for c in r)]
        record("Travel_Details has 1 data row", len(data_rows) >= 1,
               f"Found {len(data_rows)}")

        all_text = " ".join(str(c) for r in rows for c in r if c).upper()
        record("Travel_Details has G1", "G1" in all_text, all_text[:200])
        record("Travel_Details has Business Class", "BUSINESS" in all_text, all_text[:200])

        # Check price 1748.5
        numeric_vals = []
        for r in data_rows:
            for c in r:
                try:
                    numeric_vals.append(float(c))
                except (TypeError, ValueError):
                    pass
        has_price = any(abs(v - 1748.5) < 1.0 for v in numeric_vals)
        record("Business class price ~1748.5 CNY", has_price, f"Numerics: {numeric_vals}")

    # Stock_Summary sheet
    if "stock_summary" not in sheet_names_lower:
        record("Stock_Summary sheet exists", False, f"Sheets: {wb.sheetnames}")
    else:
        record("Stock_Summary sheet exists", True)
        ws2 = wb[wb.sheetnames[sheet_names_lower.index("stock_summary")]]
        rows2 = list(ws2.iter_rows(values_only=True))
        all_text2 = " ".join(str(c) for r in rows2 for c in r if c).upper()
        record("Stock_Summary has AMZN", "AMZN" in all_text2, all_text2[:200])
        record("Stock_Summary has sector info", "CONSUMER" in all_text2 or "CYCLICAL" in all_text2, all_text2[:200])

        numeric_vals2 = []
        for r in rows2:
            for c in r:
                try:
                    numeric_vals2.append(float(c))
                except (TypeError, ValueError):
                    pass
        # Check that 5-day avg ~212.58 is present
        has_avg = any(abs(v - 212.58) < 2.0 for v in numeric_vals2)
        record("Stock_Summary 5-day avg ~212.58", has_avg, f"Numerics: {numeric_vals2}")

        # Check latest close ~218.94
        has_close = any(abs(v - 218.94) < 2.0 for v in numeric_vals2)
        record("Stock_Summary latest close ~218.94", has_close, f"Numerics: {numeric_vals2}")


def check_ppt(agent_workspace):
    print("\n=== Check 2: Roadshow_Agenda.pptx ===")

    pptx_path = os.path.join(agent_workspace, "Roadshow_Agenda.pptx")
    if not os.path.exists(pptx_path):
        record("Roadshow_Agenda.pptx exists", False, f"Not found at {pptx_path}")
        return
    record("Roadshow_Agenda.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPT file readable", False, str(e))
        return
    record("PPT file readable", True)

    record("PPT has exactly 4 slides", len(prs.slides) == 4, f"Found {len(prs.slides)}")

    all_titles = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_type == 13 or hasattr(shape, "placeholder_format"):
                pass
        title_shape = slide.shapes.title
        if title_shape:
            all_titles.append(title_shape.text)

    all_titles_lower = " ".join(all_titles).lower()
    record("Slide 1 title contains 'Investor Roadshow'",
           "investor roadshow" in all_titles_lower or "roadshow" in all_titles_lower,
           f"Titles: {all_titles}")
    record("Slide 2 title contains 'Journey Details'",
           "journey" in all_titles_lower, f"Titles: {all_titles}")
    record("Slide 3 title contains 'Portfolio Overview'",
           "portfolio" in all_titles_lower, f"Titles: {all_titles}")
    record("Slide 4 title contains 'Meeting Schedule'",
           "meeting" in all_titles_lower or "schedule" in all_titles_lower, f"Titles: {all_titles}")

    # Check all text for key content
    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_text += para.text + " "
    all_text_lower = all_text.lower()
    record("PPT mentions G1 train", "g1" in all_text_lower, all_text_lower[:200])
    record("PPT mentions AMZN", "amzn" in all_text_lower, all_text_lower[:200])
    record("PPT slide 4 says TBD", "tbd" in all_text_lower, all_text_lower[-200:])


def check_email():
    print("\n=== Check 3: Email to roadshow@bank.com ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT to_addr, subject FROM email.messages
        WHERE subject ILIKE '%roadshow%' OR subject ILIKE '%travel%prep%'
    """)
    messages = cur.fetchall()
    cur.close()
    conn.close()

    all_msgs = list(messages)
    record("Roadshow email sent", len(all_msgs) >= 1,
           f"Found {len(all_msgs)} matching emails")

    if all_msgs:
        to_raw = all_msgs[0][0]
        to_str = str(to_raw).lower() if to_raw else ""
        record("Email sent to roadshow@bank.com", "roadshow@bank.com" in to_str,
               f"To: {to_str[:100]}")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_ppt(args.agent_workspace)
    check_email()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
