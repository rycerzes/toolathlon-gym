"""Evaluation for canvas-faculty-workload-review."""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=2.0):
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def get_expected_instructors():
    """Get expected instructor data from DB."""
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT u.name as instructor,
               COUNT(DISTINCT c.id) as courses_count,
               SUM(sc.cnt) as total_students,
               SUM(ac.cnt) as total_assignments,
               ROUND(SUM(sc.cnt) * 0.5, 1) as est_grading_hours
        FROM canvas.enrollments e
        JOIN canvas.users u ON u.id = e.user_id
        JOIN canvas.courses c ON c.id = e.course_id
        LEFT JOIN (SELECT course_id, COUNT(*) as cnt FROM canvas.enrollments WHERE type='StudentEnrollment' GROUP BY course_id) sc ON sc.course_id = c.id
        LEFT JOIN (SELECT course_id, COUNT(*) as cnt FROM canvas.assignments GROUP BY course_id) ac ON ac.course_id = c.id
        WHERE e.type = 'TeacherEnrollment'
        GROUP BY u.name
        ORDER BY u.name
    """)
    rows = cur.fetchall()
    cur.close()
    conn.close()
    return rows


def check_excel(agent_workspace):
    print("\n=== Checking Excel ===")
    xlsx_path = os.path.join(agent_workspace, "Faculty_Workload.xlsx")
    if not os.path.isfile(xlsx_path):
        check("Faculty_Workload.xlsx exists", False, f"Not found: {xlsx_path}")
        return
    check("Faculty_Workload.xlsx exists", True)

    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    except Exception as e:
        check("Excel readable", False, str(e))
        return
    check("Excel readable", True)

    expected = get_expected_instructors()

    # Instructor Load sheet
    il_rows = load_sheet_rows(wb, "Instructor Load")
    if il_rows is None:
        check("Sheet 'Instructor Load' exists", False, f"Available: {wb.sheetnames}")
    else:
        check("Sheet 'Instructor Load' exists", True)
        data_rows = il_rows[1:] if len(il_rows) > 1 else []
        check(f"Instructor Load has {len(expected)} rows",
              abs(len(data_rows) - len(expected)) <= 1,
              f"Found {len(data_rows)}")

        header = il_rows[0] if il_rows else []
        header_lower = [str(h).lower().replace(" ", "_") if h else "" for h in header]
        for col in ["instructor", "courses_count", "total_students", "overloaded_yn"]:
            check(f"Column '{col}' present", any(col in h for h in header_lower),
                  f"Header: {header}")

        # Spot-check Dr. Andrew Walker (3 courses, 2728 students)
        for row in data_rows:
            if row and row[0] and "andrew walker" in str(row[0]).lower():
                check("Dr. Walker courses=3", num_close(row[1], 3), f"Got {row[1]}")
                check("Dr. Walker total_students=2728", num_close(row[2], 2728, 10), f"Got {row[2]}")
                # Weekly hours = 2728 * 0.5 / 16 = 85.25 -> overloaded
                check("Dr. Walker marked overloaded",
                      str(row[-1]).strip().lower() in ("yes", "y", "true"),
                      f"Got {row[-1]}")
                break

        # Check someone who is NOT overloaded (Dr. Abigail Martin: 1 course, 365 students)
        for row in data_rows:
            if row and row[0] and "abigail martin" in str(row[0]).lower():
                check("Dr. Martin courses=1", num_close(row[1], 1), f"Got {row[1]}")
                check("Dr. Martin NOT overloaded",
                      str(row[-1]).strip().lower() in ("no", "n", "false"),
                      f"Got {row[-1]}")
                break

    # Department Summary sheet
    ds_rows = load_sheet_rows(wb, "Department Summary")
    if ds_rows is None:
        check("Sheet 'Department Summary' exists", False, f"Available: {wb.sheetnames}")
    else:
        check("Sheet 'Department Summary' exists", True)
        data_rows = ds_rows[1:] if len(ds_rows) > 1 else []
        check("Department Summary has 7 departments", len(data_rows) == 7,
              f"Found {len(data_rows)}")


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    pptx_path = os.path.join(agent_workspace, "Workload_Review.pptx")
    if not os.path.isfile(pptx_path):
        check("Workload_Review.pptx exists", False, f"Not found: {pptx_path}")
        return
    check("Workload_Review.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        check("PPTX has at least 3 slides", len(prs.slides) >= 3,
              f"Found {len(prs.slides)} slides")
        all_text = " ".join(
            shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
        ).lower()
        check("PPTX mentions workload or overloaded",
              "workload" in all_text or "overload" in all_text,
              f"Sample: {all_text[:200]}")
    except ImportError:
        check("python-pptx available", False)


def check_gsheet():
    print("\n=== Checking Google Sheet ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT id, title FROM gsheet.spreadsheets
            WHERE title ILIKE '%%workload%%' OR title ILIKE '%%faculty%%'
        """)
        sheets = cur.fetchall()
        check("Google Sheet created for workload data", len(sheets) >= 1,
              "No matching spreadsheet found")
        if sheets:
            sid = sheets[0][0]
            cur.execute("SELECT COUNT(*) FROM gsheet.cells WHERE spreadsheet_id = %s", (sid,))
            cell_count = cur.fetchone()[0]
            check("Google Sheet has data (cells)", cell_count > 10,
                  f"Found {cell_count} cells")
        cur.close()
        conn.close()
    except Exception as e:
        check("Google Sheet check", False, str(e))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_gsheet()

    total = PASS_COUNT + FAIL_COUNT
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ===")
    if FAIL_COUNT > 0:
        print(f"{FAIL_COUNT} checks failed")
        sys.exit(1)
    else:
        print("All checks passed!")
        sys.exit(0)


if __name__ == "__main__":
    main()
