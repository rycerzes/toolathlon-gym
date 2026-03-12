"""Evaluation for canvas-quiz-analysis-ppt."""
import argparse
import os
import sys
import openpyxl
from pptx import Presentation


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")
    all_errors = []

    # ---- Check Excel ----
    agent_excel = os.path.join(args.agent_workspace, "Quiz_Performance.xlsx")
    gt_excel = os.path.join(gt_dir, "Quiz_Performance.xlsx")

    if not os.path.exists(agent_excel):
        all_errors.append("Agent output Quiz_Performance.xlsx not found")
    elif not os.path.exists(gt_excel):
        all_errors.append("Groundtruth Quiz_Performance.xlsx not found")
    else:
        agent_wb = openpyxl.load_workbook(agent_excel, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_excel, data_only=True)

        # Check Quiz Details sheet
        print("  Checking Quiz Details...")
        a_rows = load_sheet_rows(agent_wb, "Quiz Details")
        g_rows = load_sheet_rows(gt_wb, "Quiz Details")
        if a_rows is None:
            all_errors.append("Sheet 'Quiz Details' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Quiz Details' not found in groundtruth")
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []
            if len(a_data) != len(g_data):
                all_errors.append(f"Quiz Details row count: agent={len(a_data)}, expected={len(g_data)}")
            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None and row[1] is not None:
                    key = (str(row[0]).strip().lower(), str(row[1]).strip().lower())
                    a_lookup[key] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = (str(g_row[0]).strip().lower(), str(g_row[1]).strip().lower())
                a_row = a_lookup.get(key)
                if a_row is None:
                    all_errors.append(f"Missing quiz row: {g_row[0]} / {g_row[1]}")
                    continue
                if len(a_row) > 2 and len(g_row) > 2:
                    if not num_close(a_row[2], g_row[2], 5):
                        all_errors.append(f"{key}.Submissions: {a_row[2]} vs {g_row[2]}")
                if len(a_row) > 3 and len(g_row) > 3:
                    if not num_close(a_row[3], g_row[3], 1.0):
                        all_errors.append(f"{key}.Avg_Score: {a_row[3]} vs {g_row[3]}")
                if len(a_row) > 4 and len(g_row) > 4:
                    if not num_close(a_row[4], g_row[4], 1.0):
                        all_errors.append(f"{key}.Min_Score: {a_row[4]} vs {g_row[4]}")
                if len(a_row) > 5 and len(g_row) > 5:
                    if not num_close(a_row[5], g_row[5], 1.0):
                        all_errors.append(f"{key}.Max_Score: {a_row[5]} vs {g_row[5]}")
            if not all_errors:
                print("    PASS")

        # Check Course Summary sheet
        print("  Checking Course Summary...")
        a_rows = load_sheet_rows(agent_wb, "Course Summary")
        g_rows = load_sheet_rows(gt_wb, "Course Summary")
        if a_rows is None:
            all_errors.append("Sheet 'Course Summary' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Course Summary' not found in groundtruth")
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []
            if len(a_data) != len(g_data):
                all_errors.append(f"Course Summary row count: agent={len(a_data)}, expected={len(g_data)}")
            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None:
                    a_lookup[str(row[0]).strip().lower()] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = str(g_row[0]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    all_errors.append(f"Missing course summary: {g_row[0]}")
                    continue
                if len(a_row) > 1 and len(g_row) > 1:
                    if not num_close(a_row[1], g_row[1], 1):
                        all_errors.append(f"{key}.Total_Quizzes: {a_row[1]} vs {g_row[1]}")
                if len(a_row) > 2 and len(g_row) > 2:
                    if not num_close(a_row[2], g_row[2], 10):
                        all_errors.append(f"{key}.Total_Submissions: {a_row[2]} vs {g_row[2]}")
                if len(a_row) > 3 and len(g_row) > 3:
                    if not num_close(a_row[3], g_row[3], 1.0):
                        all_errors.append(f"{key}.Overall_Avg_Score: {a_row[3]} vs {g_row[3]}")
            if not any("Course Summary" in e for e in all_errors):
                print("    PASS")

    # ---- Check PowerPoint ----
    agent_ppt = os.path.join(args.agent_workspace, "Quiz_Report.pptx")
    if not os.path.exists(agent_ppt):
        all_errors.append("Agent output Quiz_Report.pptx not found")
    else:
        print("  Checking Quiz_Report.pptx...")
        prs = Presentation(agent_ppt)
        slides = list(prs.slides)
        # Minimum: title + 3 courses + takeaways = 5
        if len(slides) < 5:
            all_errors.append(f"PPT has {len(slides)} slides, expected at least 5")
        else:
            # Check title slide
            title_text = ""
            for shape in slides[0].shapes:
                if shape.has_text_frame:
                    title_text += shape.text_frame.text.lower() + " "
            if "fall 2014" not in title_text or "quiz" not in title_text:
                all_errors.append(f"Title slide missing expected text. Found: {title_text[:100]}")

            # Check last slide mentions key takeaways
            last_text = ""
            for shape in slides[-1].shapes:
                if shape.has_text_frame:
                    last_text += shape.text_frame.text.lower() + " "
            if "takeaway" not in last_text and "key" not in last_text:
                # Allow "summary" or "conclusion" as alternatives
                if "summary" not in last_text and "conclusion" not in last_text:
                    all_errors.append(f"Last slide missing takeaways content")

            # Check that course names appear somewhere in the presentation
            all_ppt_text = ""
            for slide in slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        all_ppt_text += shape.text_frame.text.lower() + " "
            for course in ["creative computing", "foundations of finance", "global governance"]:
                if course not in all_ppt_text:
                    all_errors.append(f"PPT missing course: {course}")

            # Check highest/lowest mentioned
            if "global governance" not in last_text:
                all_errors.append("Last slide should mention Global Governance (highest avg)")
            if "creative computing" not in last_text:
                all_errors.append("Last slide should mention Creative Computing (lowest avg)")

        if not any("PPT" in e or "ppt" in e.lower() or "slide" in e.lower() for e in all_errors):
            print("    PASS")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:15]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
