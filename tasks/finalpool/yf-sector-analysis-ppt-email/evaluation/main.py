"""Evaluation for yf-sector-analysis-ppt-email.

Blocking checks: Sector_Analysis.xlsx (Excel data) and Sector_Analysis.pptx (PPT structure).
Non-blocking: Email DB check.
"""
import argparse
import os
import sys
import openpyxl
from pptx import Presentation


TICKERS = ["AMZN", "GOOGL", "JNJ", "JPM", "XOM"]
SECTORS = {"AMZN": "Technology", "GOOGL": "Technology", "XOM": "Energy",
           "JNJ": "Healthcare", "JPM": "Financials"}


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")
    all_errors = []

    # ---- Check Excel ----
    agent_file = os.path.join(args.agent_workspace, "Sector_Analysis.xlsx")
    gt_file = os.path.join(gt_dir, "Sector_Analysis.xlsx")

    if not os.path.exists(agent_file):
        all_errors.append("Agent output Sector_Analysis.xlsx not found")
    elif not os.path.exists(gt_file):
        all_errors.append("Groundtruth Sector_Analysis.xlsx not found")
    else:
        agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_file, data_only=True)

        # Check Stock Performance sheet
        print("  Checking Stock Performance...")
        a_rows = load_sheet_rows(agent_wb, "Stock Performance")
        g_rows = load_sheet_rows(gt_wb, "Stock Performance")
        if a_rows is None:
            all_errors.append("Sheet 'Stock Performance' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Stock Performance' not found in groundtruth")
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []

            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None:
                    a_lookup[str(row[0]).strip().upper()] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = str(g_row[0]).strip().upper()
                a_row = a_lookup.get(key)
                if a_row is None:
                    all_errors.append(f"Missing stock: {key}")
                    continue
                # Col 1: Sector
                if len(a_row) > 1 and len(g_row) > 1:
                    if not str_match(a_row[1], g_row[1]):
                        all_errors.append(f"{key}.Sector: {a_row[1]} vs {g_row[1]}")
                # Col 2: Latest_Close
                if len(a_row) > 2 and len(g_row) > 2:
                    if not num_close(a_row[2], g_row[2], 5.0):
                        all_errors.append(f"{key}.Latest_Close: {a_row[2]} vs {g_row[2]} (tol=5.0)")
                # Col 3: YTD_Return_Pct
                if len(a_row) > 3 and len(g_row) > 3:
                    if not num_close(a_row[3], g_row[3], 3.0):
                        all_errors.append(f"{key}.YTD_Return_Pct: {a_row[3]} vs {g_row[3]} (tol=3.0)")
                # Col 4: One_Year_Return_Pct
                if len(a_row) > 4 and len(g_row) > 4:
                    if not num_close(a_row[4], g_row[4], 3.0):
                        all_errors.append(f"{key}.One_Year_Return_Pct: {a_row[4]} vs {g_row[4]} (tol=3.0)")
            if not all_errors:
                print("    PASS")

        # Check Sector Summary sheet
        print("  Checking Sector Summary...")
        a_rows = load_sheet_rows(agent_wb, "Sector Summary")
        g_rows = load_sheet_rows(gt_wb, "Sector Summary")
        prev_errors = len(all_errors)
        if a_rows is None:
            all_errors.append("Sheet 'Sector Summary' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Sector Summary' not found in groundtruth")
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []

            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None:
                    a_lookup[str(row[0]).strip().lower()] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = str(g_row[0]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    all_errors.append(f"Missing sector: {g_row[0]}")
                    continue
                # Col 1: Num_Stocks
                if len(a_row) > 1 and len(g_row) > 1:
                    if not num_close(a_row[1], g_row[1], 0):
                        all_errors.append(f"{key}.Num_Stocks: {a_row[1]} vs {g_row[1]}")
                # Col 2: Avg_YTD_Return_Pct
                if len(a_row) > 2 and len(g_row) > 2:
                    if not num_close(a_row[2], g_row[2], 3.0):
                        all_errors.append(f"{key}.Avg_YTD_Return: {a_row[2]} vs {g_row[2]} (tol=3.0)")
            new_errors = len(all_errors) - prev_errors
            if new_errors == 0:
                print("    PASS")

    # ---- Check PowerPoint ----
    agent_ppt = os.path.join(args.agent_workspace, "Sector_Analysis.pptx")
    if not os.path.exists(agent_ppt):
        all_errors.append("Agent output Sector_Analysis.pptx not found")
    else:
        print("  Checking Sector_Analysis.pptx...")
        prs = Presentation(agent_ppt)
        slides = list(prs.slides)
        if len(slides) < 4:
            all_errors.append(f"PPT has {len(slides)} slides, expected at least 4")
        else:
            # Check title slide
            title_text = ""
            for shape in slides[0].shapes:
                if shape.has_text_frame:
                    title_text += shape.text_frame.text.lower() + " "
            if "sector" not in title_text:
                all_errors.append(f"Title slide missing 'sector'. Found: {title_text[:100]}")

            # Check all PPT text for key content
            all_ppt_text = ""
            for slide in slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        all_ppt_text += shape.text_frame.text.lower() + " "

            for ticker in TICKERS:
                if ticker.lower() not in all_ppt_text:
                    all_errors.append(f"PPT missing ticker: {ticker}")

            if "key findings" not in all_ppt_text and "key" not in all_ppt_text:
                if "summary" not in all_ppt_text and "conclusion" not in all_ppt_text:
                    all_errors.append("PPT missing Key Findings slide")

        if not any("ppt" in e.lower() or "slide" in e.lower() or "ticker" in e.lower() for e in all_errors[prev_errors:] if "ppt" in e.lower() or "slide" in e.lower()):
            ppt_errors = [e for e in all_errors if "ppt" in e.lower() or "slide" in e.lower() or "ticker" in e.lower()]
            if not ppt_errors:
                print("    PASS")

    # ---- Non-blocking Email check ----
    print("  Non-blocking: Email DB check...")
    try:
        import psycopg2
        conn = psycopg2.connect(host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym",
                                user="eigent", password="camel")
        cur = conn.cursor()
        cur.execute("SELECT COUNT(*) FROM email.sent_log")
        count = cur.fetchone()[0]
        print(f"    [INFO] Found {count} sent email(s) (non-blocking)")
        cur.close()
        conn.close()
    except Exception as e:
        print(f"    [INFO] Email check skipped: {e} (non-blocking)")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:15]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
