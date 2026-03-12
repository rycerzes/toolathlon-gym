"""
Evaluation for yt-fireship-scholarly-ppt-gsheet-email task.

Checks:
1. AI_Tech_Overview_2025.pptx exists with >= 4 slides and AI-related content
2. GSheet has spreadsheet with Videos sheet with >= 4 rows
3. GSheet has Papers sheet with >= 3 rows
4. GSheet has Technology_Map sheet with >= 3 rows
5. Email sent to research@ai-group.org
6. Email sent to team@lab.edu
"""
import json
import os
import sys
from argparse import ArgumentParser

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    try: return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError): return False


def str_match(a, b):
    if a is None or b is None: return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def check_pptx(agent_workspace):
    print("\n=== Check 1: PPT AI_Tech_Overview_2025.pptx ===")

    pptx_path = os.path.join(agent_workspace, "AI_Tech_Overview_2025.pptx")
    if not os.path.exists(pptx_path):
        record("AI_Tech_Overview_2025.pptx exists", False, f"Not found at {pptx_path}")
        return
    record("AI_Tech_Overview_2025.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        record("PPT has >= 4 slides", slide_count >= 4, f"Found {slide_count} slides")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += " " + shape.text
        all_text_lower = all_text.lower()

        # Keywords from real 2024+ Fireship AI videos and academic papers
        ai_terms = [
            "deepseek", "openai", "claude", "grok", "ai", "llm",
            "transformer", "language model", "machine learning", "neural",
            "gpt", "bert", "vibe coding"
        ]
        found_terms = [t for t in ai_terms if t in all_text_lower]
        record("PPT contains >= 3 AI/ML-related terms", len(found_terms) >= 3,
               f"Found terms: {found_terms}")

        has_overview = any(kw in all_text_lower for kw in ["overview", "research", "technology", "2024", "2025"])
        record("PPT contains overview/research/technology content", has_overview,
               f"Text sample: {all_text[:300]}")

    except ImportError:
        record("python-pptx available", False, "python-pptx not installed")
    except Exception as e:
        record("PPT readable", False, str(e))


def check_gsheet():
    print("\n=== Check 2: GSheet AI_Research_Overview ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Find the spreadsheet
    cur.execute("SELECT id, title FROM gsheet.spreadsheets")
    spreadsheets = cur.fetchall()

    ai_spreadsheet = None
    for ss_id, title in spreadsheets:
        if title and any(kw in title.lower() for kw in ["ai", "research", "overview", "fireship"]):
            ai_spreadsheet = (ss_id, title)
            break

    if not ai_spreadsheet and spreadsheets:
        # Take any spreadsheet
        ai_spreadsheet = spreadsheets[0]

    record("GSheet spreadsheet exists", ai_spreadsheet is not None,
           f"Spreadsheets found: {[s[1] for s in spreadsheets]}")

    if not ai_spreadsheet:
        cur.close()
        conn.close()
        return

    ss_id = ai_spreadsheet[0]

    # Check sheets
    cur.execute("SELECT id, title FROM gsheet.sheets WHERE spreadsheet_id = %s", (ss_id,))
    sheets = cur.fetchall()
    sheet_titles_lower = [(s[0], s[1].lower() if s[1] else "") for s in sheets]

    has_videos = any("video" in t for _, t in sheet_titles_lower)
    has_papers = any("paper" in t for _, t in sheet_titles_lower)
    has_tech_map = any("tech" in t or "map" in t or "technology" in t for _, t in sheet_titles_lower)

    record("GSheet has Videos sheet", has_videos, f"Sheets: {[s[1] for s in sheets]}")
    record("GSheet has Papers sheet", has_papers, f"Sheets: {[s[1] for s in sheets]}")
    record("GSheet has Technology_Map sheet", has_tech_map, f"Sheets: {[s[1] for s in sheets]}")

    def count_data_rows(sheet_id):
        cur.execute("""
            SELECT COUNT(DISTINCT row_index) FROM gsheet.cells
            WHERE sheet_id = %s AND row_index > 0
        """, (sheet_id,))
        return cur.fetchone()[0]

    if has_videos:
        video_sheet_id = next(sid for sid, t in sheet_titles_lower if "video" in t)
        video_rows = count_data_rows(video_sheet_id)
        record("Videos sheet has >= 4 data rows", video_rows >= 4,
               f"Found {video_rows} rows")

    if has_papers:
        paper_sheet_id = next(sid for sid, t in sheet_titles_lower if "paper" in t)
        paper_rows = count_data_rows(paper_sheet_id)
        record("Papers sheet has >= 3 data rows", paper_rows >= 3,
               f"Found {paper_rows} rows")

    if has_tech_map:
        tech_sheet_id = next(sid for sid, t in sheet_titles_lower if "tech" in t or "map" in t or "technology" in t)
        tech_rows = count_data_rows(tech_sheet_id)
        record("Technology_Map sheet has >= 3 data rows", tech_rows >= 3,
               f"Found {tech_rows} rows")

    cur.close()
    conn.close()


def check_xlsx_content(workspace, groundtruth_workspace="."):
    print("\n=== Check: XLSX AI_Research_Overview_local.xlsx ===")
    import openpyxl
    xlsx_path = os.path.join(workspace, "AI_Research_Overview_local.xlsx")
    if not os.path.isfile(xlsx_path):
        record("xlsx AI_Research_Overview_local.xlsx exists", False, f"Not found at {xlsx_path}")
        return False
    record("xlsx AI_Research_Overview_local.xlsx exists", True)
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        for ws in wb.worksheets:
            rows = list(ws.iter_rows(values_only=True))
            record(f"xlsx '{ws.title}' has data", len(rows) >= 2, f"{len(rows)} rows")

        # --- Groundtruth XLSX value comparison ---
        gt_path = os.path.join(groundtruth_workspace, "AI_Research_Overview_local.xlsx")
        if os.path.isfile(gt_path):
            gt_wb = openpyxl.load_workbook(gt_path, data_only=True)
            for gt_sname in gt_wb.sheetnames:
                gt_ws = gt_wb[gt_sname]
                a_ws = None
                for asn in wb.sheetnames:
                    if asn.strip().lower() == gt_sname.strip().lower():
                        a_ws = wb[asn]; break
                if a_ws is None:
                    record(f"GT sheet '{gt_sname}' exists in agent xlsx", False, f"Available: {wb.sheetnames}")
                    continue
                gt_rows = [r for r in gt_ws.iter_rows(min_row=2, values_only=True) if any(c is not None for c in r)]
                a_rows = [r for r in a_ws.iter_rows(min_row=2, values_only=True) if any(c is not None for c in r)]
                record(f"GT '{gt_sname}' row count", len(a_rows) == len(gt_rows),
                       f"Expected {len(gt_rows)}, got {len(a_rows)}")
                for ri in range(min(3, len(gt_rows))):
                    if ri >= len(a_rows): break
                    ok = True
                    for ci in range(min(len(gt_rows[ri]), len(a_rows[ri]))):
                        gv, av = gt_rows[ri][ci], a_rows[ri][ci]
                        if gv is None: continue
                        if isinstance(gv, (int, float)):
                            if not num_close(av, gv, max(abs(gv)*0.1, 1.0)): ok = False; break
                        else:
                            if not str_match(av, gv): ok = False; break
                    record(f"GT '{gt_sname}' row {ri+1} values", ok,
                           f"gt={gt_rows[ri][:4]}, agent={a_rows[ri][:4] if ri < len(a_rows) else 'missing'}")
            gt_wb.close()

        wb.close()
    except Exception as e:
        record("xlsx readable", False, str(e))
    return True


def check_emails():
    print("\n=== Check 3: Emails sent ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
    messages = cur.fetchall()
    cur.close()
    conn.close()

    def to_addresses(to_addr):
        if isinstance(to_addr, list):
            return " ".join(str(r).lower() for r in to_addr)
        elif to_addr:
            try:
                parsed = json.loads(str(to_addr))
                return " ".join(str(r).lower() for r in parsed) if isinstance(parsed, list) else str(to_addr).lower()
            except Exception:
                return str(to_addr).lower()
        return ""

    to_research = [m for m in messages if "research@ai-group.org" in to_addresses(m[2])]
    to_team = [m for m in messages if "team@lab.edu" in to_addresses(m[2])]

    record("Email sent to research@ai-group.org", len(to_research) >= 1,
           f"Total messages: {len(messages)}")
    record("Email sent to team@lab.edu", len(to_team) >= 1,
           f"Total messages: {len(messages)}")

    if to_research:
        subj, _, _, body = to_research[0]
        content = ((subj or "") + " " + (body or "")).lower()
        # Check for real 2024+ AI keywords (DeepSeek, OpenAI, Claude, Grok, AI)
        has_ai = any(kw in content for kw in [
            "deepseek", "openai", "claude", "grok", "ai", "technology",
            "research", "overview", "presentation", "llm"
        ])
        record("Research email mentions AI/technology content", has_ai, f"Subject: {subj}")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_xlsx_content(args.agent_workspace, args.groundtruth_workspace)
    check_gsheet()
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
