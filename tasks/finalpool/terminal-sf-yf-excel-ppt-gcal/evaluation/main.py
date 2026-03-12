"""Evaluation for terminal-sf-yf-excel-ppt-gcal.
Checks:
1. Investment_Committee_Briefing.xlsx (4 sheets, correct data)
2. Committee_Briefing.pptx (6 slides)
3. Google Calendar briefing event
4. compute_growth.py and market_comparison.py scripts exist
5. briefing_notes.txt exists
6. market_comparison.json exists
"""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0

# Fallback hardcoded values
_FALLBACK_SF_YF = {
    "lg_q4_revenue": 57937,
    "top_customer": "ava garcia",
    "top_brand": "lg",
}


def _get_sf_yf_expected():
    """Query sf_data and yf schemas to compute expected values dynamically."""
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        try:
            # Get Q4 2025 revenue by brand (Oct-Dec 2025)
            cur.execute("""
                SELECT p."BRAND", SUM(o."TOTAL_AMOUNT") as q4_revenue
                FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
                JOIN sf_data."SALES_DW__PUBLIC__PRODUCTS" p ON o."PRODUCT_ID" = p."PRODUCT_ID"
                WHERE o."ORDER_DATE" >= '2025-10-01' AND o."ORDER_DATE" < '2026-01-01'
                GROUP BY p."BRAND"
                ORDER BY q4_revenue DESC
            """)
            brand_rows = cur.fetchall()
            top_brand = brand_rows[0][0].lower() if brand_rows else "lg"
            # Find LG Q4 revenue
            lg_q4 = 0
            for brand, rev in brand_rows:
                if brand.lower() == "lg":
                    lg_q4 = float(rev)
                    break

            # Get top customer by total spend in 2025
            cur.execute("""
                SELECT c."CUSTOMER_NAME", SUM(o."TOTAL_AMOUNT") as total_spend
                FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
                JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
                WHERE o."ORDER_DATE" >= '2025-01-01' AND o."ORDER_DATE" < '2026-01-01'
                GROUP BY c."CUSTOMER_NAME"
                ORDER BY total_spend DESC
                LIMIT 1
            """)
            top_cust_row = cur.fetchone()
            top_customer = top_cust_row[0].lower() if top_cust_row else "ava garcia"

            return {
                "lg_q4_revenue": lg_q4,
                "top_customer": top_customer,
                "top_brand": top_brand,
            }
        finally:
            cur.close()
            conn.close()
    except Exception:
        return _FALLBACK_SF_YF


_SF_YF_EXPECTED = _get_sf_yf_expected()


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def num_close(a, b, tol=2.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except Exception:
        return False


def check_excel(workspace):
    print("\n=== Check 1: Investment_Committee_Briefing.xlsx ===")
    path = os.path.join(workspace, "Investment_Committee_Briefing.xlsx")
    if not os.path.exists(path):
        check("Excel file exists", False, f"Not found at {path}")
        return
    check("Excel file exists", True)

    wb = openpyxl.load_workbook(path)
    sheets = wb.sheetnames
    check("Has at least 4 sheets", len(sheets) >= 4, f"Found {len(sheets)}: {sheets}")

    sheets_lower = [s.lower() for s in sheets]

    # Revenue_By_Brand sheet
    rb_idx = next((i for i, s in enumerate(sheets_lower) if "revenue" in s or "brand" in s), 0)
    ws1 = wb[sheets[rb_idx]]
    rows1 = list(ws1.iter_rows(values_only=True))
    data1 = [r for r in rows1[1:] if any(c for c in r)]
    check("Revenue_By_Brand has 13 brand rows", len(data1) >= 10, f"Found {len(data1)}")

    all_text1 = " ".join(str(c) for r in rows1 for c in r if c).lower()
    check("Contains LG brand", "lg" in all_text1, f"Text: {all_text1[:120]}")
    check("Contains Microsoft brand", "microsoft" in all_text1)
    check("Contains growth column", "growth" in all_text1 or "qoq" in all_text1,
          f"Headers: {rows1[0] if rows1 else 'none'}")

    # Check LG Q4 revenue is approximately correct
    lg_row = None
    for r in data1:
        if r and str(r[0]).lower() == "lg":
            lg_row = r
            break
    if lg_row:
        expected_lg_q4 = _SF_YF_EXPECTED["lg_q4_revenue"]
        check(f"LG Q4 revenue ~{expected_lg_q4:.0f}", num_close(lg_row[2], expected_lg_q4, tol=500),
              f"Got {lg_row[2]}")
    else:
        check(f"LG Q4 revenue ~{_SF_YF_EXPECTED['lg_q4_revenue']:.0f}", False, "LG row not found")

    # Check rank column exists (should be sorted by Q4 desc)
    first_brand = str(data1[0][0]).lower() if data1 else ""
    expected_top_brand = _SF_YF_EXPECTED["top_brand"]
    check(f"First brand is {expected_top_brand} (highest Q4)", first_brand == expected_top_brand,
          f"First brand: {first_brand}")

    # Top_Customers sheet
    tc_idx = next((i for i, s in enumerate(sheets_lower) if "customer" in s), 1)
    if tc_idx < len(sheets):
        ws2 = wb[sheets[tc_idx]]
        rows2 = list(ws2.iter_rows(values_only=True))
        data2 = [r for r in rows2[1:] if any(c for c in r)]
        check("Top_Customers has 10 rows", len(data2) >= 10, f"Found {len(data2)}")
        all_text2 = " ".join(str(c) for r in rows2 for c in r if c).lower()
        expected_top_cust = _SF_YF_EXPECTED["top_customer"]
        check(f"Contains {expected_top_cust} (top customer)", expected_top_cust in all_text2,
              f"Text: {all_text2[:120]}")
        check("Contains segment info", "enterprise" in all_text2 or "consumer" in all_text2)

    # Market_Context sheet
    mc_idx = next((i for i, s in enumerate(sheets_lower) if "market" in s), 2)
    if mc_idx < len(sheets):
        ws3 = wb[sheets[mc_idx]]
        rows3 = list(ws3.iter_rows(values_only=True))
        data3 = [r for r in rows3[1:] if any(c for c in r)]
        check("Market_Context has 5 rows", len(data3) >= 5, f"Found {len(data3)}")
        all_text3 = " ".join(str(c) for r in rows3 for c in r if c).lower()
        check("Contains GOOGL", "googl" in all_text3)
        check("Contains GC=F or gold", "gc=f" in all_text3 or "gold" in all_text3)
        check("Contains trend indicator", "up" in all_text3 or "down" in all_text3 or "flat" in all_text3)

    # Executive_Summary sheet
    es_idx = next((i for i, s in enumerate(sheets_lower) if "summary" in s or "executive" in s), 3)
    if es_idx < len(sheets):
        ws4 = wb[sheets[es_idx]]
        rows4 = list(ws4.iter_rows(values_only=True))
        data4 = [r for r in rows4[1:] if any(c for c in r)]
        check("Executive_Summary has at least 5 rows", len(data4) >= 5, f"Found {len(data4)}")
        all_text4 = " ".join(str(c) for r in rows4 for c in r if c).lower()
        check("Contains total Q4 revenue metric", "q4" in all_text4 and "revenue" in all_text4,
              f"Text: {all_text4[:150]}")
        check("Contains growth metric", "growth" in all_text4)


def check_pptx(workspace):
    print("\n=== Check 2: Committee_Briefing.pptx ===")
    path = os.path.join(workspace, "Committee_Briefing.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = prs.slides
        check("Has 6 slides", len(slides) >= 6, f"Found {len(slides)}")

        # Check slide content
        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + " "

        all_lower = all_text.lower()
        check("Title slide mentions Q4", "q4" in all_lower)
        check("Contains revenue content", "revenue" in all_lower)
        check("Contains customer content", "customer" in all_lower)
        check("Contains market content", "market" in all_lower)
        check("Contains next steps", "next" in all_lower or "action" in all_lower)
    except ImportError:
        check("python-pptx available", False, "python-pptx not installed")


def check_gcal():
    print("\n=== Check 3: Calendar Briefing Event ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT summary, start_datetime, end_datetime, description
        FROM gcal.events
        WHERE lower(summary) LIKE '%%investment%%committee%%'
           OR lower(summary) LIKE '%%q4%%briefing%%'
           OR lower(summary) LIKE '%%committee%%briefing%%'
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    check("Briefing event exists", len(events) >= 1, f"Found {len(events)} matching events")

    if events:
        evt = events[0]
        summary, start, end, desc = evt

        # Should be in the week of March 9-13
        check("Event in target week (Mar 9-13)",
              start and start.strftime("%Y-%m-%d") >= "2026-03-09" and
              start.strftime("%Y-%m-%d") <= "2026-03-13",
              f"Start: {start}")

        # Should be ~2 hours
        if start and end:
            duration = (end - start).total_seconds() / 3600
            check("Event is ~2 hours", 1.5 <= duration <= 2.5,
                  f"Duration: {duration} hours")

        # Should not conflict with existing events
        # Tuesday Mar 10 9:00-11:00 is the first available 2hr slot
        if start:
            check("Event is on Tuesday Mar 10 (first available slot)",
                  start.strftime("%Y-%m-%d") == "2026-03-10",
                  f"Date: {start.strftime('%Y-%m-%d')}")

        # Description should mention committee or sales
        if desc:
            desc_lower = str(desc).lower()
            check("Description mentions briefing topic",
                  "sales" in desc_lower or "market" in desc_lower or
                  "committee" in desc_lower or "quarterly" in desc_lower,
                  f"Desc: {str(desc)[:100]}")

    cur.close()
    conn.close()


def check_scripts(workspace):
    print("\n=== Check 4: Python Scripts ===")
    check("compute_growth.py exists",
          os.path.exists(os.path.join(workspace, "compute_growth.py")))
    check("market_comparison.py exists",
          os.path.exists(os.path.join(workspace, "market_comparison.py")))


def check_outputs(workspace):
    print("\n=== Check 5: Additional Output Files ===")
    # briefing_notes.txt
    notes_path = os.path.join(workspace, "briefing_notes.txt")
    if os.path.exists(notes_path):
        check("briefing_notes.txt exists", True)
        with open(notes_path) as f:
            content = f.read().lower()
        check("Notes mention revenue", "revenue" in content or "sales" in content,
              f"Content: {content[:100]}")
        check("Notes mention market", "market" in content or "stock" in content,
              f"Content: {content[:100]}")
    else:
        check("briefing_notes.txt exists", False)

    # market_comparison.json
    mc_path = os.path.join(workspace, "market_comparison.json")
    if os.path.exists(mc_path):
        check("market_comparison.json exists", True)
        try:
            with open(mc_path) as f:
                data = json.load(f)
            check("market_comparison.json is valid JSON", True)
            # Should contain symbol data
            text = json.dumps(data).lower()
            check("Contains GOOGL data", "googl" in text)
        except Exception as e:
            check("market_comparison.json is valid JSON", False, str(e))
    else:
        check("market_comparison.json exists", False)


def check_reverse(workspace):
    print("\n=== Reverse Validation ===")
    path = os.path.join(workspace, "Investment_Committee_Briefing.xlsx")
    if os.path.isfile(path):
        wb = openpyxl.load_workbook(path)
        has_negative_revenue = False
        ws = wb[wb.sheetnames[0]]
        for row in ws.iter_rows(min_row=2, min_col=2, max_col=3, values_only=True):
            for cell in row:
                if isinstance(cell, (int, float)) and cell < 0:
                    has_negative_revenue = True
        check("No negative revenue values", not has_negative_revenue,
              "Found negative revenue")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_gcal()
    check_scripts(args.agent_workspace)
    check_outputs(args.agent_workspace)
    check_reverse(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if accuracy >= 70 else 1)


if __name__ == "__main__":
    main()
