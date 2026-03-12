"""
Evaluation script for fetch-sf-sales-forecast-ppt-gcal task.

Checks:
1. Sales_Forecast_Data.xlsx with Q1_Actuals, Q2_Forecast, Segment_Mix sheets
2. Q2_Sales_Forecast.pptx with forecast content
3. Calendar event for board presentation
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=5000.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_contains(haystack, needle):
    if haystack is None or needle is None:
        return False
    return needle.strip().lower() in str(haystack).strip().lower()


def check_excel(agent_workspace):
    """Check Sales_Forecast_Data.xlsx."""
    print("\n=== Checking Excel Output ===")

    agent_file = os.path.join(agent_workspace, "Sales_Forecast_Data.xlsx")
    if not os.path.isfile(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False

    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(agent_file, data_only=True)
    except Exception as e:
        record("Excel file readable", False, str(e))
        return False

    all_ok = True

    # Check Q1_Actuals sheet
    q1_sheet = None
    for name in wb.sheetnames:
        if "q1" in name.lower() or "actual" in name.lower():
            q1_sheet = name
            break

    if not q1_sheet:
        record("Q1_Actuals sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Q1_Actuals sheet exists", True)
        ws = wb[q1_sheet]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if r and r[0]] if len(rows) > 1 else []
        record(
            "Q1_Actuals has 15 rows (5 regions x 3 months)",
            len(data_rows) >= 15,
            f"Found {len(data_rows)} data rows",
        )
        if len(data_rows) < 15:
            all_ok = False

        # Spot check: Asia Pacific Jan revenue ~57557
        for r in data_rows:
            if r[0] and "asia" in str(r[0]).lower() and r[1] in (1, "1", "January"):
                ok = num_close(r[3] if len(r) > 3 else 0, 57557.80, tol=5000)
                record("Asia Pacific Jan revenue ~57558", ok, f"Got {r[3] if len(r) > 3 else 'N/A'}")
                if not ok:
                    all_ok = False
                break

    # Check Q2_Forecast sheet
    q2_sheet = None
    for name in wb.sheetnames:
        if "q2" in name.lower() or "forecast" in name.lower():
            q2_sheet = name
            break

    if not q2_sheet:
        record("Q2_Forecast sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Q2_Forecast sheet exists", True)
        ws2 = wb[q2_sheet]
        rows2 = list(ws2.iter_rows(values_only=True))
        data_rows2 = [r for r in rows2[1:] if r and r[0]] if len(rows2) > 1 else []
        record(
            "Q2_Forecast has 5 region rows",
            len(data_rows2) >= 5,
            f"Found {len(data_rows2)} data rows",
        )
        if len(data_rows2) < 5:
            all_ok = False

        # Check growth rates present
        has_growth = False
        for r in data_rows2:
            if r and len(r) >= 3:
                try:
                    gr = float(r[2])
                    if 2.0 <= gr <= 10.0:
                        has_growth = True
                        break
                except (TypeError, ValueError):
                    continue
        record("Growth rates present", has_growth)

    # Check Segment_Mix sheet
    seg_sheet = None
    for name in wb.sheetnames:
        if "segment" in name.lower() or "mix" in name.lower():
            seg_sheet = name
            break

    if not seg_sheet:
        record("Segment_Mix sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Segment_Mix sheet exists", True)
        ws3 = wb[seg_sheet]
        rows3 = list(ws3.iter_rows(values_only=True))
        data_rows3 = [r for r in rows3[1:] if r and r[0]] if len(rows3) > 1 else []
        record(
            "Segment_Mix has >= 20 rows (5 regions x 4 segments)",
            len(data_rows3) >= 20,
            f"Found {len(data_rows3)} data rows",
        )
        if len(data_rows3) < 20:
            all_ok = False

    wb.close()
    return all_ok


def check_pptx(agent_workspace):
    """Check Q2_Sales_Forecast.pptx."""
    print("\n=== Checking PowerPoint ===")

    pptx_file = os.path.join(agent_workspace, "Q2_Sales_Forecast.pptx")
    if not os.path.isfile(pptx_file):
        record("PowerPoint file exists", False, f"Not found: {pptx_file}")
        return False

    record("PowerPoint file exists", True)

    if Presentation is None:
        record("python-pptx available", False, "Cannot import pptx")
        return True  # File exists, can't verify content

    try:
        prs = Presentation(pptx_file)
        slides = prs.slides

        record(
            "PPT has >= 4 slides",
            len(slides) >= 4,
            f"Found {len(slides)} slides",
        )

        # Check content across all slides
        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text.lower() + " "

        record("PPT mentions Q2 forecast", "q2" in all_text and "forecast" in all_text)
        record(
            "PPT mentions regions",
            any(r in all_text for r in ["asia", "europe", "north america"]),
        )
        record(
            "PPT mentions growth",
            "growth" in all_text or "%" in all_text or "projection" in all_text,
        )

        return True
    except Exception as e:
        record("PPT readable", False, str(e))
        return False


def check_calendar():
    """Check calendar event for board presentation."""
    print("\n=== Checking Google Calendar ===")

    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute(
            "SELECT summary, description, start_datetime, end_datetime FROM gcal.events"
        )
        events = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        record("Calendar DB accessible", False, str(e))
        return False

    found = False
    for summary, description, start_dt, end_dt in events:
        summary_lower = (summary or "").lower()
        if ("board" in summary_lower or "forecast" in summary_lower or "sales" in summary_lower) and (
            "presentation" in summary_lower or "meeting" in summary_lower or "q2" in summary_lower
        ):
            found = True
            record("Board presentation event exists", True)

            # Check date is March 28, 2026
            dt_str = str(start_dt)
            record(
                "Event on March 28, 2026",
                "2026-03-28" in dt_str,
                f"Start: {dt_str}",
            )

            # Check description has forecast info
            desc_lower = (description or "").lower()
            has_info = any(
                kw in desc_lower for kw in ["revenue", "forecast", "growth", "region"]
            )
            record("Event description has forecast info", has_info)
            break

    if not found:
        record(
            "Board presentation event exists",
            False,
            f"Found {len(events)} events but none for board/forecast/sales presentation",
        )

    return found


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    excel_ok = check_excel(args.agent_workspace)
    pptx_ok = check_pptx(args.agent_workspace)
    cal_ok = check_calendar()

    print(f"\n=== SUMMARY ===")
    print(f"  Excel:    {'PASS' if excel_ok else 'FAIL'}")
    print(f"  PPT:      {'PASS' if pptx_ok else 'FAIL'}")
    print(f"  Calendar: {'PASS' if cal_ok else 'FAIL'}")
    print(f"  Passed: {PASS_COUNT}, Failed: {FAIL_COUNT}")

    overall = excel_ok and pptx_ok and cal_ok
    print(f"  Overall:  {'PASS' if overall else 'FAIL'}")

    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
