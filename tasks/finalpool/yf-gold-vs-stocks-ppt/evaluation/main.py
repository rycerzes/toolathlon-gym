"""Evaluation for yf-gold-vs-stocks-ppt."""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym", user="eigent", password="camel")
PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1; print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1; print(f"  [FAIL] {name}: {str(detail)[:300]}")


def num_close(a, b, tol=1.0):
    try: return abs(float(a) - float(b)) <= tol
    except: return False


def get_expected():
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    monthly = {}
    for symbol in ['GC=F', '^DJI']:
        cur.execute("""SELECT date, close FROM yf.stock_prices
            WHERE symbol=%s AND date>='2025-03-06' AND date<='2026-03-05' ORDER BY date""", (symbol,))
        by_month = {}
        for d, c in cur.fetchall():
            mk = d.strftime("%Y-%m")
            by_month[mk] = float(c)
        monthly[symbol] = by_month
    conn.close()
    months = sorted(set(list(monthly['GC=F'].keys()) + list(monthly['^DJI'].keys())))
    prices = []
    for m in months:
        prices.append({"month": m, "gold": monthly['GC=F'].get(m), "dji": monthly['^DJI'].get(m)})
    returns = []
    for i in range(1, len(prices)):
        gr = dr = None
        if prices[i-1]["gold"] and prices[i]["gold"] and prices[i-1]["gold"] != 0:
            gr = round((prices[i]["gold"] - prices[i-1]["gold"]) / prices[i-1]["gold"] * 100, 2)
        if prices[i-1]["dji"] and prices[i]["dji"] and prices[i-1]["dji"] != 0:
            dr = round((prices[i]["dji"] - prices[i-1]["dji"]) / prices[i-1]["dji"] * 100, 2)
        returns.append({"month": prices[i]["month"], "gold_ret": gr, "dji_ret": dr})
    return {"prices": prices, "returns": returns, "months": months}


def sheet_dicts(wb, name):
    for sn in wb.sheetnames:
        if sn.strip().lower() == name.strip().lower():
            ws = wb[sn]
            rows = list(ws.iter_rows(values_only=True))
            if len(rows) < 2: return []
            hdrs = [str(h).strip() if h else "" for h in rows[0]]
            return [{hdrs[i]: row[i] for i in range(len(hdrs))} for row in rows[1:] if not all(v is None for v in row)]
    return None


def check_excel(ws_path, exp):
    print("\n=== Checking Excel ===")
    p = os.path.join(ws_path, "Gold_vs_DJI.xlsx")
    if not os.path.isfile(p):
        record("Excel file exists", False, p); return
    record("Excel file exists", True)
    wb = openpyxl.load_workbook(p, data_only=True)

    # Monthly Prices
    d = sheet_dicts(wb, "Monthly Prices")
    if d is None:
        record("Sheet Monthly Prices", False, str(wb.sheetnames))
    else:
        record("Sheet Monthly Prices", True)
        record("Monthly Prices row count", len(d) >= 12, f"Got {len(d)}")
        for ep in exp["prices"][:3] + exp["prices"][-2:]:
            m = next((r for r in d if str(r.get("Month","")).strip() == ep["month"]), None)
            if not m:
                record(f"Month {ep['month']} present", False, "Missing"); continue
            record(f"Month {ep['month']} present", True)
            if ep["gold"]:
                record(f"Month {ep['month']} gold close",
                       num_close(m.get("Gold_Close"), ep["gold"], 20.0),
                       f"{m.get('Gold_Close')} vs {ep['gold']}")
            if ep["dji"]:
                record(f"Month {ep['month']} DJI close",
                       num_close(m.get("DJI_Close"), ep["dji"], 200.0),
                       f"{m.get('DJI_Close')} vs {ep['dji']}")

    # Returns
    d = sheet_dicts(wb, "Returns")
    if d is None:
        record("Sheet Returns", False, str(wb.sheetnames))
    else:
        record("Sheet Returns", True)
        record("Returns row count", len(d) >= 11, f"Got {len(d)}")
        for er in exp["returns"][:2] + exp["returns"][-2:]:
            m = next((r for r in d if str(r.get("Month","")).strip() == er["month"]), None)
            if not m:
                record(f"Return {er['month']} present", False, "Missing"); continue
            if er["gold_ret"] is not None:
                record(f"Return {er['month']} gold",
                       num_close(m.get("Gold_Return_Pct"), er["gold_ret"], 2.0),
                       f"{m.get('Gold_Return_Pct')} vs {er['gold_ret']}")
            if er["dji_ret"] is not None:
                record(f"Return {er['month']} DJI",
                       num_close(m.get("DJI_Return_Pct"), er["dji_ret"], 2.0),
                       f"{m.get('DJI_Return_Pct')} vs {er['dji_ret']}")
    wb.close()


def check_pptx(ws_path):
    print("\n=== Checking PPTX ===")
    p = os.path.join(ws_path, "Gold_vs_Stocks.pptx")
    if not os.path.isfile(p):
        record("PPTX file exists", False, p); return
    record("PPTX file exists", True)
    try:
        from pptx import Presentation
        prs = Presentation(p)
        slides = list(prs.slides)
        record("Slide count >= 3", len(slides) >= 3, f"Got {len(slides)}")
        if len(slides) >= 1:
            title_shape = slides[0].shapes.title
            if title_shape:
                record("Slide 1 has title", True)
                t = title_shape.text.lower()
                record("Slide 1 title mentions gold", "gold" in t, title_shape.text)
            else:
                # Check all shapes for title text
                all_text = " ".join(sh.text for sh in slides[0].shapes if sh.has_text_frame).lower()
                record("Slide 1 mentions gold", "gold" in all_text, all_text[:200])
        if len(slides) >= 3:
            all_text = " ".join(sh.text for sh in slides[2].shapes if sh.has_text_frame).lower()
            record("Slide 3 has conclusion content", "conclu" in all_text or "perform" in all_text or "better" in all_text,
                   all_text[:200])
    except ImportError:
        record("python-pptx available", False, "Cannot import pptx")
    except Exception as e:
        record("PPTX readable", False, str(e))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", default=".")
    parser.add_argument("--groundtruth_workspace", default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    exp = get_expected()
    check_excel(args.agent_workspace, exp)
    check_pptx(args.agent_workspace)
    print(f"\n=== SUMMARY: {PASS_COUNT} passed, {FAIL_COUNT} failed ===")
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump({"passed": PASS_COUNT, "failed": FAIL_COUNT, "success": FAIL_COUNT == 0}, f)
    sys.exit(0 if FAIL_COUNT == 0 else 1)

if __name__ == "__main__":
    main()
