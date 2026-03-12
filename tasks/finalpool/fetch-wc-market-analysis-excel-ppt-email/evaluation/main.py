"""Evaluation script for fetch-wc-market-analysis-excel-ppt-email."""
import os
import argparse, json, os, sys
import openpyxl

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)



DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {name}: {detail_str}")


def safe_float(val, default=None):
    try:
        if val is None:
            return default
        return float(str(val).replace(',', '').replace('%', '').replace('$', '').strip())
    except (ValueError, TypeError):
        return default


def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)


def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT
    PASS_COUNT = 0
    FAIL_COUNT = 0

    # --- Excel checks ---
    excel_path = os.path.join(agent_workspace, "Competitive_Analysis.xlsx")
    check("Competitive_Analysis.xlsx exists", os.path.exists(excel_path))

    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)
        gt_path = os.path.join(groundtruth_workspace, "Competitive_Analysis.xlsx")
        gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

        # Sheet 1: Category_Comparison
        check("Category_Comparison sheet exists", "Category_Comparison" in wb.sheetnames)
        if "Category_Comparison" in wb.sheetnames:
            ws = wb["Category_Comparison"]
            headers = [str(c.value).strip() if c.value else "" for c in ws[1]]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Category_Comparison has 8 data rows", len(data_rows) == 8, f"got {len(data_rows)}")

            for expected_col in ["Category", "Own_Products", "Own_Avg_Price", "Market_Avg_Price",
                                 "Price_Position", "Own_Revenue", "Market_Revenue",
                                 "Market_Share_Pct", "Market_Growth_Rate"]:
                check(f"Category_Comparison has {expected_col}",
                      expected_col in headers, f"headers: {headers}")

            # Check sorted alphabetically
            if data_rows:
                cats = [str(r[0]) for r in data_rows if r[0]]
                check("Category_Comparison sorted alphabetically",
                      cats == sorted(cats), f"order: {cats}")

            # Compare values with groundtruth
            if gt_wb and "Category_Comparison" in gt_wb.sheetnames:
                gt_ws = gt_wb["Category_Comparison"]
                gt_rows = {str(r[0]): r for r in gt_ws.iter_rows(min_row=2, values_only=True)}
                agent_rows = {str(r[0]): r for r in data_rows}

                match_count = 0
                total_checks = 0
                for cat in gt_rows:
                    if cat in agent_rows:
                        gt_r = gt_rows[cat]
                        ag_r = agent_rows[cat]
                        # Check numeric columns (indices 1-8)
                        for idx in [1, 2, 3, 5, 6, 7, 8]:
                            gt_val = safe_float(gt_r[idx])
                            ag_val = safe_float(ag_r[idx])
                            total_checks += 1
                            if gt_val is not None and ag_val is not None:
                                if abs(gt_val - ag_val) <= max(abs(gt_val) * 0.05, 1.0):
                                    match_count += 1
                        # Check Price_Position (index 4)
                        total_checks += 1
                        if str(gt_r[4]).strip().lower() == str(ag_r[4]).strip().lower():
                            match_count += 1

                accuracy = match_count / total_checks if total_checks > 0 else 0
                check(f"Category_Comparison data accuracy >= 75%",
                      accuracy >= 0.75, f"{match_count}/{total_checks} = {accuracy:.1%}")

        # Sheet 2: Strategic_Matrix
        check("Strategic_Matrix sheet exists", "Strategic_Matrix" in wb.sheetnames)
        if "Strategic_Matrix" in wb.sheetnames:
            ws = wb["Strategic_Matrix"]
            headers = [str(c.value).strip() if c.value else "" for c in ws[1]]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Strategic_Matrix has 8 data rows", len(data_rows) == 8, f"got {len(data_rows)}")

            for expected_col in ["Category", "Price_Position", "Market_Share_Pct",
                                 "Market_Growth_Rate", "Growth_Opportunity",
                                 "Strategic_Priority", "Recommended_Action"]:
                check(f"Strategic_Matrix has {expected_col}",
                      expected_col in headers, f"headers: {headers}")

            # Verify Growth_Opportunity and Strategic_Priority logic
            if gt_wb and "Strategic_Matrix" in gt_wb.sheetnames:
                gt_ws = gt_wb["Strategic_Matrix"]
                gt_rows = {str(r[0]): r for r in gt_ws.iter_rows(min_row=2, values_only=True)}
                agent_rows = {str(r[0]): r for r in data_rows}

                go_match = 0
                sp_match = 0
                count = 0
                for cat in gt_rows:
                    if cat in agent_rows:
                        count += 1
                        gt_go = str(gt_rows[cat][4]).strip().lower()
                        ag_go = str(agent_rows[cat][4]).strip().lower()
                        if gt_go == ag_go:
                            go_match += 1
                        gt_sp = str(gt_rows[cat][5]).strip().lower()
                        ag_sp = str(agent_rows[cat][5]).strip().lower()
                        if gt_sp == ag_sp:
                            sp_match += 1

                if count > 0:
                    check(f"Growth_Opportunity accuracy >= 75%",
                          go_match / count >= 0.75, f"{go_match}/{count}")
                    check(f"Strategic_Priority accuracy >= 75%",
                          sp_match / count >= 0.75, f"{sp_match}/{count}")

        # Sheet 3: Executive_Summary
        check("Executive_Summary sheet exists", "Executive_Summary" in wb.sheetnames)
        if "Executive_Summary" in wb.sheetnames:
            ws = wb["Executive_Summary"]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Executive_Summary has >= 8 rows", len(data_rows) >= 8, f"got {len(data_rows)}")

            labels = {str(r[0]).strip(): r[1] for r in data_rows if r[0]}

            if gt_wb and "Executive_Summary" in gt_wb.sheetnames:
                gt_ws = gt_wb["Executive_Summary"]
                gt_labels = {str(r[0]).strip(): r[1]
                             for r in gt_ws.iter_rows(min_row=2, values_only=True) if r[0]}

                for key in ["Total_Own_Products", "Total_Own_Revenue", "Total_Market_Size",
                            "Overall_Market_Share_Pct"]:
                    gt_val = safe_float(gt_labels.get(key))
                    ag_val = safe_float(labels.get(key))
                    if gt_val is not None and ag_val is not None:
                        tol = max(abs(gt_val) * 0.05, 1.0)
                        check(f"Executive_Summary {key} matches",
                              abs(gt_val - ag_val) <= tol,
                              f"expected ~{gt_val}, got {ag_val}")
                    else:
                        check(f"Executive_Summary {key} present",
                              ag_val is not None, f"missing or non-numeric")

    # --- PPT checks ---
    ppt_path = os.path.join(agent_workspace, "Strategy_Presentation.pptx")
    check("Strategy_Presentation.pptx exists", os.path.exists(ppt_path))
    if os.path.exists(ppt_path):
        from pptx import Presentation
        prs = Presentation(ppt_path)
        slide_count = len(prs.slides)
        check("PPT has >= 6 slides", slide_count >= 6, f"got {slide_count}")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        for term in ["market", "competitive", "growth", "strategy", "priority"]:
            check(f"PPT mentions '{term}'", term in all_text)

    # --- Email checks ---
    try:
        conn = get_conn()
        cur = conn.cursor()

        # Check CEO email
        cur.execute(
            "SELECT subject, to_addr, body_text FROM email.messages WHERE subject ILIKE %s",
            ('%executive summary%',)
        )
        ceo_emails = cur.fetchall()
        check("CEO email sent", len(ceo_emails) >= 1)
        if ceo_emails:
            to_str = str(ceo_emails[0][1]).lower()
            check("CEO email to ceo@company.com", "ceo@company.com" in to_str)

        # Check product team email
        cur.execute(
            "SELECT subject, to_addr, body_text FROM email.messages WHERE subject ILIKE %s",
            ('%detailed findings%',)
        )
        pt_emails = cur.fetchall()
        check("Product team email sent", len(pt_emails) >= 1)
        if pt_emails:
            to_str = str(pt_emails[0][1]).lower()
            check("Product team email to product_team@company.com",
                  "product_team@company.com" in to_str)

        conn.close()
    except Exception as e:
        check("Email check", False, str(e))

    return FAIL_COUNT == 0, f"Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
