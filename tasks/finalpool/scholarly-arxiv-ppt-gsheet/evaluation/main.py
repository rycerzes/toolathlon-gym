"""
Evaluation for scholarly-arxiv-ppt-gsheet task.
Checks: PPT file, GSheet, Word document.
"""
import argparse
import os
import sys

import psycopg2
from docx import Document
from pptx import Presentation

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0

RELEVANT_KEYWORDS = ["chain-of-thought", "self-consistency", "tree of thoughts", "chain of thought", "reasoning"]
PAPER_TITLES_LOWER = [
    "chain-of-thought prompting elicits reasoning",
    "self-consistency improves chain of thought",
    "tree of thoughts: deliberate problem solving",
]


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    ppt_path = os.path.join(agent_workspace, "LLM_Reasoning_Survey.pptx")
    if not os.path.isfile(ppt_path):
        record("PPT file LLM_Reasoning_Survey.pptx exists", False, f"Not found at: {ppt_path}")
        return
    record("PPT file LLM_Reasoning_Survey.pptx exists", True)

    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        record("PPT file readable", False, str(e))
        return
    record("PPT file readable", True)

    num_slides = len(prs.slides)
    record("PPT has at least 5 slides", num_slides >= 5, f"Found {num_slides} slides")

    # Gather all text from slides
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text.lower())
    full_text = " ".join(all_text)

    # Check title slide
    has_title = "survey" in full_text and ("reasoning" in full_text or "llm" in full_text.lower())
    record("PPT contains survey/reasoning title content", has_title, "Looked for 'survey' and 'reasoning'")

    # Check paper content
    has_chain = any(kw in full_text for kw in ["chain-of-thought", "chain of thought"])
    has_self_consistency = "self-consistency" in full_text or "self consistency" in full_text
    has_tree = "tree of thoughts" in full_text or "tree of thought" in full_text
    record("PPT mentions chain-of-thought paper", has_chain)
    record("PPT mentions self-consistency paper", has_self_consistency)
    record("PPT mentions tree of thoughts paper", has_tree)

    # Check summary slide
    has_summary = "summary" in full_text or "comparison" in full_text or "conclusion" in full_text
    record("PPT has a summary/conclusion slide", has_summary)


def check_gsheet():
    print("\n=== Checking Google Sheet ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        cur.execute("SELECT id, title FROM gsheet.spreadsheets")
        spreadsheets = cur.fetchall()

        target_ss = None
        for sid, title in spreadsheets:
            if title and ("llm" in title.lower() or "reasoning" in title.lower()) and "paper" in title.lower():
                target_ss = sid
                break

        record("GSheet 'LLM Reasoning Paper Tracker' exists",
               target_ss is not None,
               f"Found sheets: {[t for _, t in spreadsheets]}")

        if target_ss is None:
            conn.close()
            return

        # Find any sheet in this spreadsheet
        cur.execute("SELECT id, title FROM gsheet.sheets WHERE spreadsheet_id = %s", (target_ss,))
        sheets = cur.fetchall()
        record("GSheet has at least one sheet", len(sheets) > 0, f"Found: {sheets}")

        if not sheets:
            conn.close()
            return

        sheet_id = sheets[0][0]

        # Count data rows (exclude header)
        cur.execute("""
            SELECT COUNT(DISTINCT row_index) FROM gsheet.cells
            WHERE spreadsheet_id = %s AND sheet_id = %s AND row_index > 0
        """, (target_ss, sheet_id))
        data_rows = cur.fetchone()[0]
        record("GSheet has at least 3 data rows", data_rows >= 3, f"Found {data_rows} data rows")

        # Check that paper titles appear in cells
        cur.execute("""
            SELECT LOWER(value) FROM gsheet.cells
            WHERE spreadsheet_id = %s AND sheet_id = %s
        """, (target_ss, sheet_id))
        cell_values = [row[0] for row in cur.fetchall() if row[0]]
        all_cells_text = " ".join(cell_values)

        has_chain = "chain-of-thought" in all_cells_text or "chain of thought" in all_cells_text
        has_self = "self-consistency" in all_cells_text or "self consistency" in all_cells_text
        has_tree = "tree of thoughts" in all_cells_text or "tree of thought" in all_cells_text
        record("GSheet contains chain-of-thought paper entry", has_chain)
        record("GSheet contains self-consistency paper entry", has_self)
        record("GSheet contains tree of thoughts paper entry", has_tree)

        conn.close()
    except Exception as e:
        record("GSheet connection", False, str(e))


def check_word(agent_workspace):
    print("\n=== Checking Word Document ===")
    doc_path = os.path.join(agent_workspace, "LLM_Reasoning_Literature_Review.docx")
    if not os.path.isfile(doc_path):
        record("Word file LLM_Reasoning_Literature_Review.docx exists", False, f"Not found at: {doc_path}")
        return
    record("Word file LLM_Reasoning_Literature_Review.docx exists", True)

    try:
        doc = Document(doc_path)
    except Exception as e:
        record("Word file readable", False, str(e))
        return
    record("Word file readable", True)

    full_text = "\n".join(p.text for p in doc.paragraphs).lower()

    has_heading = "literature review" in full_text and ("llm" in full_text or "reasoning" in full_text)
    record("Word has 'Literature Review' heading with LLM/reasoning", has_heading)

    has_intro = len(full_text) > 300
    record("Word has substantial content (intro + sections)", has_intro, f"Text length: {len(full_text)}")

    has_chain = "chain-of-thought" in full_text or "chain of thought" in full_text
    has_self = "self-consistency" in full_text or "self consistency" in full_text
    has_tree = "tree of thoughts" in full_text or "tree of thought" in full_text
    record("Word mentions chain-of-thought method", has_chain)
    record("Word mentions self-consistency method", has_self)
    record("Word mentions tree of thoughts method", has_tree)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_gsheet()
    check_word(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ===")
    if FAIL_COUNT > 0:
        print(f"{FAIL_COUNT} checks failed")
        sys.exit(1)
    else:
        print("All checks passed!")
        sys.exit(0)


if __name__ == "__main__":
    main()
