"""Evaluation for terminal-wc-yf-ppt-notion-email.
Checks:
1. Market_Strategy_Presentation.pptx (7 slides)
2. Notion "Market Strategy Tracker" database with 5 entries
3. Email to ceo@company.com
4. Email to marketing_team@company.com
5. market_correlation.py and category_analysis.py scripts exist
6. market_correlation.json output
7. category_market_analysis.json output
"""
import argparse
import json
import os
import sys

import psycopg2
from pptx import Presentation

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def get_expected_from_db():
    """Query WC and YF schemas dynamically for expected values."""
    defaults = {
        "top_category": "electronics",
        "amzn_pct_change": 26.0,
    }
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        # Top category by revenue
        cur.execute("""
            SELECT p.categories->0->>'name' as cat, SUM((li->>'total')::numeric) as rev
            FROM wc.orders o, jsonb_array_elements(o.line_items) li
            JOIN wc.products p ON (li->>'product_id')::int = p.id
            WHERE o.status NOT IN ('cancelled','refunded','failed')
            GROUP BY cat ORDER BY rev DESC LIMIT 1
        """)
        row = cur.fetchone()
        if row and row[0]:
            defaults["top_category"] = row[0].lower()

        # AMZN percentage change (earliest to latest)
        cur.execute("""
            SELECT
                (SELECT close FROM yf.stock_prices WHERE symbol='AMZN' ORDER BY date ASC LIMIT 1) as first_close,
                (SELECT close FROM yf.stock_prices WHERE symbol='AMZN' ORDER BY date DESC LIMIT 1) as last_close
        """)
        row = cur.fetchone()
        if row and row[0] and row[1]:
            defaults["amzn_pct_change"] = float((row[1] - row[0]) / row[0] * 100)

        cur.close()
        conn.close()
    except Exception as e:
        print(f"  [WARN] DB query for expected values failed, using defaults: {e}")
    return defaults


EXPECTED = get_expected_from_db()


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def num_close(a, b, tol=2.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except Exception:
        return False


def check_pptx(workspace):
    print("\n=== Check 1: Market_Strategy_Presentation.pptx ===")
    path = os.path.join(workspace, "Market_Strategy_Presentation.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    prs = Presentation(path)
    slides = list(prs.slides)
    check("Has 7 slides", len(slides) == 7, f"Found {len(slides)}")

    # Collect all text from all slides
    all_texts = []
    for slide in slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    slide_text.append(para.text)
        all_texts.append(" ".join(slide_text).lower())

    full_text = " ".join(all_texts)

    # Slide 1: Title
    check("Slide 1 has title 'Market Positioning'",
          "market" in all_texts[0] and "positioning" in all_texts[0] if len(all_texts) > 0 else False,
          all_texts[0][:100] if all_texts else "no slides")

    # Slide 2: Market Overview with AMZN, JPM, DJI
    if len(all_texts) > 1:
        check("Slide 2 mentions AMZN", "amzn" in all_texts[1] or "amazon" in all_texts[1],
              all_texts[1][:100])
        check("Slide 2 mentions JPM", "jpm" in all_texts[1] or "jpmorgan" in all_texts[1],
              all_texts[1][:100])
        check("Slide 2 mentions DJI", "dji" in all_texts[1] or "dow" in all_texts[1],
              all_texts[1][:100])

    # Slide 3: Revenue trends
    if len(all_texts) > 2:
        check("Slide 3 has revenue data",
              "revenue" in all_texts[2] or "2025" in all_texts[2],
              all_texts[2][:100])

    # Slide 4: Category performance
    if len(all_texts) > 3:
        check("Slide 4 mentions Electronics",
              "electronics" in all_texts[3], all_texts[3][:100])
        check("Slide 4 mentions category",
              "category" in all_texts[3] or "tv" in all_texts[3] or "audio" in all_texts[3],
              all_texts[3][:100])

    # Slide 5: Correlation
    if len(all_texts) > 4:
        check("Slide 5 mentions correlation",
              "correlation" in all_texts[4], all_texts[4][:100])

    # Slide 6: Recommendations
    if len(all_texts) > 5:
        check("Slide 6 has recommendations",
              "recommend" in all_texts[5] or "pricing" in all_texts[5] or "strategy" in all_texts[5],
              all_texts[5][:100])

    # Slide 7: Next steps
    if len(all_texts) > 6:
        check("Slide 7 has next steps",
              "next" in all_texts[6] or "action" in all_texts[6] or "step" in all_texts[6],
              all_texts[6][:100])

    # Check key numbers appear somewhere (dynamically computed)
    amzn_pct_str = str(int(round(EXPECTED["amzn_pct_change"])))
    check(f"Mentions AMZN ~{amzn_pct_str}% change", amzn_pct_str in full_text,
          f"No {amzn_pct_str}% mention (expected ~{EXPECTED['amzn_pct_change']:.1f}%)")
    check(f"Mentions {EXPECTED['top_category'].title()} as top category",
          EXPECTED["top_category"] in full_text,
          f"No {EXPECTED['top_category']} mention")


def check_notion():
    print("\n=== Check 2: Notion Market Strategy Tracker ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        cur.execute("SELECT id, title FROM notion.databases")
        dbs = cur.fetchall()
        tracker_db = None
        for db_id, title in dbs:
            title_str = ""
            if isinstance(title, list):
                title_str = " ".join(
                    item.get("text", {}).get("content", "")
                    for item in title if isinstance(item, dict))
            elif isinstance(title, str):
                try:
                    parsed = json.loads(title)
                    if isinstance(parsed, list):
                        title_str = " ".join(
                            item.get("text", {}).get("content", "")
                            for item in parsed if isinstance(item, dict))
                    else:
                        title_str = str(title)
                except Exception:
                    title_str = str(title)
            else:
                title_str = str(title) if title else ""
            if "market" in title_str.lower() and "strategy" in title_str.lower() and "tracker" in title_str.lower():
                tracker_db = (db_id, title_str)
                break

        check("Market Strategy Tracker DB exists", tracker_db is not None,
              f"Databases found: {[d[1] for d in dbs]}")

        if tracker_db:
            # Check properties
            cur.execute("SELECT properties FROM notion.databases WHERE id = %s", (tracker_db[0],))
            props = cur.fetchone()[0]
            if isinstance(props, str):
                props = json.loads(props)
            prop_names = [k.lower().replace("_", " ") for k in props.keys()] if props else []
            check("Has Initiative property",
                  any("initiative" in p or "title" in str(props.get(k, {}).get("type", ""))
                      for k, p in [(k, k.lower()) for k in (props or {}).keys()]),
                  f"Props: {list((props or {}).keys())}")
            check("Has Category property",
                  any("category" in p for p in prop_names),
                  f"Props: {prop_names}")
            check("Has Market_Condition property",
                  any("market" in p and "condition" in p for p in prop_names),
                  f"Props: {prop_names}")

            # Check pages
            cur.execute("""
                SELECT id, properties FROM notion.pages
                WHERE parent->>'database_id' = %s AND NOT archived
            """, (tracker_db[0],))
            pages = cur.fetchall()
            check("Has 5 initiative entries", len(pages) == 5, f"Found {len(pages)}")

            # Check specific entries
            page_titles = []
            for pid, props in pages:
                if isinstance(props, str):
                    props = json.loads(props)
                for k, v in (props or {}).items():
                    if isinstance(v, dict) and "title" in v:
                        title_items = v.get("title", [])
                        if isinstance(title_items, list):
                            for item in title_items:
                                if isinstance(item, dict):
                                    page_titles.append(item.get("text", {}).get("content", ""))
            all_titles = " ".join(page_titles).lower()
            check("Has TV & Home Theater initiative",
                  "tv" in all_titles or "home theater" in all_titles, f"Titles: {page_titles}")
            check("Has Electronics initiative",
                  "electronics" in all_titles, f"Titles: {page_titles}")
            check("Has Audio initiative",
                  "audio" in all_titles, f"Titles: {page_titles}")
    except Exception as e:
        check("Notion check", False, str(e))
    finally:
        cur.close()
        conn.close()


def check_emails():
    print("\n=== Check 3: Emails ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        # Check CEO email
        cur.execute("""
            SELECT subject, body_text, to_addr FROM email.messages
            WHERE to_addr::text LIKE '%%ceo@company.com%%'
            ORDER BY date DESC LIMIT 1
        """)
        ceo_email = cur.fetchone()
        check("CEO email exists", ceo_email is not None,
              "No email to ceo@company.com found")
        if ceo_email:
            subj = (ceo_email[0] or "").lower()
            body = (ceo_email[1] or "").lower()
            check("CEO email subject mentions strategy or summary",
                  "strategy" in subj or "summary" in subj or "executive" in subj,
                  f"Subject: {ceo_email[0]}")
            check("CEO email body mentions AMZN",
                  "amzn" in body or "amazon" in body, "No AMZN mention")
            check("CEO email body mentions Electronics",
                  "electronics" in body, "No Electronics mention")
            check("CEO email body mentions recommendation",
                  "recommend" in body or "pricing" in body or "independent" in body or "strategy" in body,
                  "No recommendation")

        # Check marketing team email
        cur.execute("""
            SELECT subject, body_text, to_addr FROM email.messages
            WHERE to_addr::text LIKE '%%marketing_team@company.com%%'
            ORDER BY date DESC LIMIT 1
        """)
        mkt_email = cur.fetchone()
        check("Marketing team email exists", mkt_email is not None,
              "No email to marketing_team@company.com found")
        if mkt_email:
            subj = (mkt_email[0] or "").lower()
            body = (mkt_email[1] or "").lower()
            check("Marketing email subject mentions category",
                  "category" in subj or "performance" in subj or "q2" in subj,
                  f"Subject: {mkt_email[0]}")
            check("Marketing email body mentions categories",
                  "electronics" in body and ("camera" in body or "watch" in body),
                  "Missing category details")
    except Exception as e:
        check("Email check", False, str(e))
    finally:
        cur.close()
        conn.close()


def check_reverse_validation():
    print("\n=== Reverse Validation ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        # Check Notion tracker has exactly 5 entries, no noise
        cur.execute("SELECT id, title FROM notion.databases")
        dbs = cur.fetchall()
        tracker_db = None
        for db_id, title in dbs:
            title_str = ""
            if isinstance(title, list):
                title_str = " ".join(
                    item.get("text", {}).get("content", "")
                    for item in title if isinstance(item, dict))
            elif isinstance(title, str):
                try:
                    parsed = json.loads(title)
                    if isinstance(parsed, list):
                        title_str = " ".join(
                            item.get("text", {}).get("content", "")
                            for item in parsed if isinstance(item, dict))
                    else:
                        title_str = str(title)
                except Exception:
                    title_str = str(title)
            else:
                title_str = str(title) if title else ""
            if "market" in title_str.lower() and "strategy" in title_str.lower() and "tracker" in title_str.lower():
                tracker_db = (db_id, title_str)
                break

        if tracker_db:
            cur.execute("""
                SELECT id, properties FROM notion.pages
                WHERE parent->>'database_id' = %s AND NOT archived
            """, (tracker_db[0],))
            pages = cur.fetchall()
            check("Notion tracker has no extra noise entries (exactly 5)",
                  len(pages) == 5,
                  f"Found {len(pages)} entries, expected exactly 5")

            # Check no unrelated categories in notion entries
            all_props_text = " ".join(json.dumps(p[1]).lower() for p in pages if p[1])
            noise_categories = ["healthcare", "food", "real estate", "energy"]
            for cat in noise_categories:
                check(f"Notion entries do not contain noise category '{cat}'",
                      cat not in all_props_text,
                      f"Found '{cat}' in Notion data")

        # Check no emails sent to noise recipients
        noise_recipients = [
            "all-staff@company.com",
            "hr@company.com",
            "newsletter@company.com",
            "sales_team@company.com",
        ]
        for addr in noise_recipients:
            cur.execute(
                "SELECT COUNT(*) FROM email.messages WHERE to_addr::text ILIKE %s",
                (f"%{addr}%",),
            )
            cnt = cur.fetchone()[0]
            check(f"No email sent to noise recipient {addr}", cnt == 0,
                  f"Found {cnt} emails to {addr}")
    except Exception as e:
        check("Reverse validation", False, str(e))
    finally:
        cur.close()
        conn.close()


def check_scripts(workspace):
    print("\n=== Check 4: Scripts and JSON outputs ===")
    check("market_correlation.py exists",
          os.path.exists(os.path.join(workspace, "market_correlation.py")))
    check("category_analysis.py exists",
          os.path.exists(os.path.join(workspace, "category_analysis.py")))

    # Check market_correlation.json
    corr_path = os.path.join(workspace, "market_correlation.json")
    if os.path.exists(corr_path):
        check("market_correlation.json exists", True)
        try:
            with open(corr_path) as f:
                data = json.load(f)
            check("market_correlation.json has recommendation",
                  "recommendation" in data or "strategy" in str(data).lower(),
                  f"Keys: {list(data.keys()) if isinstance(data, dict) else 'not dict'}")
            # Check correlation values exist
            data_str = json.dumps(data).lower()
            check("market_correlation.json mentions AMZN correlation",
                  "amzn" in data_str, f"Content: {data_str[:200]}")
        except Exception as e:
            check("market_correlation.json valid JSON", False, str(e))
    else:
        check("market_correlation.json exists", False, f"Not found at {corr_path}")

    # Check category_market_analysis.json
    cat_path = os.path.join(workspace, "category_market_analysis.json")
    if os.path.exists(cat_path):
        check("category_market_analysis.json exists", True)
        try:
            with open(cat_path) as f:
                data = json.load(f)
            data_str = json.dumps(data).lower()
            check("category_market_analysis.json has category data",
                  "electronics" in data_str, f"Content: {data_str[:200]}")
        except Exception as e:
            check("category_market_analysis.json valid JSON", False, str(e))
    else:
        check("category_market_analysis.json exists", False, f"Not found at {cat_path}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_notion()
    check_emails()
    check_scripts(args.agent_workspace)
    check_reverse_validation()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if accuracy >= 70:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
