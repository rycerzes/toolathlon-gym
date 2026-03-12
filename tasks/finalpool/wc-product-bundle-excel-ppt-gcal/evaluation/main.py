"""Evaluation for wc-product-bundle-excel-ppt-gcal."""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=0.5):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def get_expected_data():
    """Compute expected co-purchase data from read-only DB."""
    from collections import defaultdict
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT o.id as order_id,
               (li->>'product_id')::int as product_id,
               li->>'name' as product_name,
               (li->>'price')::numeric as price
        FROM wc.orders o, jsonb_array_elements(o.line_items) li
        WHERE o.status IN ('completed', 'processing')
    """)
    rows = cur.fetchall()

    order_items = defaultdict(list)
    product_info = {}
    for order_id, pid, pname, price in rows:
        order_items[order_id].append((pid, pname, float(price)))
        product_info[pid] = (pname, float(price))

    pair_counts = defaultdict(int)
    pair_revenues = defaultdict(list)
    for order_id, items in order_items.items():
        pids = sorted(set(item[0] for item in items))
        item_map = {item[0]: item for item in items}
        for i in range(len(pids)):
            for j in range(i + 1, len(pids)):
                pa, pb = pids[i], pids[j]
                pair_counts[(pa, pb)] += 1
                combined = item_map[pa][2] + item_map[pb][2]
                pair_revenues[(pa, pb)].append(combined)

    qualified = []
    for (pa, pb), count in pair_counts.items():
        if count >= 2:
            avg_rev = round(sum(pair_revenues[(pa, pb)]) / len(pair_revenues[(pa, pb)]), 2)
            priority = round(count * avg_rev, 2)
            qualified.append({
                'pid_a': pa, 'name_a': product_info[pa][0],
                'price_a': product_info[pa][1],
                'pid_b': pb, 'name_b': product_info[pb][0],
                'price_b': product_info[pb][1],
                'count': count, 'avg_rev': avg_rev, 'priority': priority
            })

    qualified.sort(key=lambda x: x['priority'], reverse=True)
    cur.close()
    conn.close()
    return qualified


def check_excel(agent_workspace):
    """Check Bundle_Analysis.xlsx."""
    print("\n=== Checking Bundle_Analysis.xlsx ===")
    agent_file = os.path.join(agent_workspace, "Bundle_Analysis.xlsx")
    if not os.path.isfile(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False
    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(agent_file, data_only=True)
    except Exception as e:
        record("Excel readable", False, str(e))
        return False

    expected = get_expected_data()
    all_ok = True

    # --- Check Co-Purchase Matrix sheet ---
    cp_sheet = None
    for name in wb.sheetnames:
        if "co" in name.lower() and "purchase" in name.lower():
            cp_sheet = wb[name]
            break
        if "matrix" in name.lower():
            cp_sheet = wb[name]
            break
    if cp_sheet is None:
        record("Sheet 'Co-Purchase Matrix' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Co-Purchase Matrix' exists", True)
        rows = list(cp_sheet.iter_rows(min_row=2, values_only=True))
        record("Co-Purchase Matrix has data rows", len(rows) >= len(expected),
               f"Expected {len(expected)}, got {len(rows)}")

        for ep in expected:
            found = False
            for r in rows:
                if r and len(r) >= 5:
                    # Match by product names (partial match)
                    ra = str(r[0]).lower() if r[0] else ""
                    rb = str(r[1]).lower() if r[1] else ""
                    ea = ep['name_a'].lower()[:30]
                    eb = ep['name_b'].lower()[:30]
                    if (ea[:20] in ra and eb[:20] in rb) or (ea[:20] in rb and eb[:20] in ra):
                        found = True
                        ok_count = num_close(r[2], ep['count'], 0)
                        record(f"Pair count ({ep['name_a'][:25]}...)", ok_count,
                               f"Expected {ep['count']}, got {r[2]}")
                        if not ok_count:
                            all_ok = False
                        ok_priority = num_close(r[4], ep['priority'], 5.0)
                        record(f"Pair priority ({ep['name_a'][:25]}...)", ok_priority,
                               f"Expected {ep['priority']}, got {r[4]}")
                        if not ok_priority:
                            all_ok = False
                        break
            if not found:
                record(f"Pair found: {ep['name_a'][:30]}...", False, "Not found in sheet")
                all_ok = False

    # --- Check Bundle Proposals sheet ---
    bp_sheet = None
    for name in wb.sheetnames:
        if "bundle" in name.lower() and "proposal" in name.lower():
            bp_sheet = wb[name]
            break
    if bp_sheet is None:
        record("Sheet 'Bundle Proposals' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Bundle Proposals' exists", True)
        rows = list(bp_sheet.iter_rows(min_row=2, values_only=True))
        record("Bundle Proposals has >= 5 rows", len(rows) >= 5,
               f"Got {len(rows)} rows")

        # Check top bundle pricing
        for i, ep in enumerate(expected[:3]):
            combined = round(ep['price_a'] + ep['price_b'], 2)
            bundle_price = round(combined * 0.90, 2)
            found = False
            for r in rows:
                if r and len(r) >= 6:
                    rname = str(r[0]).lower() if r[0] else ""
                    # Match both product names to avoid ambiguity
                    name_a_short = ep['name_a'].lower()[:15]
                    name_b_short = ep['name_b'].lower()[:15]
                    if name_a_short in rname and name_b_short in rname:
                        found = True
                        ok_bp = num_close(r[4], bundle_price, 2.0)
                        record(f"Bundle price rank {i+1}", ok_bp,
                               f"Expected ~{bundle_price}, got {r[4]}")
                        if not ok_bp:
                            all_ok = False
                        break
            if not found:
                record(f"Bundle rank {i+1} found", False, "Not found in proposals")
                all_ok = False

    # --- Check Category Insights sheet ---
    ci_sheet = None
    for name in wb.sheetnames:
        if "category" in name.lower() or "insight" in name.lower():
            ci_sheet = wb[name]
            break
    if ci_sheet is None:
        record("Sheet 'Category Insights' exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Sheet 'Category Insights' exists", True)
        rows = list(ci_sheet.iter_rows(min_row=2, values_only=True))
        record("Category Insights has data rows", len(rows) >= 3,
               f"Got {len(rows)} rows")

        # Check that Electronics appears (most common category)
        has_electronics = any(r and r[0] and "electronics" in str(r[0]).lower() for r in rows)
        record("Electronics category in insights", has_electronics)
        if not has_electronics:
            all_ok = False

    return all_ok


def check_pptx(agent_workspace):
    """Check Bundle_Presentation.pptx."""
    print("\n=== Checking Bundle_Presentation.pptx ===")
    pptx_file = os.path.join(agent_workspace, "Bundle_Presentation.pptx")
    if not os.path.isfile(pptx_file):
        record("PPTX file exists", False, f"Not found: {pptx_file}")
        return False
    record("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_file)
    except Exception as e:
        record("PPTX readable", False, str(e))
        return False

    slide_count = len(prs.slides)
    record("PPTX has >= 5 slides", slide_count >= 5, f"Got {slide_count}")

    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text.lower() + " "

    record("PPTX mentions 'bundle'", "bundle" in all_text)
    record("PPTX mentions 'revenue' or 'price'",
           "revenue" in all_text or "price" in all_text)
    record("PPTX mentions 'co-purchase' or 'co purchase' or 'purchase'",
           "co-purchase" in all_text or "co purchase" in all_text or "purchase" in all_text)
    record("PPTX mentions 'timeline' or 'implementation' or 'launch'",
           "timeline" in all_text or "implementation" in all_text or "launch" in all_text)

    return slide_count >= 5


def check_calendar():
    """Check Google Calendar for 3 bundle launch meetings."""
    print("\n=== Checking Google Calendar ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        cur.execute("""
            SELECT summary, start_datetime, end_datetime, description
            FROM gcal.events
            WHERE LOWER(summary) LIKE '%bundle%'
               OR LOWER(summary) LIKE '%launch%'
               OR LOWER(summary) LIKE '%priority%'
               OR LOWER(summary) LIKE '%category%'
        """)
        events = cur.fetchall()
        record("At least 3 bundle-related calendar events",
               len(events) >= 3,
               f"Found {len(events)} matching events")

        if len(events) >= 3:
            # Check dates are in March 16-20 week
            march_events = []
            for ev in events:
                start_str = str(ev[1]) if ev[1] else ""
                if "2026-03-1" in start_str or "2026-03-20" in start_str:
                    march_events.append(ev)
            record("Events scheduled in March 16-20 week",
                   len(march_events) >= 3,
                   f"Found {len(march_events)} events in target week")

            # Check for 1-hour duration
            hour_events = 0
            for ev in events:
                if ev[1] and ev[2]:
                    duration = (ev[2] - ev[1]).total_seconds()
                    if 3000 <= duration <= 7200:  # between 50min and 2hr
                        hour_events += 1
            record("Events are ~1 hour each", hour_events >= 3,
                   f"{hour_events} events have ~1hr duration")

        cur.close()
        conn.close()
        return len(events) >= 3
    except Exception as e:
        record("Calendar check", False, str(e))
        return False


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    excel_ok = check_excel(args.agent_workspace)
    pptx_ok = check_pptx(args.agent_workspace)
    cal_ok = check_calendar()

    total = PASS_COUNT + FAIL_COUNT
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ===")
    if FAIL_COUNT > 0:
        print(f"{FAIL_COUNT} checks failed")

    overall = excel_ok and pptx_ok and cal_ok
    print(f"Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
